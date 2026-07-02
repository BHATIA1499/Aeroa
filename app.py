"""
Skuvvy — Flask backend (Phase 1: Auth + Supabase)
========================================================
Routes:
  Public
    GET  /                    → marketing site
    GET  /login               → login page
    GET  /signup              → signup page
    GET  /health              → health check

  Auth
    POST /auth/signup         → create account
    POST /auth/login          → sign in, set session
    POST /auth/logout         → clear session
    GET  /auth/me             → current user info (includes role)

  App (requires login)
    GET  /dashboard           → main app UI
    POST /api/upload          → parse & store CSV/XLSX
    GET  /api/analysis        → return latest analysis
    GET  /api/analysis/role   → extended role-based analysis data
    GET  /api/uploads         → list user's uploads
    POST /api/chat            → standard AI chat
    POST /api/chat/stream     → SSE streaming chat
    GET  /api/quick-insights  → 3-bullet AI insights
    GET  /api/user/role       → get current role from session
    POST /api/user/role       → set role in session + Supabase

  Stripe webhooks
    POST /webhooks/stripe     → handle subscription events
"""

import os
import io
import json
import uuid
import math
import functools
from datetime import datetime, timezone

import pandas as pd
from flask import (Flask, request, jsonify, session,
                   Response, send_from_directory, redirect, url_for)
from urllib.parse import quote
from flask_cors import CORS
from dotenv import load_dotenv
import anthropic
from supabase import create_client, Client
import stripe
from datetime import timedelta

# ── Security layer ──────────────────────────────────────────────────────────────
from security.encryption   import encrypt_file, file_sha256
from security.audit        import (AuditLogger, USER_LOGIN, USER_LOGOUT, AUTH_FAILED,
                                   FILE_UPLOAD, FILE_DELETED, REPORT_GENERATED,
                                   AI_QUERY, DATA_EXPORT, ROLE_CHANGE, USER_CREATED,
                                   RATE_LIMITED, PERMISSION_DENIED, SETTINGS_CHANGED)
from security.rbac         import require_permission, has_permission, require_company_ownership, VALID_ROLES
from security.rate_limiter import create_limiter, auth_limit
from security.file_pipeline import SecureFilePipeline
from security.ai_privacy   import prepare_for_ai, private_mode_notice
from security.headers      import apply_security_headers, generate_csrf_token, get_csrf_token
from security.validators   import sanitise, sanitise_raw, validate_email, validate_password, safe_filename
from security.auth_schemas import (SignupSchema, LoginSchema, ForgotPasswordSchema,
                                   ResetPasswordSchema, VerifyEmailSchema)
from security.lockout      import lockout, MAX_FAILS
from pydantic              import ValidationError
import time
from apscheduler.schedulers.background import BackgroundScheduler

load_dotenv()

# ── App setup ────────────────────────────────────────────────
app = Flask(__name__, static_folder="static", static_url_path="/static")

# Behind Railway's TLS-terminating proxy, Flask sees the request as http and the
# internal host. Trust the proxy's X-Forwarded-Proto / -Host so request.host_url,
# url_for(_external=True), and OAuth redirect_to values are built with the correct
# https scheme and public hostname. Without this, redirect_to becomes http://…,
# fails Supabase's https-only allow-list, and OAuth falls back to the Site URL.
from werkzeug.middleware.proxy_fix import ProxyFix
app.wsgi_app = ProxyFix(app.wsgi_app, x_proto=1, x_host=1)

app.secret_key = os.environ.get("FLASK_SECRET_KEY", os.urandom(32))
app.config["SESSION_COOKIE_SAMESITE"] = "Lax"
app.config["SESSION_COOKIE_SECURE"] = os.environ.get("FLASK_ENV") == "production"
app.config["SESSION_COOKIE_HTTPONLY"] = True
app.config["SESSION_COOKIE_NAME"]     = "__Host-session" if os.environ.get("FLASK_ENV") == "production" else "session"
app.config["PERMANENT_SESSION_LIFETIME"] = 86400   # 24 hours
CORS(app, supports_credentials=True)

# ── Clients ───────────────────────────────────────────────────
ai_client = anthropic.Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY"))
MODEL = "claude-sonnet-4-5"

SUPABASE_URL = os.environ.get("SUPABASE_URL", "")
SUPABASE_SERVICE_KEY = os.environ.get("SUPABASE_SERVICE_KEY", "")

supabase: Client = create_client(SUPABASE_URL, SUPABASE_SERVICE_KEY)

stripe.api_key = os.environ.get("STRIPE_SECRET_KEY", "")
STRIPE_WEBHOOK_SECRET = os.environ.get("STRIPE_WEBHOOK_SECRET", "")

# ── Security middleware ────────────────────────────────────────────────────────
apply_security_headers(app)
limiter       = create_limiter(app)
audit         = AuditLogger(supabase)
file_pipeline = SecureFilePipeline()

# Plan limits
PLAN_LIMITS = {
    "trial":   {"skus": 50,     "ai_msgs": 20},
    "starter": {"skus": 200,    "ai_msgs": 50},
    "growth":  {"skus": 999999, "ai_msgs": 999999},
    "studio":  {"skus": 999999, "ai_msgs": 999999},
}

# ═══════════════════════════════════════════════════════════════
# AUTH HELPERS
# ═══════════════════════════════════════════════════════════════

def get_user_db() -> Client:
    """Return a Supabase client authenticated as the current user (uses their JWT)."""
    access_token = session.get("access_token")
    refresh_token = session.get("refresh_token", "")
    if not access_token:
        return supabase
    client = create_client(SUPABASE_URL, SUPABASE_SERVICE_KEY)
    try:
        client.auth.set_session(access_token, refresh_token)
    except Exception:
        pass
    return client


def get_current_user():
    """Return profile dict from Supabase, or None."""
    uid = session.get("user_id")
    if not uid:
        return None
    # Use session-cached profile for every request — avoids a DB round-trip
    # and survives if PostgREST grants haven't been applied yet.
    cached = session.get("user_profile")
    if cached and cached.get("id") == uid:
        return cached
    try:
        db = get_user_db()
        res = db.table("profiles").select("*").eq("id", uid).single().execute()
        if res.data:
            session["user_profile"] = res.data
            return res.data
    except Exception:
        pass
    # Graceful fallback: construct minimal profile from session
    return {
        "id": uid,
        "email": session.get("user_email", ""),
        "full_name": session.get("user_name", ""),
        "plan": session.get("user_plan", "trial"),
        "trial_ends": session.get("user_trial_ends"),
        "ai_messages_used": 0,
    }


def require_auth(f):
    """Decorator: require valid session. Passes user kwarg (includes company_id, role)."""
    @functools.wraps(f)
    def decorated(*args, **kwargs):
        user = get_current_user()
        if not user:
            if request.path.startswith("/api/"):
                audit.log(AUTH_FAILED, request=request, status="FAILURE",
                          metadata={"path": request.path, "method": request.method})
                return jsonify({"error": "Authentication required", "code": "UNAUTHENTICATED"}), 401
            return redirect("/login")
        # Attach runtime role from session (user may have switched role)
        user["role"] = session.get("user_role", user.get("role", "MA"))
        return f(*args, user=user, **kwargs)
    return decorated


def check_plan_limit(user, resource):
    """Return (allowed: bool, message: str)."""
    plan = user.get("plan", "trial")

    # Trial expiry check
    trial_ends = user.get("trial_ends")
    if plan == "trial" and trial_ends:
        ends = datetime.fromisoformat(trial_ends.replace("Z", "+00:00"))
        if datetime.now(timezone.utc) > ends:
            return False, "Your free trial has ended. Please upgrade to continue."

    limit = PLAN_LIMITS.get(plan, PLAN_LIMITS["trial"])

    if resource == "ai_msg":
        # Reset monthly counter if needed
        reset_at = user.get("ai_messages_reset", "")
        if reset_at:
            reset_dt = datetime.fromisoformat(reset_at.replace("Z", "+00:00"))
            now = datetime.now(timezone.utc)
            if now.year > reset_dt.year or now.month > reset_dt.month:
                get_user_db().table("profiles").update({
                    "ai_messages_used": 0,
                    "ai_messages_reset": now.isoformat()
                }).eq("id", user["id"]).execute()
                user["ai_messages_used"] = 0

        used = user.get("ai_messages_used", 0)
        allowed = limit["ai_msgs"]
        if used >= allowed:
            return False, f"You've used all {allowed} AI messages this month. Upgrade to get more."

    return True, ""


# ═══════════════════════════════════════════════════════════════
# DATA ANALYSIS ENGINE  (same as before, cleaned up)
# ═══════════════════════════════════════════════════════════════

# ── Column finder: normalises spaces, underscores, case ───────
def _col(df, *candidates):
    """Case-insensitive, space/underscore-insensitive column matcher."""
    normalised = {c.lower().replace(" ", "").replace("_", ""): c for c in df.columns}
    for cand in candidates:
        key = cand.lower().replace(" ", "").replace("_", "")
        if key in normalised:
            return normalised[key]
    return None

def _safe_float(val, default=0.0):
    try:
        f = float(val)
        return default if math.isnan(f) else f
    except (TypeError, ValueError):
        return default

def _clean_numeric(series):
    """
    Handle every real-world numeric format:
    - Currency with commas:  "£44,400"  "$2,500.00"  (Excel/Sheets quoted export)
    - Currency symbols only: £1500  $2500  €3100
    - Percentage strings:    29.8%  45%  100%
    - Plain numbers and already-numeric columns
    """
    if pd.api.types.is_numeric_dtype(series):
        return pd.to_numeric(series, errors="coerce").fillna(0)
    cleaned = (series.astype(str).str.strip()
               .str.replace(r"[£$€¥₹,\s]", "", regex=True)   # currency symbols + thousands commas
               .str.replace(r"%$", "", regex=True)              # trailing %
               .replace({"nan": "0", "None": "0", "": "0", "-": "0", "N/A": "0", "n/a": "0"}))
    return pd.to_numeric(cleaned, errors="coerce").fillna(0)

def _fix_decimal_percent(series, col_name):
    """
    Google Sheets and some ERP exports store percentages as 0–1 decimals.
    e.g. SellThrough = 0.298 instead of 29.8
    Detect and multiply by 100 automatically.
    """
    PCT_HINTS = {"sellthrough", "sellthroughpct", "st", "marginpct", "margin",
                 "returnrate", "fullpriceratio", "fpratio"}
    normalised = col_name.lower().replace(" ", "").replace("_", "")
    if not any(h in normalised for h in PCT_HINTS):
        return series
    nums = pd.to_numeric(series, errors="coerce").dropna()
    if len(nums) > 0 and (nums <= 1.5).all() and (nums > 0).any():
        return (series * 100).round(2)
    return series

def _best_sheet(file_obj, filename):
    """
    For multi-sheet Excel files, pick the sheet that has the most
    recognisable merchandising columns rather than blindly taking sheet 0.
    """
    KEY_COLS = {"sku", "productcode", "itemcode", "revenue", "sales",
                "unitssold", "units", "qty", "stockonhand", "stock"}
    xl = pd.ExcelFile(file_obj)
    best_sheet, best_score = xl.sheet_names[0], -1
    for name in xl.sheet_names:
        try:
            df = xl.parse(name, nrows=3)
            cols = {c.lower().replace(" ", "").replace("_", "") for c in df.columns}
            score = len(cols & KEY_COLS)
            if score > best_score:
                best_score, best_sheet = score, name
        except Exception:
            continue
    return xl.parse(best_sheet)

def _parse_numbers(file_obj):
    """Parse Apple Numbers .numbers file into a DataFrame."""
    import tempfile, os, warnings
    raw = file_obj.read() if hasattr(file_obj, "read") else file_obj
    with tempfile.NamedTemporaryFile(suffix=".numbers", delete=False) as tmp:
        tmp.write(raw)
        tmp_path = tmp.name
    try:
        # numbers_parser raises on Numbers file formats newer than the installed
        # library knows about (e.g. version 26.x from the latest macOS Numbers).
        # Catch that and give the user an actionable instruction rather than an
        # opaque 500 — exporting to CSV/Excel always works.
        try:
            from numbers_parser import Document
            with warnings.catch_warnings():
                warnings.simplefilter("ignore")  # ignore "unsupported version" warnings
                doc = Document(tmp_path)
        except Exception as e:
            raise ValueError(
                "This .numbers file was saved by a newer version of Apple Numbers "
                "than we can read yet. Please open it in Numbers and choose "
                "File → Export To → CSV (or Excel), then upload that file. "
                f"(parser detail: {type(e).__name__})"
            )

        # Find the sheet/table with the most recognisable columns
        KEY_COLS = {"sku","productcode","revenue","sales","unitssold","units","qty","stockonhand"}
        best_df, best_score = None, -1
        for sheet in doc.sheets:
            for table in sheet.tables:
                rows = list(table.iter_rows())
                if len(rows) < 2:
                    continue
                headers = [str(c.value) if c.value is not None else "" for c in rows[0]]
                norm = {h.lower().replace(" ","").replace("_","") for h in headers}
                score = len(norm & KEY_COLS)
                if score > best_score:
                    best_score = score
                    data = [[c.value for c in row] for row in rows[1:]]
                    best_df = pd.DataFrame(data, columns=headers)
        if best_df is None:
            raise ValueError(
                "No data table was found in this Numbers file. If the data is on a "
                "sheet, try File → Export To → CSV (or Excel) in Numbers and upload that."
            )
        return best_df
    finally:
        os.unlink(tmp_path)


def parse_upload(file_obj, filename):
    """
    Universal file parser — handles every format users actually export:
      Excel   .xlsx / .xls   (from Excel, Google Sheets, Numbers)
      Numbers .numbers        (Apple Numbers native format)
      ODS     .ods            (LibreOffice / OpenOffice)
      CSV     .csv            (any app, any delimiter, with/without BOM)
      TSV     .tsv / .txt     (tab-delimited from reporting tools)

    Robustness built in:
      - UTF-8 BOM stripped automatically (Apple Numbers CSV export)
      - Delimiter auto-detected (comma, semicolon, tab, pipe)
      - Empty leading/trailing rows and columns dropped
      - Multi-sheet Excel: picks the sheet with most data columns
    """
    ext = filename.rsplit(".", 1)[-1].lower()

    if ext == "numbers":
        df = _parse_numbers(file_obj)

    elif ext in ("xlsx", "xls"):
        # Read bytes for multi-sheet detection
        raw = file_obj.read() if hasattr(file_obj, "read") else file_obj
        df = _best_sheet(io.BytesIO(raw), filename)

    elif ext == "ods":
        raw = file_obj.read() if hasattr(file_obj, "read") else file_obj
        df = pd.read_excel(io.BytesIO(raw), engine="odf")

    else:
        # CSV / TSV / TXT — auto-detect delimiter + BOM
        raw = file_obj.read() if hasattr(file_obj, "read") else file_obj
        # Detect and strip BOM
        if raw.startswith(b"\xef\xbb\xbf"):
            raw = raw[3:]
        text = raw.decode("utf-8", errors="replace")

        # Note: currency symbols in Excel/Sheets exports are always quoted ("£44,400")
        # so pandas correctly parses them as strings. _clean_numeric strips £,$,€
        # and commas from those strings after parsing.
        try:
            df = pd.read_csv(
                io.StringIO(text),
                sep=None,            # auto-detect: comma, semicolon, tab, pipe
                engine="python",
                skip_blank_lines=True,
                encoding_errors="replace",
            )
        except Exception:
            # Sniffer can fail on unusual files — fall back to comma delimiter
            df = pd.read_csv(
                io.StringIO(text),
                sep=",",
                skip_blank_lines=True,
                encoding_errors="replace",
            )

    # ── Post-load cleanup ──────────────────────────────────────
    # Strip column name whitespace
    df.columns = [str(c).strip() for c in df.columns]
    # Drop entirely-empty columns and rows (Google Sheets ghost cells)
    df = df.dropna(axis=1, how="all").dropna(axis=0, how="all")
    # Drop unnamed columns (Unnamed: 3, Unnamed: 4 …)
    df = df[[c for c in df.columns if not str(c).startswith("Unnamed:")]]
    # Reset index after row drops
    df = df.reset_index(drop=True)

    return df

def analyse(df):
    # ── Expanded column detection — covers all app export formats ─
    col_sku   = _col(df, "SKU","sku","ProductCode","product_code","ItemCode","item_code",
                         "Item Code","Product Code","Style","StyleCode","style_code")
    col_units = _col(df, "UnitsSold","units_sold","Units","Qty","Quantity","qty",
                         "Units Sold","Sales Units","SalesUnits","Pieces","pieces")
    col_rev   = _col(df, "Revenue","revenue","Sales","SalesValue","sales_value",
                         "Total Revenue","TotalRevenue","Net Sales","NetSales",
                         "Total Sales","TotalSales","Turnover","turnover")
    col_stock = _col(df, "StockOnHand","stock_on_hand","Stock","OnHand","on_hand",
                         "Stock On Hand","CurrentStock","current_stock","Inventory","inventory",
                         "ClosingStock","closing_stock","Units On Hand")
    col_cost  = _col(df, "CostPrice","cost_price","Cost","UnitCost","unit_cost",
                         "Cost Price","Buying Price","buying_price","COGS")
    col_margin= _col(df, "MarginPct","margin_pct","Margin","GrossMarginPct",
                         "gross_margin_pct","GM%","GM Pct","Gross Margin","GrossMargin",
                         "Margin %","Margin%","GP%")
    # Absolute margin value (pounds) — used to DERIVE margin % when the file
    # ships a "Margin £" column instead of a percentage.
    col_margin_val = _col(df, "Margin £","Margin GBP","MarginGBP","Margin Value","MarginValue",
                              "Gross Margin £","GM £","GMGBP","Gross Profit","GrossProfit",
                              "Profit £","Profit GBP","ProfitGBP","Profit")
    # Retail / selling price — used to derive margin % from (retail - cost)/retail.
    col_retail = _col(df, "Retail Price","RetailPrice","RRP","Selling Price","SellingPrice",
                          "Unit Price","UnitPrice","Ticket Price","TicketPrice","Price","SRP")
    col_st    = _col(df, "SellThrough","sell_through","SellThroughPct","ST","S/T",
                         "Sell Through","Sell Through %","SellThru","sell_thru","STR")
    col_cover = _col(df, "StockCoverWeeks","stock_cover_weeks","CoverWeeks","cover_weeks",
                         "Weeks Cover","WeeksCover","Cover","WOS","wos","Weeks of Stock")
    col_cat   = _col(df, "Category","category","Department","Dept","dept",
                         "Product Type","ProductType","Division","Class","Subclass")
    col_chan  = _col(df, "Channel","channel","SalesChannel","sales_channel",
                         "Store","Region","Location","Market")
    col_week  = _col(df, "Week","week","WeekNumber","week_number","Period","period",
                         "Week No","WeekNo","Trading Week","Wk")
    col_name  = _col(df, "ProductName","product_name","Name","Description","Desc",
                         "Product Name","Product Description","Style Name","StyleName")
    col_brand = _col(df, "Brand","brand","BrandName","brand_name","Label","Supplier")
    col_buy   = _col(df, "BuyQty","buy_qty","OpeningStock","opening_stock","InitialStock",
                         "initial_stock","IntakeQty","intake_qty","Opening Stock","Buy Qty")
    col_fp    = _col(df, "FullPriceSales","full_price_sales","FPSales","fp_sales",
                         "Full Price Revenue","FullPriceRevenue","FP Revenue")

    if not all([col_sku, col_units, col_rev]):
        missing = [n for c,n in [(col_sku,"SKU"),(col_units,"UnitsSold"),(col_rev,"Revenue")] if not c]
        found = list(df.columns[:10])
        raise ValueError(
            f"Missing required columns: {', '.join(missing)}. "
            f"Found columns: {', '.join(str(c) for c in found)}{'...' if len(df.columns)>10 else ''}. "
            f"Skuvvy needs at least SKU, UnitsSold/Units/Qty, and Revenue/Sales columns."
        )

    # ── Clean numerics — strips £$€, commas, % signs ──────────
    for col in [col_units, col_rev, col_stock, col_cost, col_fp, col_buy,
                col_margin_val, col_retail]:
        if col:
            df[col] = _clean_numeric(df[col])

    # ── Fix decimal percentages (Google Sheets / ERP exports) ──
    for col, name in [(col_margin, "MarginPct"), (col_st, "SellThrough"), (col_cover, "StockCoverWeeks")]:
        if col:
            df[col] = _clean_numeric(df[col])
            df[col] = _fix_decimal_percent(df[col], name)

    # ── Derive a margin % when the file has no explicit % column ──
    # Retail exports frequently ship an absolute "Margin £" column, or just
    # Cost/Retail prices, instead of a margin percentage. Derive the % so the
    # KPI is real rather than defaulting to 0. Order of preference:
    #   1) explicit margin % column
    #   2) Margin £ ÷ Revenue × 100
    #   3) (Retail − Cost) ÷ Retail × 100
    margin_source = "explicit" if col_margin else None
    if not col_margin:
        if col_margin_val and col_rev:
            rev_nz = df[col_rev].where(df[col_rev] != 0)        # avoid /0 → NaN
            df["_derived_margin_pct"] = (df[col_margin_val] / rev_nz * 100)
            col_margin, margin_source = "_derived_margin_pct", "margin_value_over_revenue"
        elif col_retail and col_cost:
            rp_nz = df[col_retail].where(df[col_retail] != 0)
            df["_derived_margin_pct"] = ((df[col_retail] - df[col_cost]) / rp_nz * 100)
            col_margin, margin_source = "_derived_margin_pct", "cost_and_retail"

    agg = {col_units: "sum", col_rev: "sum"}
    for c in [col_stock, col_cost, col_margin, col_st, col_cover,
              col_cat, col_chan, col_name, col_brand, col_fp, col_buy, col_margin_val]:
        if c:
            agg[c] = "last"  if c in [col_stock, col_st, col_cover] else \
                     "first" if c in [col_cat, col_chan, col_name, col_brand, col_buy] else \
                     "mean"  if c == col_margin else "sum"

    sku_df = df.groupby(col_sku, as_index=False).agg(agg)
    total_rev   = _safe_float(sku_df[col_rev].sum())
    total_units = _safe_float(sku_df[col_units].sum())
    sku_count   = len(sku_df)
    avg_st      = _safe_float(sku_df[col_st].mean()) if col_st else 0.0
    avg_cover   = _safe_float(sku_df[col_cover].mean()) if col_cover else 0.0

    # ── Margin: a simple per-SKU average AND a revenue-weighted blended % ──
    # avg_margin       = unweighted mean across the range (every SKU counts equally)
    # gross_margin_pct = blended Σmargin ÷ Σrevenue (big sellers count more) — the
    #                    "true" business gross margin. These legitimately differ.
    margin_available = bool(col_margin) or bool(col_margin_val)
    avg_margin = _safe_float(sku_df[col_margin].mean()) if col_margin else 0.0
    if col_margin_val and total_rev:
        total_margin_val = _safe_float(sku_df[col_margin_val].sum())
        gross_margin_pct = total_margin_val / total_rev * 100
        gross_profit     = total_margin_val
    elif col_margin and total_rev:
        wsum = _safe_float((sku_df[col_margin] * sku_df[col_rev]).sum())
        gross_margin_pct = wsum / total_rev
        gross_profit     = total_rev * gross_margin_pct / 100
    else:
        gross_margin_pct = 0.0
        gross_profit     = total_rev * avg_margin / 100

    # Best sellers
    best = sku_df.sort_values(col_rev, ascending=False).head(20)
    best_sellers = []
    for _, r in best.iterrows():
        e = {"sku": str(r[col_sku]), "revenue": round(_safe_float(r[col_rev]), 2),
             "units": round(_safe_float(r[col_units]))}
        if col_name:   e["name"]         = str(r[col_name])
        if col_cat:    e["category"]     = str(r[col_cat])
        if col_st:     e["sell_through"] = round(_safe_float(r[col_st]), 1)
        if col_margin: e["margin_pct"]   = round(_safe_float(r[col_margin]), 1)
        if col_stock:  e["stock"]        = round(_safe_float(r[col_stock]))
        best_sellers.append(e)

    # Reorder alerts
    reorders = []
    sev_map = {"CRITICAL": 0, "URGENT": 1, "REORDER": 2}
    for _, r in sku_df.iterrows():
        stock = _safe_float(r[col_stock]) if col_stock else None
        cover = _safe_float(r[col_cover]) if col_cover else None
        st    = _safe_float(r[col_st]) if col_st else None
        if stock == 0 or (cover is not None and cover <= 4) or (st is not None and st >= 70 and (stock is None or stock < 20)):
            urgency = "CRITICAL" if (stock == 0 or (cover is not None and cover <= 1)) else \
                      "URGENT"   if (cover is not None and cover <= 2) else "REORDER"
            e = {"sku": str(r[col_sku]), "urgency": urgency,
                 "revenue": round(_safe_float(r[col_rev]), 2)}
            if col_name:  e["name"]         = str(r[col_name])
            if col_cat:   e["category"]     = str(r[col_cat])
            if col_stock: e["stock"]        = round(stock) if stock is not None else None
            if col_cover: e["cover_weeks"]  = round(cover, 1) if cover is not None else None
            if col_st:    e["sell_through"] = round(st, 1) if st is not None else None
            reorders.append(e)
    reorders.sort(key=lambda x: (sev_map[x["urgency"]], -x["revenue"]))

    # Markdown risk
    markdowns = []
    sev_md = {"CLEARANCE": 0, "DEEP": 1, "TACTICAL": 2}
    if col_st:
        for _, r in sku_df.iterrows():
            st = _safe_float(r[col_st])
            stock = _safe_float(r[col_stock]) if col_stock else None
            if st < 30 and (stock is None or stock > 0):
                sev = "CLEARANCE" if st < 15 else "DEEP" if st < 22 else "TACTICAL"
                depth = "40–50%" if st < 15 else "25–35%" if st < 22 else "15–20%"
                e = {"sku": str(r[col_sku]), "severity": sev, "sell_through": round(st, 1),
                     "recommended_depth": depth, "revenue": round(_safe_float(r[col_rev]), 2)}
                if col_name: e["name"]     = str(r[col_name])
                if col_cat:  e["category"] = str(r[col_cat])
                markdowns.append(e)
    markdowns.sort(key=lambda x: (sev_md[x["severity"]], -x["revenue"]))

    # Category scorecard
    categories = []
    if col_cat:
        for cat, grp in sku_df.groupby(col_cat):
            cat_rev = _safe_float(grp[col_rev].sum())
            cat_st  = _safe_float(grp[col_st].mean()) if col_st else None
            cat_mg  = _safe_float(grp[col_margin].mean()) if col_margin else None
            health  = ("STRONG" if cat_st >= 40 else "HEALTHY" if cat_st >= 25
                       else "AT RISK" if cat_st >= 15 else "CRITICAL") if cat_st else "UNKNOWN"
            e = {"category": str(cat), "revenue": round(cat_rev, 2),
                 "units": round(_safe_float(grp[col_units].sum())), "sku_count": len(grp), "health": health}
            if cat_st is not None: e["avg_sell_through"] = round(cat_st, 1)
            if cat_mg is not None: e["avg_margin_pct"]   = round(cat_mg, 1)
            categories.append(e)
        categories.sort(key=lambda x: -x["revenue"])

    # Weekly trend
    weekly = []
    if col_week:
        wg = df.groupby(col_week)[col_rev].sum().reset_index().sort_values(col_week)
        weekly = [{"week": str(r[col_week]), "revenue": round(_safe_float(r[col_rev]), 2)} for _, r in wg.iterrows()]

    # Channel split
    channels = []
    if col_chan:
        cg = df.groupby(col_chan)[col_rev].sum().reset_index().sort_values(col_rev, ascending=False)
        channels = [{"channel": str(r[col_chan]),
                     "revenue": round(_safe_float(r[col_rev]), 2),
                     "share_pct": round(_safe_float(r[col_rev]) / total_rev * 100, 1) if total_rev else 0}
                    for _, r in cg.iterrows()]

    return {
        "meta": {"sku_count": sku_count, "analysed_at": datetime.utcnow().isoformat()},
        "kpis": {
            "total_revenue": round(total_rev, 2),
            "total_units": round(total_units),
            "gross_profit": round(gross_profit, 2),
            "avg_sell_through": round(avg_st, 1),
            "avg_margin_pct": round(avg_margin, 1),
            "gross_margin_pct": round(gross_margin_pct, 1),
            "margin_available": margin_available,
            "margin_source": margin_source,
            "avg_cover_weeks": round(avg_cover, 1),
            "reorder_count": len(reorders),
            "markdown_risk_count": len(markdowns),
            "critical_oos_count": sum(1 for r in reorders if r["urgency"] == "CRITICAL"),
        },
        "weekly_trend": weekly,
        "best_sellers": best_sellers,
        "reorder_alerts": reorders[:50],
        "markdown_risk": markdowns[:50],
        "category_scorecard": categories,
        "channel_split": channels,
    }


def build_system_prompt(analysis):
    k = analysis.get("kpis", {})
    return f"""You are Skuvvy Copilot — a senior fashion merchandiser and buying director sitting
beside the user, analysing their trading data with them. Be warm, sharp and commercial.

TRADING DATA
============
Revenue: £{k.get('total_revenue',0):,.0f} | Units: {k.get('total_units',0):,.0f}
Gross Profit: £{k.get('gross_profit',0):,.0f} | Sell-Through: {k.get('avg_sell_through',0):.1f}%
Gross Margin: {k.get('avg_margin_pct',0):.1f}% | Stock Cover: {k.get('avg_cover_weeks',0):.1f} wks
SKUs: {analysis.get('meta',{}).get('sku_count',0)} | Reorder alerts: {k.get('reorder_count',0)} | Markdown risk: {k.get('markdown_risk_count',0)}

TOP REORDERS: {json.dumps(analysis.get('reorder_alerts',[])[:8])}
TOP MARKDOWN RISK: {json.dumps(analysis.get('markdown_risk',[])[:8])}
BEST SELLERS: {json.dumps(analysis.get('best_sellers',[])[:8])}
CATEGORIES: {json.dumps(analysis.get('category_scorecard',[]))}

Answer like a buying director — specific SKUs, real numbers, clear actions. Never generic.

RESPONSE FORMAT (follow exactly)
================================
Write in clean, professional British English with proper sentences and grammar.
- Open with ONE short plain-sentence summary (1–2 sentences max). No heading on it.
- If there are actions, risks, or items to list, present them as bullet points — one
  idea per bullet, each starting with "- " (a hyphen and a space).
- Keep each bullet to a single concise line where possible.
- Use **bold** sparingly, only for a SKU code or a key number that deserves emphasis.
- Do NOT use markdown headings (no "#", "##", "###"). Do NOT use emoji.
- Do NOT wrap things in tables or code blocks unless the user explicitly asks for a
  table or code. Default to a brief intro plus bullets.
- Avoid decorative symbols, asterisk bullets, or repeated punctuation.

DIAGNOSTIC ANSWERS (when the user asks "why", "what's happening", "what should I do",
"should I markdown/reorder", or "what needs my attention")
==========================================================
Reason like a merchandiser by walking through cause to action. Use this four-part shape,
each as its own bullet led by a bold label (no headings, keep it tight):
- **Problem:** what happened, with the real number (e.g. Dresses revenue fell 15%).
- **Reason:** why it happened, naming specific SKUs/categories from the data.
- **Action:** the concrete commercial decision to take now.
- **Impact:** the expected outcome in £ or points where you can estimate it.
Only use this shape for diagnostic/decision questions. For simple factual look-ups, just
answer with the brief intro plus bullets. Never invent numbers — if the data doesn't
support an estimate, say what you'd need to confirm it."""


# ═══════════════════════════════════════════════════════════════
# ROLE-BASED ANALYSIS HELPERS
# ═══════════════════════════════════════════════════════════════

def _compute_health_score(avg_st, avg_margin, avg_cover, reorder_count):
    """Compute business health score 0-100."""
    # sell-through component
    if avg_st >= 50:
        st_pts = 25
    elif avg_st >= 35:
        st_pts = 20
    elif avg_st >= 20:
        st_pts = 12
    else:
        st_pts = 5

    # margin component
    if avg_margin >= 55:
        mg_pts = 25
    elif avg_margin >= 40:
        mg_pts = 20
    elif avg_margin >= 25:
        mg_pts = 12
    else:
        mg_pts = 5

    # cover component
    if 4 <= avg_cover <= 8:
        cov_pts = 25
    elif (2 <= avg_cover < 4) or (8 < avg_cover <= 12):
        cov_pts = 18
    elif 1 <= avg_cover < 2:
        cov_pts = 10
    else:
        cov_pts = 3

    # reorder component
    if reorder_count == 0:
        ro_pts = 25
    elif reorder_count <= 3:
        ro_pts = 20
    elif reorder_count <= 10:
        ro_pts = 12
    else:
        ro_pts = 5

    return int(st_pts + mg_pts + cov_pts + ro_pts)


def _gbp(v) -> str:
    """Compact GBP formatter for reasoning chips."""
    try:
        v = float(v or 0)
    except (TypeError, ValueError):
        return "£0"
    a = abs(v)
    if a >= 1_000_000: return f"£{v/1_000_000:.1f}m"
    if a >= 1_000:     return f"£{v/1_000:.0f}k"
    return f"£{v:,.0f}"


def _action_reasoning(a, ctx):
    """
    Build the 'Why Skuvvy recommends this' block for one ranked action.

    Returns a dict with:
      data_points    — the evidence behind the decision (sell-through, category
                       average, weeks of cover, stock value tied up, severity)
      expected_impact— revenue and margin-point effect
      confidence     — model confidence (0-100)
      why            — one-sentence plain-English rationale
    Every figure is drawn from the dataset; nothing is invented.
    """
    sku = a.get("sku", "")
    typ = a.get("type", "")
    md  = ctx["md_by"].get(sku)
    ro  = ctx["ro_by"].get(sku)
    cat = ctx["cat_by"].get(sku)          # when the action targets a category
    bs  = ctx["seller_by"].get(sku)

    item_cat = (md or ro or bs or {}).get("category")
    cat_obj  = ctx["cat_by"].get(item_cat) if item_cat else None
    cat_avg_st = (cat_obj or {}).get("avg_sell_through")
    avg_st   = ctx["avg_st"]
    avg_cover = ctx["avg_cover"]

    dp = []
    rev = a.get("revenue_impact", 0)

    if md:
        st = md.get("sell_through", 0)
        dp.append({"label": "Sell-through", "value": f"{st:.0f}%"})
        if cat_avg_st is not None:
            dp.append({"label": f"{item_cat} avg ST", "value": f"{cat_avg_st:.0f}%"})
        else:
            dp.append({"label": "Range avg ST", "value": f"{avg_st:.0f}%"})
        dp.append({"label": "Stock value tied up", "value": _gbp(md.get("revenue", 0))})
        dp.append({"label": "Recommended depth", "value": md.get("recommended_depth", "25–35%")})
        why = (f"{sku} is selling at just {st:.0f}% — "
               + (f"well below the {item_cat} category average of {cat_avg_st:.0f}%" if cat_avg_st is not None
                  else f"below the {avg_st:.0f}% range average")
               + f". Acting now releases {_gbp(md.get('revenue', 0))} of trapped stock before it ages further.")
    elif ro:
        cover = ro.get("cover_weeks")
        st    = ro.get("sell_through")
        if cover is not None:
            dp.append({"label": "Weeks of cover", "value": f"{cover:.1f} wks"})
        if st is not None:
            dp.append({"label": "Sell-through", "value": f"{st:.0f}%"})
        if ro.get("stock") is not None:
            dp.append({"label": "Stock on hand", "value": f"{ro.get('stock'):,} units"})
        dp.append({"label": "Revenue at risk", "value": _gbp(ro.get("revenue", 0))})
        dp.append({"label": "Urgency", "value": ro.get("urgency", "REORDER").title()})
        cover_txt = f"only {cover:.1f} weeks of cover left" if cover is not None else "cover running low"
        why = (f"{sku} has {cover_txt}"
               + (f" against {st:.0f}% sell-through" if st is not None else "")
               + f", putting {_gbp(ro.get('revenue', 0))} of demand at risk if it stocks out before the next delivery.")
    elif cat:
        dp.append({"label": "Category revenue", "value": _gbp(cat.get("revenue", 0))})
        if cat.get("avg_sell_through") is not None:
            dp.append({"label": "Avg sell-through", "value": f"{cat.get('avg_sell_through'):.0f}%"})
        dp.append({"label": "Health", "value": cat.get("health", "—").title()})
        dp.append({"label": "SKUs in range", "value": str(cat.get("sku_count", 0))})
        cst = cat.get("avg_sell_through")
        if cst is not None and cst < avg_st:
            why = (f"{sku} is rated {cat.get('health','').title()} with {cst:.0f}% sell-through, "
                   f"below the {avg_st:.0f}% range average — trimming forward commitment protects margin.")
        else:
            why = (f"{sku} is your strongest category at {_gbp(cat.get('revenue', 0))} revenue and "
                   f"{(cst or 0):.0f}% sell-through — leaning in captures demand outpacing supply.")
    else:
        # Generic / monitoring / OTB direction
        dp.append({"label": "Range avg ST", "value": f"{avg_st:.0f}%"})
        dp.append({"label": "Range avg cover", "value": f"{avg_cover:.1f} wks"})
        if bs:
            dp.append({"label": f"{sku} revenue", "value": _gbp(bs.get("revenue", 0))})
        why = a.get("description", "Recommended from the current trading position across the range.")

    mi = a.get("margin_impact", 0)
    impact = {
        "revenue": rev,
        "revenue_label": ("Revenue recovered" if rev >= 0 else "Revenue traded"),
        "margin_pts": mi,
    }
    return {"data_points": dp, "expected_impact": impact,
            "confidence": a.get("confidence", 70), "why": why}


def _category_health(cat: dict) -> int:
    """Numeric 0-100 health score for a single category (ST + margin)."""
    st = float(cat.get("avg_sell_through") or 0)
    mg = float(cat.get("avg_margin_pct") or 0)
    if not st and not mg:
        return 0
    score = (min(st, 100) / 100) * 55 + (min(mg, 60) / 60) * 45
    return int(round(max(0, min(100, score))))


def _category_cards(analysis) -> list:
    """Compact category performance cards for the Command Centre grid."""
    cats = analysis.get("category_scorecard", []) or []
    cards = []
    for c in cats:
        if not c or float(c.get("revenue") or 0) <= 0:
            continue
        cards.append({
            "category":     c.get("category", "Other"),
            "revenue":      c.get("revenue", 0),
            "units":        c.get("units", 0),
            "sku_count":    c.get("sku_count", 0),
            "sell_through": c.get("avg_sell_through"),
            "margin":       c.get("avg_margin_pct"),
            "health":       _category_health(c),
            "health_label": c.get("health", "—"),
        })
    cards.sort(key=lambda x: -(x.get("revenue") or 0))
    return cards


def _category_detail(analysis, category: str) -> dict:
    """
    Drill-down for one category: headline KPIs, top sellers, slow movers,
    and recommended actions — all derived from the dataset.
    """
    cats = analysis.get("category_scorecard", []) or []
    cat = next((c for c in cats if (c.get("category") or "").lower() == (category or "").lower()), None)
    if not cat:
        return {}

    clow = (category or "").lower()
    best = analysis.get("best_sellers", []) or []
    md   = analysis.get("markdown_risk", []) or []
    members = [s for s in best if (s.get("category") or "").lower() == clow]

    # Top sellers: highest-revenue lines in this category.
    top_sellers = sorted(members, key=lambda x: -(x.get("revenue") or 0))[:5]
    top_skus = {t.get("sku") for t in top_sellers}

    # Look-up so we can enrich markdown rows with units / margin / stock
    # that live on the best-seller records.
    best_by_sku = {s.get("sku"): s for s in best}

    # Slow movers: genuinely low sell-through lines (the markdown-risk list is the
    # canonical "slow" set, st < 30). A SKU already shown as a top seller is NEVER
    # repeated here — top sellers rank on revenue, slow movers on demand velocity,
    # so the same SKU must not appear in both lists.
    md_members = [s for s in md if (s.get("category") or "").lower() == clow
                  and s.get("sku") not in top_skus]
    if md_members:
        # Merge in best-seller fields where available so the card shows full data.
        slow_src = []
        for s in md_members:
            base = dict(best_by_sku.get(s.get("sku"), {}))
            base.update({kk: vv for kk, vv in s.items() if vv is not None})
            slow_src.append(base)
    else:
        # Fallback: lowest sell-through members not already listed as top sellers.
        slow_src = [s for s in members
                    if s.get("sku") not in top_skus and s.get("sell_through") is not None]
    slow_movers = sorted(
        slow_src,
        key=lambda x: (x.get("sell_through") if x.get("sell_through") is not None else 999)
    )[:5]

    def slim(s):
        return {
            "sku":          s.get("sku"),
            "name":         s.get("name") or s.get("sku"),
            "revenue":      s.get("revenue", 0),
            "units":        s.get("units", 0),
            "sell_through": s.get("sell_through"),
            "margin_pct":   s.get("margin_pct"),
            "stock":        s.get("stock"),
        }

    st  = cat.get("avg_sell_through")
    mg  = cat.get("avg_margin_pct")
    rev = cat.get("revenue", 0)
    health_lbl = cat.get("health", "—")
    cat_name = cat.get("category", category)

    # Recommended actions for the category
    actions = []
    if st is not None and st < 30:
        actions.append({
            "type": "Markdown",
            "text": f"Clear ageing {cat_name} stock — {st:.0f}% sell-through signals weak demand; a 20–30% markdown releases working capital.",
        })
    if st is not None and st >= 45:
        actions.append({
            "type": "Reorder",
            "text": f"Lean into {cat_name} — {st:.0f}% sell-through is outpacing supply; increase buying depth before stockouts.",
        })
    if mg is not None and mg < 40:
        actions.append({
            "type": "Margin",
            "text": f"Protect {cat_name} margin — at {mg:.0f}% it sits below the 40% benchmark; review cost prices and discount depth.",
        })
    if slow_movers and (slow_movers[0].get("sell_through") or 100) < 25:
        sm = slow_movers[0]
        actions.append({
            "type": "Review",
            "text": f"{sm.get('name') or sm.get('sku')} is the weakest line at {(sm.get('sell_through') or 0):.0f}% — consider markdown or delisting.",
        })
    if not actions:
        actions.append({
            "type": "Monitor",
            "text": f"{cat_name} is trading in line with the range — monitor weekly and hold current strategy.",
        })

    return {
        "category":     cat.get("category"),
        "revenue":      rev,
        "units":        cat.get("units", 0),
        "sku_count":    cat.get("sku_count", 0),
        "sell_through": st,
        "margin":       mg,
        "health":       _category_health(cat),
        "health_label": health_lbl,
        "top_sellers":  [slim(s) for s in top_sellers],
        "slow_movers":  [slim(s) for s in slow_movers],
        "actions":      actions,
    }


def _trading_meeting(analysis, role, kpis, health, health_detail, label, prev_metrics=None):
    """
    Monday Trading Review pack — the senior-merchandiser meeting view.

    Five sections: business overview, AI executive summary, what changed this
    week, top business questions, and the director summary export. Every figure
    is dataset-driven; week-on-week deltas use the previous upload when present.
    """
    cats = analysis.get("category_scorecard", []) or []
    weekly = analysis.get("weekly_trend", []) or []
    rev   = float(kpis.get("total_revenue") or 0)
    mg    = float(kpis.get("avg_margin_pct") or 0)
    st    = float(kpis.get("avg_sell_through") or 0)
    md_ct = int(kpis.get("markdown_risk_count") or 0)
    cover = float(kpis.get("avg_cover_weeks") or 0)
    margin_ok = kpis.get("margin_available") is not False
    top_cat = cats[0] if cats else None

    # Period label from the weekly trend, else a generic label
    if weekly:
        first, last = weekly[0].get("week"), weekly[-1].get("week")
        period = f"{first} – {last}" if first != last else str(first)
    else:
        period = "Latest trading period"

    # Deltas vs previous upload (week-on-week)
    def pct_delta(cur, prev):
        if not prev:
            return None
        return (cur - prev) / prev * 100

    rev_delta = mg_delta = st_delta = health_delta = None
    if prev_metrics:
        rev_delta    = pct_delta(rev, prev_metrics.get("revenue"))
        # Only show a margin / sell-through movement when there is a real prior
        # baseline (>0). A missing or zero previous value would otherwise produce a
        # misleading swing like "+46.2% to 0%".
        _mg_prev = prev_metrics.get("margin")
        _st_prev = prev_metrics.get("sell_through")
        mg_delta     = (mg - _mg_prev) if _mg_prev else None
        st_delta     = (st - _st_prev) if _st_prev else None
        health_delta = health - prev_metrics.get("health", health)

    def fmt_pct(v):
        # Percentage metrics (revenue %, margin, sell-through) are always shown with
        # a "%" suffix so units stay consistent across the whole app — never "pts".
        if v is None:
            return None
        return f"{'+' if v >= 0 else ''}{v:.1f}%"

    overview = [
        {"label": "Revenue",      "value": _gbp(rev),
         "delta": fmt_pct(rev_delta), "trend": ("up" if (rev_delta or 0) >= 0 else "down") if rev_delta is not None else "flat",
         "sub": "vs last week" if rev_delta is not None else "this period"},
        {"label": "Margin",       "value": (f"{mg:.1f}%" if margin_ok else "Insufficient data"),
         "delta": fmt_pct(mg_delta) if margin_ok else None,
         "trend": ("up" if (mg_delta or 0) >= 0 else "down") if (mg_delta is not None and margin_ok) else "flat",
         "sub": "blended margin"},
        {"label": "Sell-Through", "value": f"{st:.0f}%",
         "delta": fmt_pct(st_delta), "trend": ("up" if (st_delta or 0) >= 0 else "down") if st_delta is not None else "flat",
         "sub": "WoW" if st_delta is not None else "range average"},
        {"label": "Health Score", "value": f"{health}/100",
         "delta": (f"{'+' if (health_delta or 0) >= 0 else ''}{int(health_delta)} pts" if health_delta is not None else None),
         "trend": ("up" if (health_delta or 0) >= 0 else "down") if health_delta is not None else "flat",
         "sub": health_detail.get("label", "")},
    ]

    # What changed this week
    changes = []
    changes.append({
        "title": "Revenue Movement",
        "value": fmt_pct(rev_delta) if rev_delta is not None else _gbp(rev),
        "trend": ("up" if (rev_delta or 0) >= 0 else "down") if rev_delta is not None else "flat",
        "driver": (f"Driven by {top_cat['category']} (£{top_cat['revenue']:,.0f})" if top_cat
                   else "Across the range"),
    })
    if margin_ok:
        md_cause = ("Markdown dependency on slow movers" if md_ct > 3
                    else "Stable pricing across the range")
        changes.append({
            "title": "Margin Movement",
            "value": fmt_pct(mg_delta) if mg_delta is not None else f"{mg:.1f}%",
            "trend": ("up" if (mg_delta or 0) >= 0 else "down") if mg_delta is not None else "flat",
            "driver": md_cause,
        })
    else:
        changes.append({
            "title": "Margin Movement", "value": "Insufficient data", "trend": "flat",
            "driver": "This dataset has no margin column — add one to track margin movement.",
        })
    changes.append({
        "title": "Stock Position",
        "value": (f"{md_ct} SKUs flagged" if md_ct else "Clean"),
        "trend": "down" if md_ct > 5 else "flat",
        "driver": (f"Slow/overstocked lines at ~{cover:.0f} wks average cover" if md_ct
                   else f"Healthy cover at ~{cover:.0f} wks across the range"),
    })

    # Top business questions — clickable prompts that open the AI assistant
    questions = [
        "Why are we behind target?",
        "Which categories need action?",
        "Where are we losing margin?",
        "Which products should we reorder?",
        "Which products should we markdown?",
    ]

    return {
        "header": {
            "title":     "Monday Trading Review",
            "dataset":   label,
            "period":    period,
            "generated": datetime.utcnow().strftime("%d %b %Y"),
        },
        "overview":      overview,
        "exec_summary":  role.get("executive_summary", ""),
        "changes":       changes,
        "questions":     questions,
    }


def _build_role_analysis(analysis):
    """Build extended role-based analysis data from base analysis."""
    k = analysis.get("kpis", {})
    best_sellers    = analysis.get("best_sellers", [])
    reorders        = analysis.get("reorder_alerts", [])
    markdowns       = analysis.get("markdown_risk", [])
    categories      = analysis.get("category_scorecard", [])
    weekly_trend    = analysis.get("weekly_trend", [])

    total_rev    = k.get("total_revenue", 0)
    total_units  = k.get("total_units", 0)
    avg_st       = k.get("avg_sell_through", 0)
    avg_margin   = k.get("avg_margin_pct", 0)
    avg_cover    = k.get("avg_cover_weeks", 0)
    reorder_count = k.get("reorder_count", 0)
    sku_count    = analysis.get("meta", {}).get("sku_count", 0)

    # Bottom 10 sellers (ascending by units)
    bottom_sellers = sorted(best_sellers, key=lambda x: x.get("units", 0))[:10]

    # Stock metrics
    total_stock_units = sum(s.get("stock", 0) for s in best_sellers if s.get("stock") is not None)
    # Estimate stock value: if we have cost data use it, else use revenue * 0.5 / 12
    total_stock_value = round(total_rev * 0.5 / 12, 2) if total_rev else 0

    # Missing data count
    missing_data_count = sum(
        1 for s in best_sellers
        if not s.get("stock") or not s.get("margin_pct") or not s.get("sell_through")
    )
    data_quality_pct = round((sku_count - missing_data_count) / max(sku_count, 1) * 100, 1)

    # Open PO metrics
    critical_urgent = [r for r in reorders if r.get("urgency") in ("CRITICAL", "URGENT")]
    open_po_count = len(critical_urgent)
    open_po_value = round(sum(r.get("revenue", 0) * 0.15 for r in critical_urgent), 2)

    # Health score (explainable breakdown — total is the canonical score)
    health_detail = _health_breakdown(k)
    health_score  = health_detail["total"]

    # Plan data (actuals * multipliers)
    plan_data = {
        "revenue_plan": round(total_rev * 1.05, 2),
        "margin_plan":  round(avg_margin * 1.08, 1),
        "stock_plan":   round(total_stock_units * 0.92, 0),
    }

    # Vs plan
    vs_plan = {
        "revenue_pct": round(total_rev / plan_data["revenue_plan"] * 100, 1) if plan_data["revenue_plan"] else 0,
        "margin_pct":  round(avg_margin / plan_data["margin_plan"] * 100, 1) if plan_data["margin_plan"] else 0,
        "stock_pct":   round(total_stock_units / max(plan_data["stock_plan"], 1) * 100, 1),
    }

    cash_trapped = total_stock_value

    # Weekly sparkline — last 8 revenue values from weekly_trend, or simulated
    if weekly_trend:
        sparkline = [w["revenue"] for w in weekly_trend[-8:]]
        while len(sparkline) < 8:
            sparkline.insert(0, sparkline[0] * 0.95 if sparkline else total_rev / 52)
    else:
        base = total_rev / 52
        sparkline = [round(base * (0.85 + 0.03 * i + (0.05 if i % 3 == 0 else 0)), 2) for i in range(8)]

    # Top risks (generated from data)
    top_risks = []
    # Risk 1: markdown risk
    if markdowns:
        worst_md = markdowns[0]
        top_risks.append(
            f"{len(markdowns)} SKUs at markdown risk — {worst_md['sku']} has only "
            f"{worst_md.get('sell_through', 0):.0f}% sell-through, risking £{worst_md.get('revenue', 0):,.0f} revenue"
        )
    else:
        top_risks.append("Monitor sell-through weekly — low ST can quickly escalate to markdown territory")

    # Risk 2: stock cover
    if avg_cover > 12:
        top_risks.append(
            f"Average stock cover of {avg_cover:.1f} weeks is above target — overstocked range risks cash and margin"
        )
    elif avg_cover < 3:
        top_risks.append(
            f"Average stock cover of {avg_cover:.1f} weeks is critically low — risk of lost sales across the range"
        )
    else:
        if reorder_count > 5:
            top_risks.append(
                f"{reorder_count} SKUs need urgent reordering — risk of stockouts impacting revenue momentum"
            )
        else:
            top_risks.append("Margin erosion risk if markdown depth increases — review pricing strategy")

    # Risk 3: weak categories
    weak_cats = [c for c in categories if c.get("health") in ("CRITICAL", "AT RISK")]
    if weak_cats:
        top_risks.append(
            f"{len(weak_cats)} categories rated AT RISK or CRITICAL — "
            f"{weak_cats[0]['category']} needs immediate attention with £{weak_cats[0]['revenue']:,.0f} revenue exposed"
        )
    else:
        top_risks.append(
            f"Gross margin at {avg_margin:.1f}% — if this falls below 40% it will impact profitability targets"
        )

    # Top opportunities (generated from data)
    top_opportunities = []
    strong_cats = [c for c in categories if c.get("health") == "STRONG"]
    if strong_cats:
        top_opportunities.append(
            f"Accelerate intake in {strong_cats[0]['category']} — strong sell-through of "
            f"{strong_cats[0].get('avg_sell_through', 0):.0f}% signals demand outpacing supply"
        )
    else:
        top_opportunities.append(
            f"Focus buying power on top performers — best sellers driving £{best_sellers[0]['revenue']:,.0f} each"
            if best_sellers else "Identify hero SKUs and concentrate OTB to drive sell-through"
        )

    if reorder_count > 0:
        top_opportunities.append(
            f"Rapid replenishment of {open_po_count} critical/urgent SKUs could unlock "
            f"£{open_po_value * 6:,.0f} in additional revenue over 6 weeks"
        )
    else:
        top_opportunities.append(
            "Strong stock position — opportunity to negotiate better terms with suppliers given healthy cover"
        )

    if avg_margin > 45:
        top_opportunities.append(
            f"Gross margin at {avg_margin:.1f}% is above benchmark — explore premium pricing on hero SKUs "
            f"to further improve profitability"
        )
    else:
        top_opportunities.append(
            "Markdown execution on slow movers will release working capital and improve overall range margin"
        )

    # AI actions (6 ranked actions for Merchandiser view)
    ai_actions = []
    rank = 1

    # Action 1: top markdown candidate
    if markdowns:
        top_md = markdowns[0]
        ai_actions.append({
            "rank": rank, "type": "Markdown",
            "sku": top_md["sku"],
            "description": f"Reduce price {top_md.get('recommended_depth','25-35%')} to clear excess stock and free working capital",
            "revenue_impact": round(top_md.get("revenue", 0) * 0.3, 0),
            "margin_impact": -8.5,
            "confidence": 87,
            "priority": "Critical" if top_md["severity"] == "CLEARANCE" else "Urgent"
        })
        rank += 1

    # Action 2: top reorder
    if reorders:
        top_ro = reorders[0]
        ai_actions.append({
            "rank": rank, "type": "Reorder",
            "sku": top_ro["sku"],
            "description": "Replenish immediately to maintain 6-week cover and protect revenue",
            "revenue_impact": round(top_ro.get("revenue", 0) * 1.2, 0),
            "margin_impact": 2.3,
            "confidence": 92,
            "priority": "Critical" if top_ro["urgency"] == "CRITICAL" else "Urgent"
        })
        rank += 1

    # Action 3: second markdown
    if len(markdowns) > 1:
        md2 = markdowns[1]
        ai_actions.append({
            "rank": rank, "type": "Markdown",
            "sku": md2["sku"],
            "description": f"Tactical {md2.get('recommended_depth','15-20%')} price reduction to stimulate demand",
            "revenue_impact": round(md2.get("revenue", 0) * 0.15, 0),
            "margin_impact": -4.2,
            "confidence": 79,
            "priority": "Urgent" if md2["severity"] in ("CLEARANCE", "DEEP") else "Recommended"
        })
        rank += 1

    # Action 4: category action
    if weak_cats:
        wc = weak_cats[0]
        ai_actions.append({
            "rank": rank, "type": "Cancel PO",
            "sku": wc["category"],
            "description": f"Review and reduce forward orders for {wc['category']} — low sell-through signals weak demand",
            "revenue_impact": round(wc.get("revenue", 0) * -0.1, 0),
            "margin_impact": 3.1,
            "confidence": 74,
            "priority": "Recommended"
        })
        rank += 1
    elif len(reorders) > 1:
        ro2 = reorders[1]
        ai_actions.append({
            "rank": rank, "type": "Reorder",
            "sku": ro2["sku"],
            "description": "Secondary replenishment — prevent stockout before next delivery window",
            "revenue_impact": round(ro2.get("revenue", 0) * 0.8, 0),
            "margin_impact": 1.8,
            "confidence": 85,
            "priority": "Urgent"
        })
        rank += 1

    # Action 5: monitor strong category
    if strong_cats:
        sc = strong_cats[0]
        ai_actions.append({
            "rank": rank, "type": "Monitor",
            "sku": sc["category"],
            "description": f"Accelerate {sc['category']} intake — strong sell-through of {sc.get('avg_sell_through',0):.0f}% signals untapped demand",
            "revenue_impact": round(sc.get("revenue", 0) * 0.25, 0),
            "margin_impact": 1.5,
            "confidence": 81,
            "priority": "Recommended"
        })
        rank += 1
    else:
        ai_actions.append({
            "rank": rank, "type": "Monitor",
            "sku": best_sellers[0]["sku"] if best_sellers else "Top SKU",
            "description": "Monitor best sellers daily — high velocity risks stockout without adequate cover",
            "revenue_impact": round(total_rev * 0.05, 0),
            "margin_impact": 0.5,
            "confidence": 76,
            "priority": "Recommended"
        })
        rank += 1

    # Action 6: OTB / buying direction
    ai_actions.append({
        "rank": rank, "type": "Reorder",
        "sku": "OTB Review",
        "description": f"Reallocate OTB from underperforming categories to top performers — estimated {avg_margin:.0f}% margin uplift available",
        "revenue_impact": round(total_rev * 0.08, 0),
        "margin_impact": 2.0,
        "confidence": 70,
        "priority": "Recommended"
    })

    # Fill to 6 if short
    while len(ai_actions) < 6:
        ai_actions.append({
            "rank": len(ai_actions) + 1, "type": "Monitor",
            "sku": "Range Review",
            "description": "Conduct weekly range review to identify emerging trends and risks",
            "revenue_impact": 0,
            "margin_impact": 0,
            "confidence": 65,
            "priority": "Recommended"
        })

    # ── Attach 'Why Skuvvy recommends this' reasoning to every action ─────────
    _reason_ctx = {
        "md_by":     {m.get("sku"): m for m in markdowns},
        "ro_by":     {r.get("sku"): r for r in reorders},
        "cat_by":    {c.get("category"): c for c in categories},
        "seller_by": {b.get("sku"): b for b in best_sellers},
        "avg_st":    avg_st,
        "avg_cover": avg_cover,
    }
    for _a in ai_actions:
        _a["reasoning"] = _action_reasoning(_a, _reason_ctx)

    # Executive summary (generated inline, not from Claude API)
    perf_word = "above" if vs_plan["revenue_pct"] >= 100 else "below"
    top_cat = categories[0] if categories else None
    weak_cat = next((c for c in categories if c.get("health") in ("AT RISK", "CRITICAL")), None)
    critical_count = k.get("critical_oos_count", 0)

    summary_parts = [
        f"Trading is performing {perf_word} plan at {vs_plan['revenue_pct']:.0f}% of revenue target.",
    ]
    if top_cat:
        summary_parts.append(
            f"{top_cat['category']} is the top performing category at £{top_cat['revenue']:,.0f} revenue"
            + (f" and {top_cat.get('avg_sell_through', 0):.0f}% sell-through." if top_cat.get('avg_sell_through') else ".")
        )
    if weak_cat:
        summary_parts.append(
            f"{weak_cat['category']} is flagged as {weak_cat['health']} with only "
            f"{weak_cat.get('avg_sell_through', 0):.0f}% sell-through — immediate action required."
        )
    if critical_count > 0:
        summary_parts.append(
            f"{critical_count} SKUs are critically out-of-stock — emergency reorder recommended to prevent further revenue loss."
        )
    elif reorder_count > 0:
        summary_parts.append(
            f"{reorder_count} SKUs require reorder action to maintain adequate cover across the range."
        )
    if avg_margin > 0:
        action_verb = "Prioritise markdown execution" if len(markdowns) > 3 else "Maintain pricing discipline"
        est_uplift = round((avg_margin * 1.05 - avg_margin) / avg_margin * 100, 1) if avg_margin else 0
        summary_parts.append(
            f"{action_verb} to optimise margin by approximately {est_uplift:.1f}% over the next 4 weeks."
        )

    executive_summary = " ".join(summary_parts)

    return {
        "bottom_sellers": bottom_sellers,
        "total_stock_units": round(total_stock_units),
        "total_stock_value": total_stock_value,
        "missing_data_count": missing_data_count,
        "data_quality_pct": data_quality_pct,
        "open_po_count": open_po_count,
        "open_po_value": open_po_value,
        "health_score": health_score,
        "health_breakdown": health_detail,
        "plan_data": plan_data,
        "vs_plan": vs_plan,
        "cash_trapped": cash_trapped,
        "top_risks": top_risks,
        "top_opportunities": top_opportunities,
        "ai_actions": ai_actions[:6],
        "executive_summary": executive_summary,
        "weekly_sparkline": sparkline,
    }


# ═══════════════════════════════════════════════════════════════
# PUBLIC ROUTES
# ═══════════════════════════════════════════════════════════════

@app.route("/")
def index():
    return send_from_directory(".", "aeroa_fixed.html")

@app.route("/login")
def login_page():
    if get_current_user():
        return redirect("/dashboard")
    return send_from_directory(".", "login.html")

@app.route("/signup")
def signup_page():
    if get_current_user():
        return redirect("/dashboard")
    return send_from_directory(".", "signup.html")

@app.route("/health")
def health():
    key = os.environ.get("ANTHROPIC_API_KEY", "")
    key_preview = (key[:12] + "...") if len(key) > 12 else f"MISSING or too short ({len(key)} chars)"
    return jsonify({"status": "ok", "timestamp": datetime.utcnow().isoformat(), "ai_key": key_preview, "ai_model": MODEL})


@app.route("/api/test-ai")
def test_ai():
    """Diagnostic: test Anthropic connectivity."""
    try:
        resp = ai_client.messages.create(
            model=MODEL, max_tokens=20,
            messages=[{"role": "user", "content": "Say hi"}]
        )
        return jsonify({"ok": True, "reply": resp.content[0].text, "model": MODEL})
    except Exception as e:
        return jsonify({"ok": False, "error": str(e), "model": MODEL}), 500


# ═══════════════════════════════════════════════════════════════
# AUTH ROUTES
# ═══════════════════════════════════════════════════════════════

@app.route("/auth/signup", methods=["POST"])
@auth_limit()  # brute-force / abuse protection: 10/min per IP
def auth_signup():
    # Server-side validation (format + length + strength), independent of the
    # frontend. Sanitisation reuses security.validators under the hood.
    try:
        data = SignupSchema.model_validate(request.get_json(silent=True) or {})
    except ValidationError:
        # Generic message — never disclose which field failed.
        audit.log(AUTH_FAILED, request=request, status="FAILURE",
                  metadata={"stage": "signup_validation"})
        return jsonify({"error": "Please enter a valid email and a password with "
                                 "at least 8 characters, one uppercase letter and one number."}), 400
    email, password, name = data.email, data.password, data.full_name

    # The "verify" response below is returned IDENTICALLY whether or not the
    # email is already registered, so signup cannot be used to enumerate accounts.
    _generic_verify = jsonify({"ok": True, "verify": True,
                               "message": "Check your email for a 6-digit verification code"})
    try:
        # Use sign_up (not admin.create_user) so Supabase sends the OTP verification email
        res = supabase.auth.sign_up({
            "email": email,
            "password": password,
            "options": {
                "data": {"full_name": name},
            },
        })
        if not res.user:
            # Do not reveal failure detail — keep the response uniform.
            return _generic_verify

        # Empty identities => email already registered. Return the SAME generic
        # response instead of confirming existence (anti-enumeration).
        return _generic_verify
    except Exception as e:
        msg = str(e).lower()
        if "already registered" in msg or "already been registered" in msg or "duplicate" in msg:
            # Already exists — stay uniform, do not confirm.
            return _generic_verify
        app.logger.error(f"Signup error: {e}", exc_info=True)
        return jsonify({"error": "Signup failed — please try again"}), 500


@app.route("/auth/login", methods=["POST"])
@auth_limit()  # brute-force protection: 10/min per IP
def auth_login():
    def _deny():
        # ONE generic response for invalid format, wrong credentials, AND lockout.
        # An attacker must not be able to tell these cases apart.
        return jsonify({"error": "Incorrect email or password"}), 401

    # Server-side validation, independent of the frontend.
    try:
        data = LoginSchema.model_validate(request.get_json(silent=True) or {})
    except ValidationError:
        audit.log(AUTH_FAILED, request=request, status="FAILURE",
                  metadata={"stage": "login_validation"})
        return _deny()
    email, password = data.email, data.password

    # Per-account lockout, layered on the per-IP rate limit. A locked account
    # returns the SAME message as a wrong password — the lockout is never disclosed.
    locked, _retry = lockout.status(email)
    if locked:
        time.sleep(min(lockout.delay_for_current(email), 5.0))  # match failure timing
        audit.log(AUTH_FAILED, request=request, status="BLOCKED",
                  metadata={"email": email, "reason": "locked"})
        return _deny()

    try:
        res = supabase.auth.sign_in_with_password({"email": email, "password": password})
    except Exception as e:
        msg = str(e).lower()
        if "invalid" in msg or "credentials" in msg:
            res = None  # bad creds -> fall through to the failure path
        else:
            app.logger.error(f"Login error: {e}", exc_info=True)
            return jsonify({"error": "Login failed — please try again"}), 500

    if res and getattr(res, "user", None) and getattr(res, "session", None):
        lockout.record_success(email)  # clear the failed-attempt counter
        meta = res.user.user_metadata or {}
        session["user_id"]         = res.user.id
        session["user_email"]      = email
        session["user_name"]       = meta.get("full_name", email.split("@")[0])
        session["user_plan"]       = "trial"
        session["access_token"]    = res.session.access_token
        session["refresh_token"]   = res.session.refresh_token
        csrf_token = generate_csrf_token()
        audit.log(USER_LOGIN, request=request, user_id=res.user.id,
                  company_id=session.get("company_id"),
                  metadata={"email": email})
        return jsonify({"ok": True, "redirect": "/dashboard", "csrf_token": csrf_token})

    # ── Failed attempt: count it, apply progressive delay, maybe lock ─────────
    fails, just_locked = lockout.record_failure(email)
    if just_locked:
        # Best-effort recovery email via Supabase's existing reset flow. It reads
        # "reset your password" (NOT "account locked"), so it gives the real user
        # a reset link without disclosing the lockout to an attacker.
        try:
            supabase.auth.reset_password_for_email(
                email, {"redirect_to": "https://aeroa-ai.up.railway.app/reset-password"})
        except Exception:
            pass
    time.sleep(min(lockout.delay_for_current(email), 5.0))  # progressive delay
    audit.log(AUTH_FAILED, request=request, status="FAILURE",
              metadata={"email": email, "fails": fails, "locked": just_locked})
    return _deny()


# ─────────────────────────────────────────────────────────────────────────────
# Google OAuth (Supabase social login)
# ─────────────────────────────────────────────────────────────────────────────
# Flow (backend-driven, no secrets in the browser):
#   1. /auth/google         -> 302 to Supabase's hosted authorize endpoint
#   2. Supabase <-> Google   -> redirects back to /auth/callback with the
#                               session tokens in the URL *hash* (implicit flow)
#   3. /auth/callback        -> tiny JS page reads the hash and POSTs the tokens
#   4. /auth/oauth-session   -> validates the token, sets the Flask session,
#                               upserts the profile (mirrors /auth/verify-email)
# The Supabase Google provider must be enabled and {base}/auth/callback must be
# in Supabase's allow-listed redirect URLs (dashboard config).

@app.route("/auth/google")
def auth_google():
    """Kick off Google sign-in by redirecting to Supabase's authorize endpoint."""
    if not SUPABASE_URL:
        return redirect("/login?error=oauth_unavailable")
    base = request.host_url.rstrip("/")
    redirect_to = quote(f"{base}/auth/callback", safe="")
    authorize_url = (f"{SUPABASE_URL}/auth/v1/authorize"
                     f"?provider=google&redirect_to={redirect_to}")
    audit.log(USER_LOGIN, request=request, status="PENDING",
              metadata={"provider": "google", "stage": "authorize"})
    return redirect(authorize_url, code=302)


@app.route("/auth/callback")
def auth_callback():
    """Serve the small page that reads OAuth tokens from the URL hash."""
    return send_from_directory(".", "oauth-callback.html")


@app.route("/auth/oauth-session", methods=["POST"])
@auth_limit()  # abuse protection: 10/min per IP
def auth_oauth_session():
    """Validate Supabase OAuth tokens and establish the Flask session."""
    body = request.get_json(silent=True) or {}
    access_token  = (body.get("access_token") or "").strip()
    refresh_token = (body.get("refresh_token") or "").strip()
    if not access_token:
        return jsonify({"error": "Sign-in failed. Please try again."}), 400
    try:
        # get_user(jwt) validates the token WITHOUT mutating the shared client.
        ures = supabase.auth.get_user(access_token)
        user = getattr(ures, "user", None)
        if not user or not getattr(user, "id", None):
            audit.log(AUTH_FAILED, request=request, status="FAILURE",
                      metadata={"provider": "google", "stage": "token_validate"})
            return jsonify({"error": "Sign-in failed. Please try again."}), 401

        email = (getattr(user, "email", "") or "").lower()
        meta  = getattr(user, "user_metadata", None) or {}
        name  = (meta.get("full_name") or meta.get("name")
                 or (email.split("@")[0] if email else "there"))
        trial_ends = (datetime.now(timezone.utc) + timedelta(days=14)).isoformat()

        # Create or update the profile (non-fatal — session still works without it).
        try:
            get_user_db().table("profiles").upsert({
                "id": user.id, "email": email, "full_name": name,
                "plan": "trial", "trial_ends": trial_ends,
            }).execute()
        except Exception:
            pass

        session["user_id"]         = user.id
        session["user_email"]      = email
        session["user_name"]       = name
        session["user_plan"]       = "trial"
        session["user_trial_ends"] = trial_ends
        session["access_token"]    = access_token
        session["refresh_token"]   = refresh_token
        csrf_token = generate_csrf_token()
        audit.log(USER_LOGIN, request=request, user_id=user.id,
                  metadata={"email": email, "provider": "google"})
        return jsonify({"ok": True, "redirect": "/dashboard", "csrf_token": csrf_token})
    except Exception as e:
        app.logger.error(f"OAuth session error: {e}")
        return jsonify({"error": "Sign-in failed. Please try again."}), 400


@app.route("/auth/logout", methods=["POST"])
def auth_logout():
    uid = session.get("user_id")
    cid = session.get("company_id")
    audit.log(USER_LOGOUT, request=request, user_id=uid, company_id=cid)
    session.clear()
    return jsonify({"ok": True, "redirect": "/login"})


@app.route("/auth/forgot-password", methods=["POST"])
@auth_limit()  # abuse / enumeration protection: 10/min per IP
def auth_forgot_password():
    # ONE generic response whether or not the email is registered — this is the
    # only thing the endpoint may ever say, so it can't be used to enumerate.
    _generic = jsonify({"ok": True,
                        "message": "If that email is registered, you'll receive a reset link"})
    try:
        data = ForgotPasswordSchema.model_validate(request.get_json(silent=True) or {})
    except ValidationError:
        audit.log(AUTH_FAILED, request=request, status="FAILURE",
                  metadata={"stage": "forgot_password_validation"})
        # Still return the uniform message: a malformed email must not look
        # different from an unregistered one.
        return _generic
    email = data.email
    try:
        redirect_url = "https://aeroa-ai.up.railway.app/reset-password"
        supabase.auth.reset_password_for_email(email, {"redirect_to": redirect_url})
    except Exception as e:
        app.logger.error(f"Forgot password error: {e}")
    # Always the same response — never disclose whether the email exists.
    return _generic


@app.route("/reset-password")
def reset_password_page():
    return send_from_directory(".", "reset-password.html")


@app.route("/auth/reset-password", methods=["POST"])
@auth_limit()  # abuse protection: 10/min per IP
def do_reset_password():
    # Server-side validation: token presence + password STRENGTH (this is a new
    # password, so the full strength policy applies — unlike login).
    try:
        data = ResetPasswordSchema.model_validate(request.get_json(silent=True) or {})
    except ValidationError:
        audit.log(AUTH_FAILED, request=request, status="FAILURE",
                  metadata={"stage": "reset_password_validation"})
        return jsonify({"error": "Please choose a password with at least 8 characters, "
                                 "one uppercase letter and one number."}), 400
    access_token  = data.access_token
    refresh_token = data.refresh_token
    new_password  = data.password
    try:
        supabase.auth.set_session(access_token, refresh_token)
        supabase.auth.update_user({"password": new_password})
        return jsonify({"ok": True, "message": "Password updated — you can now log in"})
    except Exception as e:
        app.logger.error(f"Reset password error: {e}")
        return jsonify({"error": "Reset link has expired. Please request a new one."}), 400


@app.route("/auth/verify-email", methods=["POST"])
@auth_limit()  # OTP brute-force protection: 10/min per IP
def auth_verify_email():
    # Server-side validation: email format + short numeric token. Caps the OTP
    # length so the code can't be used as a brute-force oracle with junk input.
    try:
        data = VerifyEmailSchema.model_validate(request.get_json(silent=True) or {})
    except ValidationError:
        audit.log(AUTH_FAILED, request=request, status="FAILURE",
                  metadata={"stage": "verify_email_validation"})
        return jsonify({"error": "Invalid or expired code — please try again"}), 400
    email = data.email
    token = data.token
    name  = data.full_name
    try:
        # Try "signup" type first (Supabase v2 sends OTP with type=signup for new accounts)
        res = None
        last_error = None
        for otp_type in ("signup", "email"):
            try:
                res = supabase.auth.verify_otp({"email": email, "token": token, "type": otp_type})
                if res.user and res.session:
                    break
            except Exception as e2:
                last_error = e2
                continue

        if res and res.user and res.session:
            from datetime import timezone
            trial_ends = (datetime.now(timezone.utc) + __import__('datetime').timedelta(days=14)).isoformat()
            # Create or update profile in DB
            try:
                get_user_db().table("profiles").upsert({
                    "id": res.user.id,
                    "email": email,
                    "full_name": name or email.split("@")[0],
                    "plan": "trial",
                    "trial_ends": trial_ends,
                }).execute()
            except Exception:
                pass
            session["user_id"]         = res.user.id
            session["user_email"]      = email
            session["user_name"]       = name or email.split("@")[0]
            session["user_plan"]       = "trial"
            session["user_trial_ends"] = trial_ends
            session["access_token"]    = res.session.access_token
            session["refresh_token"]   = res.session.refresh_token
            csrf_token = generate_csrf_token()
            return jsonify({"ok": True, "redirect": "/dashboard", "csrf_token": csrf_token})
        return jsonify({"error": "Invalid or expired code — try again"}), 400
    except Exception as e:
        app.logger.error(f"Verify email error: {e}")
        return jsonify({"error": "Invalid or expired code — please try again"}), 400


@app.route("/auth/me")
def auth_me():
    user = get_current_user()
    if not user:
        return jsonify({"error": "Not authenticated"}), 401
    return jsonify({
        "id":        user["id"],
        "email":     user["email"],
        "full_name": user.get("full_name", ""),
        "plan":      user.get("plan", "trial"),
        "trial_ends": user.get("trial_ends"),
        "ai_messages_used": user.get("ai_messages_used", 0),
        "role": session.get("user_role", user.get("role", "AM")),
        "csrf_token": get_csrf_token(),
    })


# ═══════════════════════════════════════════════════════════════
# ROLE ROUTES
# ═══════════════════════════════════════════════════════════════

@app.route("/api/user/role", methods=["GET"])
@require_auth
def get_user_role(user):
    """Return current role from session (default AM)."""
    role = session.get("user_role", user.get("role", "AM"))
    return jsonify({"role": role})


@app.route("/api/user/role", methods=["POST"])
@require_auth
def set_user_role(user):
    """Set role in session and save to Supabase profiles.role column."""
    body = request.get_json(silent=True) or {}
    role = (body.get("role") or "AM").strip()

    valid_roles = {"MA", "AM", "Merchandiser", "Director"}
    if role not in valid_roles:
        return jsonify({"error": f"Invalid role. Must be one of: {', '.join(valid_roles)}"}), 400

    old_role = session.get("user_role", user.get("role", "MA"))
    session["user_role"] = role

    # Try to persist to Supabase (non-fatal if column doesn't exist yet)
    try:
        get_user_db().table("profiles").update({"role": role}).eq("id", user["id"]).execute()
        # Invalidate the session cache so next get_current_user() reloads
        session.pop("user_profile", None)
    except Exception as e:
        app.logger.warning(f"Could not persist role to Supabase: {e}")

    audit.log(ROLE_CHANGE, request=request, user_id=user["id"],
              company_id=user.get("company_id"),
              metadata={"from_role": old_role, "to_role": role})

    return jsonify({"ok": True, "role": role})


@app.route("/api/user/onboarding", methods=["POST"])
@require_auth
def save_onboarding(user):
    """Persist onboarding answers (role, categories, challenges). Role maps to the
    existing profiles.role column; the rest is stored best-effort in an
    `onboarding` JSON column if present (non-fatal if the column doesn't exist)."""
    body = request.get_json(silent=True) or {}
    role = (body.get("role") or "").strip()
    categories = body.get("categories") or []
    challenges = body.get("challenges") or []

    role_map = {  # onboarding labels → internal role codes
        "Assistant Merchandiser": "MA", "Merchandising Assistant": "MA",
        "Merchandiser": "Merchandiser", "Freelancer": "Merchandiser",
        "Consultant": "Director",
    }
    internal_role = role_map.get(role, session.get("user_role", "AM"))
    session["user_role"] = internal_role
    session["onboarding"] = {"role": role, "categories": categories, "challenges": challenges}

    try:
        get_user_db().table("profiles").update({
            "role": internal_role,
            "onboarding": {"role": role, "categories": categories, "challenges": challenges},
            "onboarded": True,
        }).eq("id", user["id"]).execute()
        session.pop("user_profile", None)
    except Exception as e:
        # Most likely the optional columns aren't in the schema yet — that's fine,
        # the frontend also remembers completion locally.
        app.logger.info(f"Onboarding persist (non-fatal): {e}")
        try:
            get_user_db().table("profiles").update({"role": internal_role}).eq("id", user["id"]).execute()
            session.pop("user_profile", None)
        except Exception:
            pass

    return jsonify({"ok": True, "role": internal_role})


# ═══════════════════════════════════════════════════════════════
# PROTECTED APP ROUTES
# ═══════════════════════════════════════════════════════════════

@app.route("/dashboard")
@require_auth
def dashboard(user):
    return send_from_directory(".", "dashboard.html")


def _get_retention_hours(company_id) -> int:
    """Return configured retention hours for a company (default 24h)."""
    if not company_id:
        return 24
    try:
        res = supabase.table("companies").select("retention_hours") \
            .eq("id", str(company_id)).single().execute()
        if res.data:
            return int(res.data.get("retention_hours", 24))
    except Exception:
        pass
    return 24


@app.route("/api/upload", methods=["POST"])
@require_auth
@limiter.limit("30 per hour")
def upload(user):
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400

    f = request.files["file"]
    if not f.filename:
        return jsonify({"error": "Empty filename"}), 400

    filename   = safe_filename(f.filename)
    file_bytes = f.read()
    company_id = user.get("company_id")

    # ── Stage 1-4: Security pipeline ──────────────────────────────
    scan = file_pipeline.validate(file_bytes, filename)
    if not scan.valid:
        audit.log(FILE_UPLOAD, request=request, user_id=user["id"],
                  company_id=company_id, resource=filename,
                  status="FAILURE", metadata={"reason": scan.error})
        return jsonify({"error": scan.error}), 422

    # ── Stage 5: Parse and analyse ────────────────────────────────
    try:
        df           = parse_upload(io.BytesIO(file_bytes), filename)
        result       = analyse(df)
        sku_count    = result["meta"]["sku_count"]

        plan         = user.get("plan", "trial")
        sku_limit    = PLAN_LIMITS.get(plan, PLAN_LIMITS["trial"])["skus"]
        if sku_count > sku_limit:
            return jsonify({
                "error": f"Your {plan.title()} plan supports up to {sku_limit} SKUs. "
                         f"This file has {sku_count}. Please upgrade.",
                "upgrade_required": True,
            }), 402

        # ── Stage 6: Store encrypted ───────────────────────────────
        _, file_hash   = encrypt_file(file_bytes)          # hash pre-encryption
        retention_hrs  = _get_retention_hours(company_id)
        expires_at     = (datetime.now(timezone.utc) + timedelta(hours=retention_hrs)).isoformat()

        upload_row = get_user_db().table("uploads").insert({
            "user_id":        user["id"],
            "company_id":     company_id,
            "filename":       filename,
            "sku_count":      sku_count,
            "analysis":       result,
            "file_hash":      file_hash,
            "file_size":      len(file_bytes),
            "file_encrypted": True,
            "expires_at":     expires_at,
        }).execute()

        upload_id = upload_row.data[0]["id"] if upload_row.data else None
        session["current_upload_id"] = upload_id

        audit.log(FILE_UPLOAD, request=request, user_id=user["id"],
                  company_id=company_id,
                  resource=f"upload:{upload_id}",
                  metadata={
                      "filename":  filename,
                      "sku_count": sku_count,
                      "file_size": len(file_bytes),
                      "expires_at": expires_at,
                  })

        return jsonify({
            "status":    "ok",
            "upload_id": upload_id,
            "filename":  filename,
            "sku_count": sku_count,
            "kpis":      result["kpis"],
            "expires_at": expires_at,
            "security": {
                "encrypted":      True,
                "retention_hours": retention_hrs,
                "file_hash":      file_hash[:16] + "…",
            },
        })

    except ValueError as e:
        return jsonify({"error": str(e)}), 422
    except Exception as e:
        app.logger.error(f"Upload error: {e}", exc_info=True)
        # Surface the real reason (type + short message) instead of an opaque
        # message — parsing failures are diagnostic, not sensitive, and the
        # generic text left us blind to format/library issues.
        detail = f"{type(e).__name__}: {str(e)[:200]}"
        return jsonify({
            "error": f"Couldn't process this file. {detail}",
            "detail": detail,
        }), 500


@app.route("/api/analysis")
@require_auth
def get_analysis(user):
    upload_id = request.args.get("upload_id") or session.get("current_upload_id")
    db = get_user_db()
    company_id = user.get("company_id")

    if upload_id:
        try:
            q = db.table("uploads").select("analysis").eq("id", upload_id).eq("user_id", user["id"])
            # Only filter by company_id when it is actually set (avoids "None" string mismatch)
            if company_id is not None:
                q = q.eq("company_id", company_id)
            res = q.single().execute()
            if res.data:
                audit.log(DATA_EXPORT, request=request, user_id=user["id"],
                          company_id=company_id,
                          resource=f"analysis:{upload_id or 'latest'}")
                return jsonify(res.data["analysis"])
        except Exception:
            pass

    # Fallback: return latest upload for this user
    try:
        q = db.table("uploads").select("analysis").eq("user_id", user["id"])
        if company_id is not None:
            q = q.eq("company_id", company_id)
        res = q.order("created_at", desc=True).limit(1).execute()
        if res.data:
            audit.log(DATA_EXPORT, request=request, user_id=user["id"],
                      company_id=company_id,
                      resource=f"analysis:{upload_id or 'latest'}")
            return jsonify(res.data[0]["analysis"])
    except Exception:
        pass

    return jsonify({"error": "No data uploaded yet"}), 404


@app.route("/api/analysis/role")
@require_auth
def get_role_analysis(user):
    """Return extended role-based analysis data."""
    upload_id = request.args.get("upload_id") or session.get("current_upload_id")
    db = get_user_db()
    analysis = None

    if upload_id:
        try:
            res = db.table("uploads").select("analysis") \
                .eq("id", upload_id).eq("user_id", user["id"]).single().execute()
            if res.data:
                analysis = res.data["analysis"]
        except Exception:
            pass

    if not analysis:
        try:
            res = db.table("uploads").select("analysis") \
                .eq("user_id", user["id"]).order("created_at", desc=True).limit(1).execute()
            if res.data:
                analysis = res.data[0]["analysis"]
        except Exception:
            pass

    if not analysis:
        return jsonify({"error": "No data uploaded yet"}), 404

    try:
        role_data = _build_role_analysis(analysis)
        return jsonify(role_data)
    except Exception as e:
        app.logger.error(f"Role analysis error: {e}", exc_info=True)
        return jsonify({"error": "Failed to compute role analysis"}), 500


@app.route("/api/uploads")
@require_auth
def list_uploads(user):
    try:
        res = get_user_db().table("uploads") \
            .select("id, filename, sku_count, created_at") \
            .eq("user_id", user["id"]) \
            .order("created_at", desc=True).limit(20).execute()
        return jsonify({"uploads": res.data or []})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ═══════════════════════════════════════════════════════════════
# DATA WORKSPACE  —  Persistent file library, switcher, compare,
# timeline. Every upload becomes a reusable, searchable asset.
#
# Workspace metadata (display name, tags, archived, trashed) is
# stored inside the analysis JSONB blob under the "_workspace" key
# so no DB migration is required. PostgREST JSON-path selection
# (analysis->_workspace, analysis->kpis) keeps listing lightweight.
# ═══════════════════════════════════════════════════════════════

WORKSPACE_META_KEY = "_workspace"


def _file_type_from_name(filename: str) -> str:
    ext = (filename or "").rsplit(".", 1)[-1].lower() if "." in (filename or "") else ""
    return {
        "csv": "CSV", "xlsx": "Excel", "xls": "Excel",
        "numbers": "Numbers", "ods": "Spreadsheet",
        "tsv": "TSV", "txt": "Text",
    }.get(ext, (ext.upper() or "File"))


def _health_breakdown(kpis: dict) -> dict:
    """
    Explainable Business Health Score (0-100).

    Decomposes the score into five transparent components a merchandiser can
    reason about, each grounded in a real KPI:

      1. Revenue Performance  /30  — revenue quality (velocity × margin)
      2. Sell-Through         /20  — how fast stock is converting to sales
      3. Margin Protection    /20  — profitability vs a healthy benchmark
      4. Stock Efficiency     /20  — weeks of cover vs the ideal 4-8 wk band
      5. Risk Exposure        /10  — penalty for OOS / markdown / reorder load

    Returns {"total", "label", "components":[{key,name,score,max,detail,status}]}.
    The total is the canonical health score used across the product so the
    headline number and this breakdown always agree.
    """
    if not isinstance(kpis, dict):
        kpis = {}
    st       = float(kpis.get("avg_sell_through") or 0)
    mg       = float(kpis.get("gross_margin_pct") or kpis.get("avg_margin_pct") or 0)
    cover    = float(kpis.get("avg_cover_weeks") or 0)
    reorder  = int(kpis.get("reorder_count") or 0)
    critical = int(kpis.get("critical_oos_count") or 0)
    markdown = int(kpis.get("markdown_risk_count") or 0)
    margin_ok = kpis.get("margin_available") is not False

    def status(score, mx):
        r = score / mx if mx else 0
        return "strong" if r >= 0.75 else "watch" if r >= 0.45 else "weak"

    # 1. Revenue Performance /30 — composite of sell-through velocity and margin
    rev_pts = round(min(1.0, st / 85.0) * 18 + min(1.0, mg / 60.0) * 12)
    rev_pts = max(0, min(30, rev_pts))

    # 2. Sell-Through /20
    if   st >= 70: st_pts = 20
    elif st >= 55: st_pts = 16
    elif st >= 40: st_pts = 12
    elif st >= 25: st_pts = 7
    else:          st_pts = 3

    # 3. Margin Protection /20
    if not margin_ok:
        mg_pts = 10  # neutral — not penalised for missing data
        mg_detail = "Insufficient margin data in this dataset — scored neutrally."
        mg_status = "watch"
    else:
        if   mg >= 60: mg_pts = 20
        elif mg >= 50: mg_pts = 16
        elif mg >= 40: mg_pts = 12
        elif mg >= 30: mg_pts = 7
        else:          mg_pts = 3
        mg_detail = f"Blended margin of {mg:.1f}% vs a 55% healthy benchmark."
        mg_status = status(mg_pts, 20)

    # 4. Stock Efficiency /20 — ideal cover band is 4-8 weeks
    if   4 <= cover <= 8:                        cov_pts = 18
    elif (2 <= cover < 4) or (8 < cover <= 12):  cov_pts = 13
    elif (1 <= cover < 2) or (12 < cover <= 16): cov_pts = 7
    else:                                        cov_pts = 3
    if cover <= 0:
        cov_detail = "No stock-cover signal in this dataset."
    elif cover < 4:
        cov_detail = f"{cover:.1f} wks cover — running lean, stockout risk on top sellers."
    elif cover <= 8:
        cov_detail = f"{cover:.1f} wks cover — healthy, within the ideal 4-8 wk band."
    else:
        cov_detail = f"{cover:.1f} wks cover — overstocked, cash tied up in slow stock."

    # 5. Risk Exposure /10 — start clean, penalise live risks
    risk_pts = 10
    risk_pts -= min(5, critical * 2)        # critical OOS hurts most
    risk_pts -= min(3, reorder // 5)        # reorder backlog
    risk_pts -= min(2, markdown // 5)       # markdown dependency
    risk_pts = max(0, risk_pts)
    risk_bits = []
    if critical: risk_bits.append(f"{critical} critical OOS")
    if reorder:  risk_bits.append(f"{reorder} reorder alerts")
    if markdown: risk_bits.append(f"{markdown} markdown risks")
    risk_detail = ("Live risks: " + ", ".join(risk_bits) + ".") if risk_bits \
                  else "No critical stock or margin risks detected."

    components = [
        {"key": "revenue", "name": "Revenue Performance", "score": rev_pts, "max": 30,
         "status": status(rev_pts, 30),
         "detail": f"Revenue quality from {st:.0f}% sell-through and {mg:.0f}% margin."},
        {"key": "sell_through", "name": "Sell-Through", "score": st_pts, "max": 20,
         "status": status(st_pts, 20),
         "detail": f"Range sell-through of {st:.1f}% — how fast stock is converting."},
        {"key": "margin", "name": "Margin Protection", "score": mg_pts, "max": 20,
         "status": mg_status, "detail": mg_detail},
        {"key": "stock", "name": "Stock Efficiency", "score": cov_pts, "max": 20,
         "status": status(cov_pts, 20), "detail": cov_detail},
        {"key": "risk", "name": "Risk Exposure", "score": risk_pts, "max": 10,
         "status": status(risk_pts, 10), "detail": risk_detail},
    ]
    total = max(0, min(100, sum(c["score"] for c in components)))
    label = "Healthy" if total >= 80 else "Watch" if total >= 60 else "At Risk"
    return {"total": total, "label": label, "components": components}


def _ws_health_score(kpis: dict) -> int:
    """Canonical Business Health Score (0-100) — the breakdown total."""
    if not isinstance(kpis, dict):
        return 0
    return int(_health_breakdown(kpis)["total"])


def _dataset_type_from_name(fname: str, tags=None) -> str:
    """Infer the semantic dataset type (not the file extension) from name/tags."""
    n = (fname or "").lower()
    tagset = " ".join(tags or []).lower()
    hay = n + " " + tagset
    if "buy" in hay and "plan" in hay:
        return "Buy Plan"
    if "markdown" in hay or "clearance" in hay:
        return "Markdown Plan"
    if "supplier" in hay or "sourcing" in hay:
        return "Supplier Report"
    if "stock" in hay or "inventory" in hay:
        return "Stock Report"
    if "range" in hay or "assortment" in hay:
        return "Range Plan"
    if "trad" in hay or "weekly" in hay or "wk" in hay or "week" in hay:
        return "Trading Report"
    return "Trading Report"


def _ws_stock_value(kpis: dict) -> int:
    """Stock value at cost for a dataset (uses computed KPI, else derives from revenue)."""
    if not isinstance(kpis, dict):
        return 0
    sv = kpis.get("total_stock_value")
    if sv:
        return int(round(sv))
    rev = float(kpis.get("total_revenue") or 0)
    return int(round(rev * 0.5 / 12)) if rev else 0


def _ws_row_to_card(row: dict) -> dict:
    """Map a uploads row (with analysis->_workspace + analysis->kpis) to a file card."""
    ws    = row.get("_workspace") or row.get("ws") or {}
    if not isinstance(ws, dict):
        ws = {}
    kpis  = row.get("kpis") or {}
    if not isinstance(kpis, dict):
        kpis = {}
    fname = row.get("filename") or "Untitled"
    tags  = ws.get("tags") or []
    return {
        "id":            row.get("id"),
        "filename":      fname,
        "display_name":  ws.get("display_name") or fname,
        "upload_date":   row.get("created_at"),
        "analysed_date": ws.get("analysed_at") or row.get("created_at"),
        "file_type":     _file_type_from_name(fname),
        "dataset_type":  ws.get("dataset_type") or _dataset_type_from_name(fname, tags),
        "record_count":  row.get("sku_count") or 0,
        "file_size":     row.get("file_size") or 0,
        "status":        ws.get("status") or "Analysed",
        "tags":          tags,
        "archived":      bool(ws.get("archived")),
        "trashed":       bool(ws.get("trashed")),
        "revenue":       kpis.get("total_revenue") or 0,
        "margin":        kpis.get("avg_margin_pct") or 0,
        "sell_through":  kpis.get("avg_sell_through") or 0,
        "stock_value":   _ws_stock_value(kpis),
        "health_score":  _ws_health_score(kpis),
    }


def _ws_fetch_rows(user, *, limit=200):
    """Fetch lightweight workspace rows (no full analysis blob)."""
    db = get_user_db()
    return db.table("uploads") \
        .select("id, filename, sku_count, created_at, file_size, "
                "ws:analysis->_workspace, kpis:analysis->kpis") \
        .eq("user_id", user["id"]) \
        .order("created_at", desc=True).limit(limit).execute()


@app.route("/api/workspace/files")
@require_auth
def workspace_files(user):
    """List every persisted file for this user with workspace metadata."""
    q          = (request.args.get("q") or "").strip().lower()
    view       = (request.args.get("view") or "active").lower()   # active|archived|trash|all
    tag_filter = (request.args.get("tag") or "").strip().lower()
    try:
        res   = _ws_fetch_rows(user)
        cards = [_ws_row_to_card(r) for r in (res.data or [])]

        out = []
        for c in cards:
            if view == "active"   and (c["archived"] or c["trashed"]):  continue
            if view == "archived" and (not c["archived"] or c["trashed"]): continue
            if view == "trash"    and not c["trashed"]:                  continue
            if view == "all"      and c["trashed"]:                      continue
            if q and q not in (c["display_name"] + " " + c["filename"]).lower() \
                 and not any(q in t.lower() for t in c["tags"]):
                continue
            if tag_filter and not any(tag_filter == t.lower() for t in c["tags"]):
                continue
            out.append(c)

        # Tag cloud across active+archived (exclude trash)
        all_tags = {}
        for c in cards:
            if c["trashed"]:
                continue
            for t in c["tags"]:
                all_tags[t] = all_tags.get(t, 0) + 1

        return jsonify({
            "files": out,
            "tags": sorted(all_tags.keys()),
            "current_upload_id": session.get("current_upload_id"),
            "counts": {
                "active":   sum(1 for c in cards if not c["archived"] and not c["trashed"]),
                "archived": sum(1 for c in cards if c["archived"] and not c["trashed"]),
                "trash":    sum(1 for c in cards if c["trashed"]),
                "total":    len(cards),
            },
        })
    except Exception as e:
        app.logger.error(f"workspace_files error: {e}", exc_info=True)
        return jsonify({"error": str(e)}), 500


def _ws_mutate(user, upload_id, mutator):
    """Read analysis, mutate its _workspace dict, write back. Returns updated card."""
    db = get_user_db()
    res = db.table("uploads").select("analysis, filename, sku_count, created_at, file_size") \
        .eq("id", upload_id).eq("user_id", user["id"]).single().execute()
    if not res.data:
        return None
    analysis = res.data.get("analysis") or {}
    if not isinstance(analysis, dict):
        analysis = {}
    ws = analysis.get(WORKSPACE_META_KEY)
    if not isinstance(ws, dict):
        ws = {}
    mutator(ws)
    ws["updated_at"] = datetime.now(timezone.utc).isoformat()
    analysis[WORKSPACE_META_KEY] = ws
    db.table("uploads").update({"analysis": analysis}) \
        .eq("id", upload_id).eq("user_id", user["id"]).execute()
    return _ws_row_to_card({
        "id": upload_id,
        "filename": res.data.get("filename"),
        "sku_count": res.data.get("sku_count"),
        "created_at": res.data.get("created_at"),
        "file_size": res.data.get("file_size"),
        "_workspace": ws,
        "kpis": (analysis.get("kpis") or {}),
    })


@app.route("/api/workspace/file/<upload_id>/update", methods=["POST"])
@require_auth
def workspace_update_file(user, upload_id):
    """Rename, tag, archive/unarchive, trash/restore a file."""
    body = request.get_json(silent=True) or {}

    def mutator(ws):
        if "display_name" in body:
            name = (body.get("display_name") or "").strip()
            if name:
                ws["display_name"] = name[:120]
        if "tags" in body and isinstance(body["tags"], list):
            ws["tags"] = [str(t).strip()[:32] for t in body["tags"] if str(t).strip()][:12]
        if "archived" in body:
            ws["archived"] = bool(body["archived"])
        if "trashed" in body:
            ws["trashed"] = bool(body["trashed"])

    try:
        card = _ws_mutate(user, upload_id, mutator)
        if not card:
            return jsonify({"error": "File not found"}), 404
        audit.log(DATA_EXPORT, request=request, user_id=user["id"],
                  company_id=user.get("company_id"),
                  resource=f"workspace:update:{upload_id}",
                  metadata={k: body[k] for k in ("display_name", "archived", "trashed") if k in body})
        return jsonify({"ok": True, "file": card})
    except Exception as e:
        app.logger.error(f"workspace_update error: {e}", exc_info=True)
        return jsonify({"error": str(e)}), 500


@app.route("/api/workspace/switch", methods=["POST"])
@require_auth
def workspace_switch(user):
    """Make a previously-uploaded file the active dataset — no re-upload."""
    body = request.get_json(silent=True) or {}
    upload_id = body.get("upload_id")
    if not upload_id:
        return jsonify({"error": "upload_id required"}), 400
    try:
        db = get_user_db()
        res = db.table("uploads").select("id, filename") \
            .eq("id", upload_id).eq("user_id", user["id"]).single().execute()
        if not res.data:
            return jsonify({"error": "File not found"}), 404
        session["current_upload_id"] = upload_id
        return jsonify({"ok": True, "upload_id": upload_id,
                        "filename": res.data.get("filename")})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


def _ws_key_rec(card: dict) -> str:
    """Synthesize a one-line 'next move' from a card's KPIs (no AI cost)."""
    st = float(card.get("sell_through") or 0)
    mg = float(card.get("margin") or 0)
    if st >= 70:
        return "Reorder best-sellers — strong sell-through, protect availability"
    if st < 45:
        return "Trigger early markdown on slow stock to release cash"
    if mg < 40:
        return "Review pricing & markdown depth to protect margin"
    return "Hold and monitor — trading in line with plan"


def _ws_compare_metrics(kpis: dict) -> dict:
    kpis = kpis or {}
    return {
        "revenue":      float(kpis.get("total_revenue") or 0),
        "units":        float(kpis.get("total_units") or 0),
        "margin":       float(kpis.get("avg_margin_pct") or 0),
        "sell_through": float(kpis.get("avg_sell_through") or 0),
        "reorders":     float(kpis.get("reorder_count") or 0),
        "health":       _ws_health_score(kpis),
    }


@app.route("/api/workspace/compare")
@require_auth
def workspace_compare(user):
    """Compare two datasets (A vs B): revenue, margin, sell-through, category movement."""
    id_a = request.args.get("a")
    id_b = request.args.get("b")
    if not id_a or not id_b:
        return jsonify({"error": "Two file ids (a, b) required"}), 400
    try:
        db = get_user_db()

        def load(uid):
            r = db.table("uploads") \
                .select("id, filename, created_at, sku_count, "
                        "ws:analysis->_workspace, kpis:analysis->kpis, "
                        "cats:analysis->category_scorecard") \
                .eq("id", uid).eq("user_id", user["id"]).single().execute()
            return r.data

        a = load(id_a)
        b = load(id_b)
        if not a or not b:
            return jsonify({"error": "One or both files not found"}), 404

        def label(row):
            ws = row.get("ws") or {}
            return (ws.get("display_name") if isinstance(ws, dict) else None) or row.get("filename")

        ma = _ws_compare_metrics(a.get("kpis"))
        mb = _ws_compare_metrics(b.get("kpis"))
        deltas = {}
        for key in ma:
            av, bv = ma[key], mb[key]
            deltas[key] = {
                "a": av, "b": bv, "diff": bv - av,
                "pct": ((bv - av) / av * 100) if av else None,
            }

        # Category movement
        def cat_map(row):
            cats = row.get("cats") or []
            m = {}
            if isinstance(cats, list):
                for c in cats:
                    if isinstance(c, dict):
                        nm = c.get("category") or c.get("name")
                        if nm:
                            m[nm] = c
            return m

        cma, cmb = cat_map(a), cat_map(b)
        cat_moves = []
        for nm in sorted(set(cma) | set(cmb)):
            ca, cb = cma.get(nm, {}), cmb.get(nm, {})
            rev_a = float(ca.get("revenue") or 0)
            rev_b = float(cb.get("revenue") or 0)
            cat_moves.append({
                "category": nm,
                "rev_a": rev_a, "rev_b": rev_b, "rev_diff": rev_b - rev_a,
                "st_a": float(ca.get("sell_through") or 0),
                "st_b": float(cb.get("sell_through") or 0),
                "status": "new" if not ca else ("dropped" if not cb else "changed"),
            })
        cat_moves.sort(key=lambda x: abs(x["rev_diff"]), reverse=True)

        # AI-style written commentary (synthesized server-side, no AI cost)
        commentary = []
        rev = deltas.get("revenue", {})
        if rev.get("a") is not None:
            rp = rev.get("pct")
            commentary.append({
                "text": f"Revenue {'grew' if rev['diff'] >= 0 else 'fell'} from "
                        f"£{ma['revenue']:,.0f} to £{mb['revenue']:,.0f}"
                        + (f" ({'+' if rp >= 0 else ''}{rp:.1f}%)." if rp is not None else "."),
                "tone": "good" if rev["diff"] >= 0 else "bad",
            })
        mgd = deltas.get("margin", {})
        # Only comment on margin when BOTH datasets actually carry margin data —
        # otherwise a missing column reads as a false "collapse to 0%".
        _ma_av = bool((a.get("kpis") or {}).get("margin_available", ma.get("margin")))
        _mb_av = bool((b.get("kpis") or {}).get("margin_available", mb.get("margin")))
        if mgd and _ma_av and _mb_av:
            commentary.append({
                "text": f"Margin moved {'up' if mgd['diff'] >= 0 else 'down'} from "
                        f"{ma['margin']:.0f}% to {mb['margin']:.0f}% — "
                        + ("healthy full-price mix." if mgd["diff"] >= 0 else "watch markdown depth."),
                "tone": "good" if mgd["diff"] >= 0 else "warn",
            })
        hd = deltas.get("health", {})
        if hd:
            commentary.append({
                "text": f"Business Health {'improved' if hd['diff'] >= 0 else 'declined'} "
                        f"from {ma['health']:.0f} to {mb['health']:.0f}.",
                "tone": "good" if hd["diff"] >= 0 else "bad",
            })
        if cat_moves:
            up = [c["category"] for c in cat_moves if c["rev_diff"] > 0][:2]
            dn = [c["category"] for c in cat_moves if c["rev_diff"] < 0][:2]
            if up or dn:
                parts = []
                if up: parts.append(f"{', '.join(up)} strengthened")
                if dn: parts.append(f"{', '.join(dn)} weakened")
                commentary.append({
                    "text": " while ".join(parts) + " — rebalance intake toward the winning categories.",
                    "tone": "warn",
                })

        return jsonify({
            "a": {"id": a["id"], "label": label(a), "date": a.get("created_at"), "metrics": ma},
            "b": {"id": b["id"], "label": label(b), "date": b.get("created_at"), "metrics": mb},
            "deltas": deltas,
            "category_movement": cat_moves[:25],
            "commentary": commentary,
        })
    except Exception as e:
        app.logger.error(f"workspace_compare error: {e}", exc_info=True)
        return jsonify({"error": str(e)}), 500


def _entity_metrics_from_sku(rows: dict) -> dict:
    """Merge the per-SKU view across best_sellers / markdown_risk / reorder_alerts
    into a single metrics record (best_sellers is richest so it wins ties)."""
    return {
        "revenue":      _safe_float(rows.get("revenue")),
        "units":        _safe_float(rows.get("units")) if rows.get("units") is not None else None,
        "sell_through": _safe_float(rows.get("sell_through")) if rows.get("sell_through") is not None else None,
        "margin":       _safe_float(rows.get("margin_pct")) if rows.get("margin_pct") is not None else None,
        "stock":        _safe_float(rows.get("stock")) if rows.get("stock") is not None else None,
        "name":         rows.get("name"),
        "category":     rows.get("category"),
    }


def _build_sku_index(dataset: dict) -> dict:
    """sku -> merged metrics record, pooling best_sellers, markdown_risk, reorder_alerts."""
    pool = {}
    for src in ("best_sellers", "markdown_risk", "reorder_alerts"):
        for r in (dataset.get(src) or []):
            if not isinstance(r, dict):
                continue
            sku = r.get("sku")
            if not sku:
                continue
            cur = pool.setdefault(sku, {})
            # best_sellers carries the fullest record, so only fill gaps —
            # never overwrite a non-null value already captured.
            for k, v in r.items():
                if v is not None and cur.get(k) in (None, ""):
                    cur[k] = v
    return {sku: _entity_metrics_from_sku(rec) for sku, rec in pool.items()}


def _build_category_index(dataset: dict) -> dict:
    """category name -> metrics record from category_scorecard."""
    out = {}
    for c in (dataset.get("category_scorecard") or []):
        if not isinstance(c, dict):
            continue
        nm = c.get("category") or c.get("name")
        if not nm:
            continue
        out[nm] = {
            "revenue":      _safe_float(c.get("revenue")),
            "units":        _safe_float(c.get("units")) if c.get("units") is not None else None,
            "sell_through": _safe_float(c.get("avg_sell_through")) if c.get("avg_sell_through") is not None else None,
            "margin":       _safe_float(c.get("avg_margin_pct")) if c.get("avg_margin_pct") is not None else None,
            "sku_count":    c.get("sku_count"),
            "name":         nm,
            "category":     nm,
        }
    return out


@app.route("/api/workspace/compare/entities")
@require_auth
def workspace_compare_entities(user):
    """Searchable index of every SKU and category present in either dataset, with
    last-period (A) vs this-period (B) metrics + deltas. Powers the Compare search
    box so a merchandiser can track whether an action on a specific SKU/category
    actually moved the numbers week-on-week."""
    id_a = request.args.get("a")
    id_b = request.args.get("b")
    if not id_a or not id_b:
        return jsonify({"error": "Two file ids (a, b) required"}), 400
    try:
        db = get_user_db()

        def load(uid):
            r = db.table("uploads") \
                .select("id, filename, created_at, "
                        "ws:analysis->_workspace, kpis:analysis->kpis, "
                        "cats:analysis->category_scorecard, "
                        "best:analysis->best_sellers, "
                        "md:analysis->markdown_risk, "
                        "reorders:analysis->reorder_alerts") \
                .eq("id", uid).eq("user_id", user["id"]).single().execute()
            d = r.data or {}
            return {
                "id": d.get("id"),
                "filename": d.get("filename"),
                "created_at": d.get("created_at"),
                "ws": d.get("ws"),
                "kpis": d.get("kpis") or {},
                "category_scorecard": d.get("cats") or [],
                "best_sellers": d.get("best") or [],
                "markdown_risk": d.get("md") or [],
                "reorder_alerts": d.get("reorders") or [],
            }

        a = load(id_a)
        b = load(id_b)
        if not a.get("id") or not b.get("id"):
            return jsonify({"error": "One or both files not found"}), 404

        def label(row):
            ws = row.get("ws") or {}
            return (ws.get("display_name") if isinstance(ws, dict) else None) or row.get("filename")

        def diff_pct(av, bv):
            if av is None or bv is None:
                return None
            return ((bv - av) / av * 100) if av else None

        def build_entities(idx_a, idx_b):
            out = []
            for key in sorted(set(idx_a) | set(idx_b)):
                ra, rb = idx_a.get(key), idx_b.get(key)
                base = rb or ra
                metrics = {}
                for fld in ("revenue", "units", "sell_through", "margin"):
                    av = ra.get(fld) if ra else None
                    bv = rb.get(fld) if rb else None
                    metrics[fld] = {
                        "a": av, "b": bv,
                        "diff": (bv - av) if (av is not None and bv is not None) else None,
                        "pct": diff_pct(av, bv),
                    }
                out.append({
                    "key": key,
                    "name": base.get("name") or key,
                    "category": base.get("category"),
                    "status": "new" if not ra else ("dropped" if not rb else "changed"),
                    "metrics": metrics,
                })
            return out

        ma_av = bool(a["kpis"].get("margin_available", True))
        mb_av = bool(b["kpis"].get("margin_available", True))

        return jsonify({
            "a": {"id": a["id"], "label": label(a), "date": a.get("created_at")},
            "b": {"id": b["id"], "label": label(b), "date": b.get("created_at")},
            "margin_available": ma_av and mb_av,
            "skus": build_entities(_build_sku_index(a), _build_sku_index(b)),
            "categories": build_entities(_build_category_index(a), _build_category_index(b)),
        })
    except Exception as e:
        app.logger.error(f"workspace_compare_entities error: {e}", exc_info=True)
        return jsonify({"error": str(e)}), 500


@app.route("/api/workspace/timeline")
@require_auth
def workspace_timeline(user):
    """Historical performance trend across all persisted files (oldest → newest)."""
    try:
        res = _ws_fetch_rows(user)
        points = []
        for r in (res.data or []):
            card = _ws_row_to_card(r)
            if card["trashed"]:
                continue
            points.append({
                "id":           card["id"],
                "label":        card["display_name"],
                "date":         card["upload_date"],
                "revenue":      card["revenue"],
                "margin":       card["margin"],
                "sell_through": card["sell_through"],
                "health":       card["health_score"],
                "records":      card["record_count"],
                "key_rec":      _ws_key_rec(card),
            })
        points.sort(key=lambda p: p["date"] or "")

        # Local trend commentary (no AI cost) comparing latest vs previous
        commentary = []
        if len(points) >= 2:
            cur, prev = points[-1], points[-2]
            def trend(name, c, p, unit=""):
                if p == 0:
                    return None
                d = c - p
                pct = d / p * 100
                arrow = "up" if d > 0 else ("down" if d < 0 else "flat")
                return {"metric": name, "direction": arrow,
                        "change": round(d, 1), "pct": round(pct, 1), "unit": unit}
            for nm, key, unit in [("Revenue", "revenue", "£"),
                                  ("Margin", "margin", "%"),
                                  ("Sell-Through", "sell_through", "%"),
                                  ("Health Score", "health", "")]:
                t = trend(nm, cur[key], prev[key], unit)
                if t:
                    commentary.append(t)

        return jsonify({"points": points, "commentary": commentary,
                        "count": len(points)})
    except Exception as e:
        app.logger.error(f"workspace_timeline error: {e}", exc_info=True)
        return jsonify({"error": str(e)}), 500


@app.route("/api/workspace/command")
@require_auth
def workspace_command(user):
    """
    Decision Mode — the Command Centre.
    Answers 'What should I do next?' for a dataset: health score, top 5 ranked
    actions (with revenue/margin impact + confidence), biggest opportunity/risk,
    a decision simulator, commercial impact, and the Business Memory engine
    (cross-dataset comparison vs the previous upload).
    """
    try:
        upload_id = request.args.get("upload_id") or session.get("current_upload_id")
        analysis, upload_id = _get_analysis_for_user(user, upload_id)
        if not analysis or not analysis.get("kpis"):
            return jsonify({"error": "No analysis available for this dataset"}), 404

        kpis  = analysis.get("kpis") or {}
        role  = _build_role_analysis(analysis)
        ws    = analysis.get(WORKSPACE_META_KEY) or {}
        label = (ws.get("display_name") if isinstance(ws, dict) else None) \
                or analysis.get("filename") or "Dataset"
        health = _ws_health_score(kpis)
        health_detail = _health_breakdown(kpis)

        # Top 5 ranked actions from the role engine's ai_actions
        actions = []
        for a in (role.get("ai_actions") or [])[:5]:
            mi = a.get("margin_impact", 0)
            actions.append({
                "rank":          a.get("rank"),
                "action":        f"{a.get('type','Action')} — {a.get('sku','')}".strip(" —"),
                "reason":        a.get("description", ""),
                "rev_impact":    a.get("revenue_impact", 0),
                "margin_impact": f"{mi:+.1f}%" if isinstance(mi, (int, float)) else str(mi),
                "confidence":    a.get("confidence", 70),
                "reasoning":     a.get("reasoning"),
            })

        def first_item(lst):
            if isinstance(lst, list) and lst:
                v = lst[0]
                if isinstance(v, dict):
                    return {"t": v.get("title") or v.get("text") or "", "s": v.get("detail") or "", "impact": v.get("impact") or ""}
                return {"t": str(v), "s": "", "impact": ""}
            return {"t": "—", "s": "", "impact": ""}

        top_opp  = first_item(role.get("top_opportunities"))
        top_risk = first_item(role.get("top_risks"))

        # Commercial impact roll-up
        rev_up   = sum(max(0, a.get("revenue_impact", 0)) for a in (role.get("ai_actions") or []))
        rev_risk = abs(sum(min(0, a.get("revenue_impact", 0)) for a in (role.get("ai_actions") or [])))
        commercial = {
            "revenue_at_stake": round(rev_risk or kpis.get("markdown_revenue", 0) or 0),
            "margin_protect":   round((kpis.get("total_revenue", 0) or 0) * 0.04),
            "total_upside":     round(rev_up),
        }

        # ── Business Memory: compare vs previous upload ──────────────────
        memory = {"improved": [], "declined": [], "new_risks": [], "resolved_risks": [],
                  "new_opportunities": [], "recurring": [], "recommended": []}
        prev_metrics = None
        try:
            rows = (_ws_fetch_rows(user).data or [])
            ordered = [r for r in rows if not ((r.get("ws") or {}) if isinstance(r.get("ws"), dict) else {}).get("trashed")]
            ordered.sort(key=lambda r: r.get("created_at") or "")
            ids = [r.get("id") for r in ordered]
            prev_analysis = None
            if upload_id in ids:
                i = ids.index(upload_id)
                if i > 0:
                    prev_analysis, _ = _get_analysis_for_user(user, ids[i - 1])
            if prev_analysis:
                pk = prev_analysis.get("kpis") or {}
                cur_m, prev_m = _ws_compare_metrics(kpis), _ws_compare_metrics(pk)
                prev_metrics = prev_m
                labels = {"revenue": "Revenue", "margin": "Margin", "sell_through": "Sell-through", "health": "Business Health"}
                for k, lbl in labels.items():
                    diff = cur_m[k] - prev_m[k]
                    if abs(diff) < 0.5:
                        continue
                    if k == "revenue":
                        txt = f"{lbl} {'+' if diff>=0 else ''}£{diff:,.0f}"
                    else:
                        txt = f"{lbl} {'+' if diff>=0 else ''}{diff:.1f}{'%' if k!='health' else ' pts'}"
                    (memory["improved"] if diff >= 0 else memory["declined"]).append(txt)

                prev_role  = _build_role_analysis(prev_analysis)
                cur_risks  = {str(x.get('title') if isinstance(x, dict) else x) for x in (role.get("top_risks") or [])}
                prev_risks = {str(x.get('title') if isinstance(x, dict) else x) for x in (prev_role.get("top_risks") or [])}
                memory["new_risks"]      = sorted(cur_risks - prev_risks)[:3]
                memory["resolved_risks"] = sorted(prev_risks - cur_risks)[:3]
                recurring = sorted(cur_risks & prev_risks)
                memory["recurring"] = [{"issue": r, "weeks": 2, "note": "Seen in the last two uploads — prioritise action."} for r in recurring[:2]]
                memory["new_opportunities"] = [str(x.get('title') if isinstance(x, dict) else x) for x in (role.get("top_opportunities") or [])][:3]
        except Exception as me:
            app.logger.warning(f"workspace_command memory error: {me}")
        memory["recommended"] = [a["action"] for a in actions[:2]]

        return jsonify({
            "upload_id":     upload_id,
            "dataset_label": label,
            "health_score":  health,
            "health_label":  ("Healthy" if health >= 80 else "Watch" if health >= 60 else "At Risk"),
            "health_breakdown": health_detail,
            "categories":    _category_cards(analysis),
            "revenue":       kpis.get("total_revenue", 0),
            "margin":        kpis.get("avg_margin_pct", 0),
            "sell_through":  kpis.get("avg_sell_through", 0),
            "exec_summary":  role.get("executive_summary", ""),
            "top_opportunity": top_opp,
            "top_risk":      top_risk,
            "actions":       actions,
            "simulator": {
                "base_revenue": kpis.get("total_revenue", 0),
                "base_margin":  kpis.get("avg_margin_pct", 0),
                "base_health":  health,
                "base_stock":   _ws_stock_value(kpis),
                "levers": [
                    {"key": "markdown", "label": "Markdown depth",              "min": 0, "max": 40, "step": 5, "default": 20, "unit": "%"},
                    {"key": "reorder",  "label": "Replenishment / reorder depth", "min": 0, "max": 50, "step": 5, "default": 25, "unit": "%"},
                    {"key": "po_cancel", "label": "Purchase order cancellation",  "min": 0, "max": 50, "step": 5, "default": 0,  "unit": "%"},
                    {"key": "intake",   "label": "Future intake reduction",      "min": 0, "max": 40, "step": 5, "default": 10, "unit": "%"},
                ],
            },
            "commercial_impact": commercial,
            "memory": memory,
            "trading_meeting": _trading_meeting(
                analysis, role, kpis, health, health_detail, label, prev_metrics),
        })
    except Exception as e:
        app.logger.error(f"workspace_command error: {e}", exc_info=True)
        return jsonify({"error": str(e)}), 500


@app.route("/api/workspace/category")
@require_auth
def workspace_category(user):
    """Category drill-down — overview KPIs, top sellers, slow movers, actions."""
    try:
        category  = request.args.get("category", "").strip()
        if not category:
            return jsonify({"error": "category is required"}), 400
        upload_id = request.args.get("upload_id") or session.get("current_upload_id")
        analysis, upload_id = _get_analysis_for_user(user, upload_id)
        if not analysis or not analysis.get("kpis"):
            return jsonify({"error": "No analysis available for this dataset"}), 404
        detail = _category_detail(analysis, category)
        if not detail:
            return jsonify({"error": f"Category '{category}' not found"}), 404
        return jsonify(detail)
    except Exception as e:
        app.logger.error(f"workspace_category error: {e}", exc_info=True)
        return jsonify({"error": str(e)}), 500


def _get_analysis_for_user(user, upload_id=None):
    """Helper to fetch analysis JSON from Supabase."""
    try:
        db = get_user_db()
        uid = upload_id or session.get("current_upload_id")
        company_id = user.get("company_id")  # keep as None if not set — don't cast to string
        if uid:
            q = db.table("uploads").select("analysis, id") \
                .eq("id", uid).eq("user_id", user["id"])
            if company_id is not None:
                q = q.eq("company_id", str(company_id))
            res = q.single().execute()
            if res.data:
                return res.data["analysis"], res.data["id"]

        q = db.table("uploads").select("analysis, id") \
            .eq("user_id", user["id"]).order("created_at", desc=True).limit(1)
        if company_id is not None:
            q = q.eq("company_id", str(company_id))
        res = q.execute()
        if res.data:
            return res.data[0]["analysis"], res.data[0]["id"]
    except Exception:
        pass
    return None, None


def _company_private_mode(company_id) -> bool:
    """Check if a company has Private Processing Mode enabled."""
    if not company_id:
        return False
    try:
        res = supabase.table("companies").select("private_processing_mode") \
            .eq("id", str(company_id)).single().execute()
        if res.data:
            return bool(res.data.get("private_processing_mode", False))
    except Exception:
        pass
    return False


def _increment_ai_usage(user):
    get_user_db().table("profiles").update(
        {"ai_messages_used": (user.get("ai_messages_used", 0) + 1)}
    ).eq("id", user["id"]).execute()


@app.route("/api/chat", methods=["POST"])
@require_auth
def chat(user):
    allowed, msg = check_plan_limit(user, "ai_msg")
    if not allowed:
        return jsonify({"error": msg, "upgrade_required": True}), 402

    body = request.get_json(silent=True) or {}
    user_msg  = (body.get("message") or "").strip()
    upload_id = body.get("upload_id")
    if not user_msg:
        return jsonify({"error": "Empty message"}), 400

    analysis, uid = _get_analysis_for_user(user, upload_id)
    if not analysis:
        return jsonify({"error": "No data uploaded yet. Please upload a trading file first."}), 404

    private_mode = _company_private_mode(user.get("company_id"))
    masked_analysis = prepare_for_ai(analysis, private_mode=private_mode)

    # Fetch recent chat history from Supabase
    db = get_user_db()
    try:
        hist_res = db.table("chat_messages") \
            .select("role, content").eq("user_id", user["id"]) \
            .eq("upload_id", uid).order("created_at", desc=False).limit(20).execute()
        history = [{"role": r["role"], "content": r["content"]} for r in (hist_res.data or [])]
    except Exception:
        history = []

    history.append({"role": "user", "content": user_msg})

    try:
        response = ai_client.messages.create(
            model=MODEL, max_tokens=1024,
            system=build_system_prompt(masked_analysis),
            messages=history,
        )
        reply = response.content[0].text

        # Persist both messages
        db.table("chat_messages").insert([
            {"user_id": user["id"], "upload_id": uid, "role": "user",      "content": user_msg},
            {"user_id": user["id"], "upload_id": uid, "role": "assistant", "content": reply},
        ]).execute()

        _increment_ai_usage(user)
        audit.log(AI_QUERY, request=request, user_id=user["id"],
                  company_id=user.get("company_id"),
                  resource=f"upload:{uid}",
                  metadata={"private_mode": private_mode})
        return jsonify({"reply": reply})

    except anthropic.AuthenticationError as e:
        app.logger.error(f"Anthropic auth error (bad API key?): {e}")
        return jsonify({"error": "AI configuration error — please contact support."}), 502
    except anthropic.APIError as e:
        app.logger.error(f"Anthropic API error: {e}")
        return jsonify({"error": f"AI service error: {e}"}), 502
    except Exception as e:
        app.logger.error(f"Chat unexpected error: {e}")
        return jsonify({"error": f"Unexpected error: {e}"}), 500


@app.route("/api/chat/stream", methods=["POST"])
@require_auth
def chat_stream(user):
    allowed, msg = check_plan_limit(user, "ai_msg")
    if not allowed:
        def err():
            yield f"data: {json.dumps({'error': msg, 'upgrade_required': True})}\n\n"
        return Response(err(), mimetype="text/event-stream")

    body = request.get_json(silent=True) or {}
    user_msg  = (body.get("message") or "").strip()
    upload_id = body.get("upload_id")

    analysis, uid = _get_analysis_for_user(user, upload_id)
    if not analysis:
        def err():
            yield f"data: {json.dumps({'error': 'No data uploaded yet'})}\n\n"
        return Response(err(), mimetype="text/event-stream")

    private_mode = _company_private_mode(user.get("company_id"))
    masked_analysis = prepare_for_ai(analysis, private_mode=private_mode)

    stream_db = get_user_db()
    try:
        hist_res = stream_db.table("chat_messages") \
            .select("role, content").eq("user_id", user["id"]) \
            .eq("upload_id", uid).order("created_at", desc=False).limit(20).execute()
        history = [{"role": r["role"], "content": r["content"]} for r in (hist_res.data or [])]
    except Exception:
        history = []

    history.append({"role": "user", "content": user_msg})

    def generate():
        full_reply = []
        try:
            with ai_client.messages.stream(
                model=MODEL, max_tokens=1024,
                system=build_system_prompt(masked_analysis),
                messages=history,
            ) as stream:
                for text in stream.text_stream:
                    full_reply.append(text)
                    yield f"data: {json.dumps({'delta': text})}\n\n"

            reply = "".join(full_reply)
            stream_db.table("chat_messages").insert([
                {"user_id": user["id"], "upload_id": uid, "role": "user",      "content": user_msg},
                {"user_id": user["id"], "upload_id": uid, "role": "assistant", "content": reply},
            ]).execute()
            _increment_ai_usage(user)
            audit.log(AI_QUERY, request=request, user_id=user["id"],
                      company_id=user.get("company_id"),
                      resource=f"upload:{uid}",
                      metadata={"private_mode": private_mode})
            yield f"data: {json.dumps({'done': True})}\n\n"

        except Exception as e:
            app.logger.error(f"Stream error: {e}")
            yield f"data: {json.dumps({'error': str(e)})}\n\n"

    return Response(generate(), mimetype="text/event-stream",
                    headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})


@app.route("/api/quick-insights")
@require_auth
def quick_insights(user):
    analysis, uid = _get_analysis_for_user(user)
    if not analysis:
        return jsonify({"insights": [
            "Upload your trading CSV or Excel to unlock AI insights.",
            "Track sell-through, reorder alerts, and markdown risk in one place.",
            "Ask the AI anything about your range — grounded in your actual numbers.",
        ]})

    private_mode = _company_private_mode(user.get("company_id"))
    masked_analysis = prepare_for_ai(analysis, private_mode=private_mode)

    k = masked_analysis["kpis"]
    reorders  = masked_analysis.get("reorder_alerts", [])
    markdowns = masked_analysis.get("markdown_risk", [])
    cats      = masked_analysis.get("category_scorecard", [])

    prompt = f"""Based on this trading data, generate exactly 3 concise action-oriented insights for a fashion buyer. Each one sentence, starting with an action verb. Return as a JSON array of strings only.

Revenue £{k.get('total_revenue',0):,.0f} | ST {k.get('avg_sell_through',0):.1f}% | Margin {k.get('avg_margin_pct',0):.1f}%
{k.get('reorder_count',0)} reorder alerts, {k.get('critical_oos_count',0)} OOS
Top reorder: {reorders[0] if reorders else 'none'}
Top markdown: {markdowns[0] if markdowns else 'none'}
Best category: {cats[0] if cats else 'none'}
Weakest category: {cats[-1] if len(cats)>1 else 'none'}

Return ONLY a JSON array, no markdown."""

    try:
        res = ai_client.messages.create(
            model=MODEL, max_tokens=300,
            messages=[{"role": "user", "content": prompt}],
        )
        insights = json.loads(res.content[0].text.strip())
        audit.log(AI_QUERY, request=request, user_id=user["id"],
                  company_id=user.get("company_id"),
                  resource=f"upload:{uid}",
                  metadata={"private_mode": private_mode})
        return jsonify({"insights": insights[:3]})
    except Exception:
        return jsonify({"insights": [
            f"Raise POs immediately for {k.get('critical_oos_count',0)} out-of-stock SKUs before revenue is lost.",
            f"Review {k.get('markdown_risk_count',0)} slow-selling SKUs for tactical markdowns to protect margin.",
            f"Overall sell-through of {k.get('avg_sell_through',0):.1f}% indicates overbought range — tighten next season's OTB.",
        ]})


# ═══════════════════════════════════════════════════════════════
# BILLING — plan info + Stripe Checkout
# ═══════════════════════════════════════════════════════════════

# Display metadata for the upgrade modal. Hard limits live in PLAN_LIMITS; this
# is purely the marketing copy shown to the user.
BILLING_PLANS = {
    "starter": {"name": "Starter", "price": "£49/mo",
                "skus": "200 SKUs / upload", "ai": "50 Copilot questions / mo",
                "blurb": "For solo merchandisers getting started."},
    "growth":  {"name": "Growth",  "price": "£149/mo",
                "skus": "Unlimited SKUs", "ai": "Unlimited Copilot",
                "blurb": "For busy merchandisers and freelancers across clients."},
    "studio":  {"name": "Studio",  "price": "£349/mo",
                "skus": "Unlimited SKUs", "ai": "Unlimited Copilot",
                "blurb": "For consultants and small teams.", "best": True},
}


def _trial_state(user):
    """Compute trial/subscription state for the frontend (display only — hard
    limits are still enforced server-side in check_plan_limit / upload / chat)."""
    plan = (user or {}).get("plan", "trial")
    trial_ends = (user or {}).get("trial_ends")
    is_paid = plan in ("starter", "growth", "studio")
    days_left, expired = None, False
    if trial_ends and not is_paid:
        try:
            end = datetime.fromisoformat(str(trial_ends).replace("Z", "+00:00"))
            now = datetime.now(end.tzinfo) if end.tzinfo else datetime.utcnow()
            days_left = max(0, (end - now).days + (1 if (end - now).seconds else 0))
            expired = end <= now
        except Exception:
            days_left = None
    return {
        "plan": plan,
        "is_paid": is_paid,
        "trial_ends": trial_ends,
        "days_left": days_left,
        "expired": bool(expired and not is_paid),
        "premium_unlocked": is_paid,
    }


@app.route("/api/billing/plans")
@require_auth
def billing_plans(user):
    """Plan catalogue + the caller's current trial/subscription state."""
    return jsonify({"plans": BILLING_PLANS, "state": _trial_state(user)})


@app.route("/api/billing/checkout", methods=["POST"])
@require_auth
def billing_checkout(user):
    """Create a Stripe Checkout Session for the chosen plan and return its URL.
    The user completes payment on Stripe's hosted page; the webhook upgrades the
    profile. We never handle card details ourselves."""
    if not stripe.api_key:
        return jsonify({"error": "Billing is not configured yet. Please contact support to upgrade."}), 503
    body = request.get_json(silent=True) or {}
    plan = (body.get("plan") or "").strip().lower()
    price_env = {"starter": "STRIPE_PRICE_STARTER",
                 "growth":  "STRIPE_PRICE_GROWTH",
                 "studio":  "STRIPE_PRICE_STUDIO"}.get(plan)
    price_id = os.environ.get(price_env, "") if price_env else ""
    if not price_id:
        return jsonify({"error": "That plan isn't available for self-serve checkout yet."}), 400
    base = request.host_url.rstrip("/")
    try:
        kwargs = {
            "mode": "subscription",
            "line_items": [{"price": price_id, "quantity": 1}],
            "success_url": f"{base}/dashboard?upgraded=1",
            "cancel_url": f"{base}/dashboard?upgrade_cancelled=1",
            "client_reference_id": user["id"],
            "allow_promotion_codes": True,
        }
        cust = user.get("stripe_customer_id")
        if cust:
            kwargs["customer"] = cust
        elif user.get("email"):
            kwargs["customer_email"] = user["email"]
        sess = stripe.checkout.Session.create(**kwargs)
        return jsonify({"url": sess.url})
    except Exception as e:
        app.logger.error(f"checkout error: {e}")
        return jsonify({"error": "Could not start checkout. Please try again."}), 500


# ═══════════════════════════════════════════════════════════════
# STRIPE WEBHOOKS
# ═══════════════════════════════════════════════════════════════

@app.route("/webhooks/stripe", methods=["POST"])
def stripe_webhook():
    payload = request.get_data()
    sig     = request.headers.get("Stripe-Signature", "")

    # Verify the signature against the raw bytes, then parse the raw JSON
    # ourselves. The webhook payload is already plain JSON, so json.loads gives
    # a fully-nested plain dict — avoiding StripeObject quirks where .get() is
    # misread as a field lookup (KeyError('get')).
    try:
        stripe.Webhook.construct_event(payload, sig, STRIPE_WEBHOOK_SECRET)
    except Exception as e:
        return jsonify({"error": str(e)}), 400

    try:
        event = json.loads(payload)
    except Exception as e:
        return jsonify({"error": f"payload parse: {e}"}), 400

    try:
        return _process_stripe_event(event)
    except Exception as e:
        import traceback as _tb
        app.logger.error(f"stripe_webhook error: {e}\n{_tb.format_exc()}")
        # Return 500 so Stripe retries; the handler is idempotent.
        return jsonify({"error": "internal"}), 500


def _process_stripe_event(event):
    # `event` is a plain nested dict (parsed via json.loads in the caller), so
    # dict.get() works throughout — avoiding StripeObject .get() quirks.

    # Idempotency: skip only if we've already *successfully* processed this
    # event. A row that exists but is still processed=false means a prior
    # attempt failed part-way, so we must allow it to be retried.
    try:
        existing = supabase.table("stripe_events").select("processed").eq("id", event["id"]).execute()
        if existing.data and existing.data[0].get("processed"):
            return jsonify({"ok": True})
        supabase.table("stripe_events").upsert(
            {"id": event["id"], "type": event["type"], "processed": False}
        ).execute()
    except Exception:
        pass

    etype = event["type"]
    data  = event["data"]["object"]

    PRICE_TO_PLAN = {
        os.environ.get("STRIPE_PRICE_STARTER", ""): "starter",
        os.environ.get("STRIPE_PRICE_GROWTH",  ""): "growth",
        os.environ.get("STRIPE_PRICE_STUDIO",  ""): "studio",
    }

    if etype == "checkout.session.completed":
        customer_id = data.get("customer")
        sub_id      = data.get("subscription")
        email       = (data.get("customer_details") or {}).get("email", "") or ""
        ref_id      = data.get("client_reference_id")

        # Stripe does NOT include line_items in the event payload by default,
        # so fetch the purchased price explicitly to know which plan was bought.
        price_id = ""
        try:
            items = stripe.checkout.Session.list_line_items(data["id"], limit=1)
            if items and items.data:
                price_id = items.data[0].price.id
        except Exception as e:
            app.logger.error(f"line_items fetch error: {e}")
        # Fallback: read the price off the subscription itself.
        if not price_id and sub_id:
            try:
                sub = stripe.Subscription.retrieve(sub_id)
                price_id = sub["items"]["data"][0]["price"]["id"]
            except Exception as e:
                app.logger.error(f"subscription fetch error: {e}")

        plan = PRICE_TO_PLAN.get(price_id, "starter")

        update = {
            "plan": plan,
            "stripe_customer_id": customer_id,
            "stripe_subscription_id": sub_id,
        }
        # Prefer the reliable client_reference_id (the user's profile id);
        # fall back to email only if it's missing.
        try:
            if ref_id:
                supabase.table("profiles").update(update).eq("id", ref_id).execute()
            elif email:
                supabase.table("profiles").update(update).eq("email", email.lower()).execute()
            else:
                app.logger.error("Webhook: no client_reference_id or email to match profile")
        except Exception as e:
            app.logger.error(f"Webhook profile update error: {e}")

    elif etype in ("customer.subscription.updated", "customer.subscription.deleted"):
        customer_id = data.get("customer")
        status      = data.get("status")
        plan        = "trial" if (etype == "customer.subscription.deleted" or status == "canceled") else \
                      PRICE_TO_PLAN.get(data.get("items", {}).get("data", [{}])[0].get("price", {}).get("id", ""), "starter")
        try:
            supabase.table("profiles").update({"plan": plan}) \
                .eq("stripe_customer_id", customer_id).execute()
        except Exception as e:
            app.logger.error(f"Webhook subscription error: {e}")

    try:
        supabase.table("stripe_events").update({"processed": True}).eq("id", event["id"]).execute()
    except Exception as e:
        app.logger.error(f"stripe_events mark processed error: {e}")
    return jsonify({"ok": True})


# ═══════════════════════════════════════════════════════════════
# AI ACTION CENTRE ENGINE
# ═══════════════════════════════════════════════════════════════

def _run_action_engine(analysis):
    """
    Commercial action prioritisation engine.
    Scans analysis data and generates ranked actions across 6 types:
    Markdown | Reorder | Intake Reduction | PO Cancellation | Stock Transfer | Supplier Review
    """
    kpis         = analysis.get("kpis", {})
    best_sellers = analysis.get("best_sellers", [])
    reorders     = analysis.get("reorder_alerts", [])
    markdowns    = analysis.get("markdown_risk", [])
    categories   = analysis.get("category_scorecard", [])

    total_rev    = max(kpis.get("total_revenue", 1), 1)
    avg_margin   = kpis.get("avg_margin_pct", 50)
    avg_st       = kpis.get("avg_sell_through", 35)
    avg_cover    = kpis.get("avg_cover_weeks", 8)
    total_units  = kpis.get("total_units", 1)

    actions = []
    aid = 0

    # ── 1. MARKDOWN RECOMMENDATIONS ──────────────────────────────
    # Trigger: sell_through < 30%, low sell-through signals overstock
    for item in markdowns[:10]:
        aid += 1
        st       = item.get("sell_through", 0)
        revenue  = item.get("revenue", 0)
        sku      = item.get("sku", "Unknown")
        cat      = item.get("category", "General") or "General"
        severity = item.get("severity", "TACTICAL")
        rec_d    = item.get("recommended_depth", "15–20%")

        depth    = 45 if severity == "CLEARANCE" else 30 if severity == "DEEP" else 18
        stock_val        = revenue * (1 - avg_margin / 100) * max(1, (100 - st) / 50)
        cash_recovery    = round(stock_val * (depth / 100) * 0.82)
        rev_impact       = round(revenue * (1 - st / 100) * 0.68)
        margin_impact    = round(-(depth * 0.55), 1)
        confidence       = 92 if severity == "CLEARANCE" else 84 if severity == "DEEP" else 72

        rev_w   = min(30, revenue / total_rev * 3000)
        st_w    = (30 - min(st, 30)) / 30 * 35
        urg_w   = 22 if severity == "CLEARANCE" else 14 if severity == "DEEP" else 7
        priority_score   = min(99, int(rev_w + st_w + urg_w + confidence * 0.08))
        priority_label   = "Critical" if severity == "CLEARANCE" else "High" if severity == "DEEP" else "Medium"

        # Estimated weeks cover derived from ST
        est_cover = max(1, round(100 / max(st, 0.5) * 4, 1))

        actions.append({
            "id": aid, "type": "Markdown",
            "priority_label": priority_label, "priority_score": priority_score,
            "title": f"{cat} Overstock — {sku}",
            "sku": sku, "category": cat,
            "issue": (
                f"Sell-through at {st:.0f}% — well below the 30% markdown trigger threshold. "
                f"Estimated {est_cover} weeks cover remaining at current rate of sale."
            ),
            "root_cause": (
                f"Original intake volume exceeded consumer demand in {cat}. "
                f"Rate of sale has not recovered to plan, and the stock position is classified as {severity.title()}. "
                f"Continued full-price trading will result in further cover deterioration."
            ),
            "recommended_action": (
                f"Apply {rec_d} markdown to {sku} immediately. "
                f"Prioritise digital/online channel for maximum velocity. "
                f"Review in-store placement and consider promotional support to accelerate clearance."
            ),
            "expected_outcome": (
                f"Accelerate sell-through from {st:.0f}% to a projected 70–85% within 4–6 weeks. "
                f"Recover £{cash_recovery:,} in cash. Reduce cover to a sustainable 4–6 week position."
            ),
            "business_impact": (
                f"Releases £{cash_recovery:,} working capital currently trapped in slow-moving stock. "
                f"Prevents the need for deeper, more damaging markdowns later in the season. "
                f"Clears floor and warehouse space for higher-margin intake."
            ),
            "revenue_impact": rev_impact,
            "margin_impact": margin_impact,
            "cash_recovery": cash_recovery,
            "confidence": confidence,
        })

    # ── 2. REORDER RECOMMENDATIONS ──────────────────────────────
    # Trigger: sell_through > 70%, weeks_cover < 4, stockout risk
    for item in reorders[:8]:
        aid += 1
        sku      = item.get("sku", "Unknown")
        cat      = item.get("category", "General") or "General"
        revenue  = item.get("revenue", 0)
        urgency  = item.get("urgency", "REORDER")
        cover    = item.get("cover_weeks") or 3
        st       = item.get("sell_through") or 72

        wkly_rev     = revenue / max(cover, 0.5)
        missed_rev   = round(wkly_rev * 8 * 0.62)    # 8 weeks of missed revenue
        cash_recovery = 0                               # investment, not recovery
        margin_impact = round(avg_margin, 1)            # maintain current margin

        confidence   = 93 if urgency == "CRITICAL" else 86 if urgency == "URGENT" else 74
        cover_w      = max(0, (4 - min(cover, 4)) / 4 * 38)
        rev_w        = min(22, revenue / total_rev * 2200)
        urg_w        = 26 if urgency == "CRITICAL" else 15 if urgency == "URGENT" else 5
        priority_score = min(99, int(rev_w + cover_w + urg_w + confidence * 0.08))
        priority_label = "Critical" if urgency == "CRITICAL" else "High" if urgency == "URGENT" else "Medium"

        time_str = "days" if urgency == "CRITICAL" else "1–2 weeks" if urgency == "URGENT" else "3–4 weeks"
        actions.append({
            "id": aid, "type": "Reorder",
            "priority_label": priority_label, "priority_score": priority_score,
            "title": f"Stockout Risk — {sku}",
            "sku": sku, "category": cat,
            "issue": (
                f"Stock cover at {cover:.1f} weeks with sell-through at {st:.0f}%. "
                f"Current rate of sale will exhaust stock within {time_str}."
            ),
            "root_cause": (
                f"Strong consumer demand has depleted {sku} inventory faster than forecast. "
                f"No replenishment is currently scheduled. "
                f"Cover has fallen below the 4-week minimum trigger threshold."
            ),
            "recommended_action": (
                f"Place emergency replenishment order for {sku} to restore 6–8 week cover. "
                f"Engage supplier to expedite lead time. "
                f"Allocate available stock to highest-revenue channels in the interim."
            ),
            "expected_outcome": (
                f"Prevent stockout and recover £{missed_rev:,} in projected lost revenue. "
                f"Maintain full-price sell-through position and protect {cat} category performance."
            ),
            "business_impact": (
                f"Stockout prevention protects £{missed_rev:,} revenue at full margin ({avg_margin:.0f}%). "
                f"Avoids lost customer trust and market share in {cat}. "
                f"Maintains category momentum heading into remainder of season."
            ),
            "revenue_impact": missed_rev,
            "margin_impact": margin_impact,
            "cash_recovery": cash_recovery,
            "confidence": confidence,
        })

    # ── 3. INTAKE REDUCTION ──────────────────────────────────────
    # Trigger: category with health AT RISK / CRITICAL and ST < 25%
    for cat_data in categories:
        health  = cat_data.get("health", "")
        cat_st  = cat_data.get("avg_sell_through")
        cat_mg  = cat_data.get("avg_margin_pct")
        cat     = cat_data.get("category", "Unknown")
        cat_rev = cat_data.get("revenue", 0)

        if health not in ("AT RISK", "CRITICAL"):
            continue
        if cat_st is None or cat_st > 28:
            continue

        aid += 1
        est_cover       = max(12, round(100 / max(cat_st, 0.5) * 4))
        excess_cover    = max(0, est_cover - 12)
        intake_saving   = round(cat_rev * 0.22)
        cash_recovery   = round(intake_saving * 0.18)    # working capital preserved
        margin_impact   = round((cat_mg or avg_margin) * 0.08, 1)
        confidence      = 79 if health == "CRITICAL" else 68

        rev_w      = min(18, cat_rev / total_rev * 1800)
        st_w       = (28 - min(cat_st, 28)) / 28 * 28
        priority_score = min(88, int(rev_w + st_w + confidence * 0.25))
        priority_label = "High" if health == "CRITICAL" else "Medium"

        actions.append({
            "id": aid, "type": "Intake Reduction",
            "priority_label": priority_label, "priority_score": priority_score,
            "title": f"Reduce {cat} Forward Buying",
            "sku": "Category Level", "category": cat,
            "issue": (
                f"{cat} category health is {health} with {cat_st:.0f}% sell-through. "
                f"Estimated {est_cover}+ weeks cover — exceeds the 20-week intake reduction trigger."
            ),
            "root_cause": (
                f"Buy volume in {cat} has consistently outpaced consumer demand. "
                f"Forward buying commitments risk compounding the existing overstock position "
                f"if not addressed before intake lands."
            ),
            "recommended_action": (
                f"Reduce forward buying in {cat} by 20–25% for next season's intake. "
                f"Review all open and uncommitted orders. "
                f"Identify and cancel or defer non-essential intake to bring cover to 8–12 week target."
            ),
            "expected_outcome": (
                f"Prevent £{intake_saving:,} of excess stock accumulation. "
                f"Bring {cat} cover into the sustainable 8–12 week target range. "
                f"Preserve £{cash_recovery:,} in OTB budget for higher-performing categories."
            ),
            "business_impact": (
                f"Reduces cash tied in slow-moving {cat} stock. "
                f"Protects OTB budget and enables reallocation to stronger categories. "
                f"Improves overall business sell-through rate and margin efficiency."
            ),
            "revenue_impact": 0,
            "margin_impact": margin_impact,
            "cash_recovery": cash_recovery,
            "confidence": confidence,
        })

    # ── 4. PO CANCELLATION ──────────────────────────────────────
    # Trigger: CRITICAL health categories with ST < 18%
    for cat_data in categories:
        health  = cat_data.get("health", "")
        cat_st  = cat_data.get("avg_sell_through")
        cat     = cat_data.get("category", "Unknown")
        cat_rev = cat_data.get("revenue", 0)

        if health != "CRITICAL" or (cat_st is not None and cat_st >= 18):
            continue

        aid += 1
        po_value      = round(cat_rev * 0.14)
        cash_recovery = po_value
        confidence    = 73

        rev_w       = min(16, cat_rev / total_rev * 1600)
        st_w        = (18 - min(cat_st or 10, 18)) / 18 * 28
        priority_score = min(84, int(rev_w + st_w + confidence * 0.25))

        actions.append({
            "id": aid, "type": "PO Cancellation",
            "priority_label": "High", "priority_score": priority_score,
            "title": f"Cancel Incoming {cat} Purchase Orders",
            "sku": "Category Level", "category": cat,
            "issue": (
                f"{cat} current stock cannot realistically clear before new intake is scheduled to arrive. "
                f"Sell-through at {cat_st or 0:.0f}% with {health} health status."
            ),
            "root_cause": (
                f"Existing overstock in {cat} combined with open purchase orders will create a compounding "
                f"inventory problem. Incoming stock will depress sell-through further and require deeper markdowns."
            ),
            "recommended_action": (
                f"Immediately review all open and uncommitted {cat} purchase orders. "
                f"Cancel or defer any non-confirmed POs within cancellation windows. "
                f"Engage supplier to negotiate deferral on committed orders."
            ),
            "expected_outcome": (
                f"Prevent £{po_value:,} of additional unnecessary stock entering the business. "
                f"Protect cash flow and OTB budget. Reduce forward markdown exposure."
            ),
            "business_impact": (
                f"Cancelling or deferring {cat} POs saves £{po_value:,} in cash outflow. "
                f"Avoids compounding existing overstock and prevents further margin erosion. "
                f"Frees OTB for reallocation to higher-performing categories."
            ),
            "revenue_impact": 0,
            "margin_impact": 0,
            "cash_recovery": cash_recovery,
            "confidence": confidence,
        })

    # ── 5. STOCK TRANSFER ────────────────────────────────────────
    # Identify rebalancing: high-stock vs low-stock categories
    low_cats  = [c for c in categories if (c.get("avg_sell_through") or 0) > 55
                 and c.get("health") in ("STRONG", "HEALTHY")]
    high_cats = [c for c in categories if (c.get("avg_sell_through") or 100) < 28
                 and c.get("health") in ("AT RISK", "CRITICAL")]

    if low_cats and high_cats:
        from_c = high_cats[0]
        to_c   = low_cats[0]
        aid   += 1
        transfer_val   = round(min(from_c.get("revenue", 0), to_c.get("revenue", 0)) * 0.09)
        rev_impact     = round(transfer_val * 0.55)
        cash_recovery  = 0
        confidence     = 66

        priority_score = min(74, int(min(14, transfer_val / total_rev * 1400) + confidence * 0.42))

        actions.append({
            "id": aid, "type": "Stock Transfer",
            "priority_label": "Medium", "priority_score": priority_score,
            "title": f"Rebalance: {from_c['category']} → {to_c['category']}",
            "sku": "Cross-Category",
            "category": f"{from_c['category']} → {to_c['category']}",
            "issue": (
                f"{from_c['category']} is overstocked at {from_c.get('avg_sell_through', 0):.0f}% sell-through "
                f"while {to_c['category']} is outperforming at {to_c.get('avg_sell_through', 0):.0f}% sell-through."
            ),
            "root_cause": (
                f"Imbalanced stock allocation across categories. {from_c['category']} received excess depth "
                f"relative to demand; {to_c['category']} may face constraint on growth potential."
            ),
            "recommended_action": (
                f"Investigate cross-channel stock transfer opportunities between {from_c['category']} and "
                f"{to_c['category']}. Assess if product formats, sizes or colourways are transferable. "
                f"Review store-by-store stock allocation for rebalancing potential."
            ),
            "expected_outcome": (
                f"Better match of stock to consumer demand signals. "
                f"Estimated £{rev_impact:,} revenue improvement through optimised stock deployment."
            ),
            "business_impact": (
                f"Reduces markdown requirement in {from_c['category']} while supporting full-price "
                f"revenue growth in {to_c['category']}. Improves overall stock efficiency across the business."
            ),
            "revenue_impact": rev_impact,
            "margin_impact": round(avg_margin * 0.04, 1),
            "cash_recovery": cash_recovery,
            "confidence": confidence,
        })

    # ── 6. SUPPLIER REVIEW ───────────────────────────────────────
    # Trigger: category with margin < 38% AND sell-through < 40%
    reviewed_cats = set()
    for cat_data in categories:
        cat_mg  = cat_data.get("avg_margin_pct")
        cat_st  = cat_data.get("avg_sell_through")
        cat     = cat_data.get("category", "Unknown")
        cat_rev = cat_data.get("revenue", 0)

        if cat in reviewed_cats:
            continue
        if cat_mg is None or cat_mg > 38:
            continue
        if cat_st is None or cat_st > 42:
            continue

        reviewed_cats.add(cat)
        aid += 1
        target_margin    = 48
        margin_gap       = max(0, target_margin - cat_mg)
        potential_gain   = round(cat_rev * margin_gap / 100)
        confidence       = 63

        rev_w       = min(12, cat_rev / total_rev * 1200)
        mg_w        = margin_gap * 0.9
        priority_score = min(72, int(rev_w + mg_w + confidence * 0.22))

        actions.append({
            "id": aid, "type": "Supplier Review",
            "priority_label": "Medium", "priority_score": priority_score,
            "title": f"{cat} — Commercial Renegotiation",
            "sku": "Category Level", "category": cat,
            "issue": (
                f"{cat} delivering {cat_mg:.0f}% margin with {cat_st:.0f}% sell-through. "
                f"Both metrics are below commercial viability thresholds of 40% margin and 40% sell-through."
            ),
            "root_cause": (
                f"Cost prices in {cat} may not reflect current market conditions or range productivity. "
                f"Poor sell-through combined with low margin signals a fundamental commercial viability issue "
                f"that requires supplier-level review."
            ),
            "recommended_action": (
                f"Initiate cost price renegotiation with {cat} suppliers. "
                f"Target minimum {target_margin}% margin. "
                f"Review product range architecture and pricing strategy. "
                f"If renegotiation is unsuccessful, consider exiting the category."
            ),
            "expected_outcome": (
                f"Potential margin improvement of £{potential_gain:,} if cost prices reach the {target_margin}% threshold. "
                f"Long-term commercial viability of the {cat} category established."
            ),
            "business_impact": (
                f"Improves {cat} commercial contribution to the business. "
                f"Either establishes a viable margin floor or provides data to support a category exit decision. "
                f"Protects OTB investment in productive, margin-accretive categories."
            ),
            "revenue_impact": round(potential_gain * 0.35),
            "margin_impact": round(margin_gap, 1),
            "cash_recovery": 0,
            "confidence": confidence,
        })

    # ── Sort by priority score and rank ─────────────────────────
    actions.sort(key=lambda x: x["priority_score"], reverse=True)
    for i, a in enumerate(actions):
        a["rank"] = i + 1

    # ── Summary stats ────────────────────────────────────────────
    total_opportunity = sum(a["revenue_impact"] for a in actions)
    total_cash        = sum(a["cash_recovery"] for a in actions)
    avg_confidence    = round(sum(a["confidence"] for a in actions) / max(len(actions), 1))
    critical_count    = sum(1 for a in actions if a["priority_label"] == "Critical")
    high_count        = sum(1 for a in actions if a["priority_label"] == "High")

    return {
        "actions": actions,
        "summary": {
            "total_actions": len(actions),
            "critical_count": critical_count,
            "high_count": high_count,
            "medium_count": len(actions) - critical_count - high_count,
            "total_opportunity": round(total_opportunity),
            "total_cash_recovery": round(total_cash),
            "avg_confidence": avg_confidence,
        }
    }


@app.route("/action-centre")
@require_auth
def action_centre_page(user):
    return send_from_directory(".", "action-centre.html")


@app.route("/api/action-centre")
@require_auth
def get_action_centre(user):
    db = get_user_db()
    upload_id = request.args.get("upload_id") or session.get("current_upload_id")

    analysis = None
    if upload_id:
        try:
            res = db.table("uploads").select("analysis") \
                .eq("id", upload_id).eq("user_id", user["id"]).single().execute()
            if res.data:
                analysis = res.data["analysis"]
        except Exception:
            pass

    if not analysis:
        try:
            res = db.table("uploads").select("analysis") \
                .eq("user_id", user["id"]).order("created_at", desc=True).limit(1).execute()
            if res.data:
                analysis = res.data[0]["analysis"]
        except Exception:
            pass

    if not analysis:
        return jsonify({"actions": [], "summary": {}, "no_data": True})

    try:
        result = _run_action_engine(analysis)
        audit.log(AI_QUERY, request=request, user_id=user["id"],
                  company_id=user.get("company_id"),
                  resource="action-centre",
                  metadata={"action_count": len(result.get("actions", []))})
        return jsonify(result)
    except Exception as e:
        app.logger.error(f"Action centre error: {e}", exc_info=True)
        return jsonify({"error": "Failed to generate actions"}), 500


# ═══════════════════════════════════════════════════════════════
# AI TRADE REPORT ENGINE
# ═══════════════════════════════════════════════════════════════

def _build_report_prompt(analysis):
    """Build the Claude prompt for trade report generation."""
    analysis = prepare_for_ai(analysis)
    k    = analysis.get("kpis", {})
    cats = analysis.get("category_scorecard", [])
    ro   = analysis.get("reorder_alerts", [])
    md   = analysis.get("markdown_risk", [])
    bs   = analysis.get("best_sellers", [])
    wk   = analysis.get("weekly_trend", [])

    # Build rich data block for Claude
    rev   = k.get("total_revenue", 0)
    units = k.get("total_units", 0)
    gp    = k.get("gross_profit", 0)
    st    = k.get("avg_sell_through", 0)
    mg    = k.get("avg_margin_pct", 0)
    cover = k.get("avg_cover_weeks", 0)
    sku_c = analysis.get("meta", {}).get("sku_count", 0)

    # Derive weekly trend direction
    wk_vals = [w.get("revenue", 0) for w in wk[-6:]] if wk else []
    trend = "improving" if (len(wk_vals) >= 2 and wk_vals[-1] > wk_vals[0]) else \
            "declining" if (len(wk_vals) >= 2 and wk_vals[-1] < wk_vals[0]) else "stable"

    cat_summary = "\n".join(
        f"  - {c['category']}: £{c['revenue']:,.0f} rev, "
        f"{c.get('avg_sell_through','N/A')}% ST, "
        f"{c.get('avg_margin_pct','N/A')}% margin, "
        f"Health: {c.get('health','N/A')}"
        for c in cats[:8]
    ) if cats else "  No category data"

    reorder_summary = "\n".join(
        f"  - {r['sku']} ({r.get('category','')}) [{r['urgency']}]: "
        f"{r.get('cover_weeks','?')} wks cover, {r.get('sell_through','?')}% ST"
        for r in ro[:5]
    ) if ro else "  None"

    markdown_summary = "\n".join(
        f"  - {m['sku']} ({m.get('category','')}): {m.get('sell_through','?')}% ST, {m['severity']}"
        for m in md[:5]
    ) if md else "  None"

    top_sellers = "\n".join(
        f"  - {b['sku']}: £{b['revenue']:,.0f}, {b.get('sell_through','?')}% ST"
        for b in bs[:5]
    ) if bs else "  No data"

    data_block = f"""
PERIOD TRADING SUMMARY
======================
Revenue: £{rev:,.0f}
Units Sold: {units:,}
Gross Profit: £{gp:,.0f}
Gross Margin %: {mg:.1f}%
Sell-Through: {st:.1f}%
Avg Stock Cover: {cover:.1f} weeks
SKU Count: {sku_c}
Revenue Trend (recent): {trend}
Reorder Alerts: {k.get('reorder_count',0)} total ({k.get('critical_oos_count',0)} critical)
Markdown Risk SKUs: {k.get('markdown_risk_count',0)}

CATEGORY PERFORMANCE:
{cat_summary}

REORDER / STOCKOUT RISKS:
{reorder_summary}

MARKDOWN RISKS:
{markdown_summary}

TOP SELLING SKUs:
{top_sellers}
"""

    prompt = f"""You are a senior Merchandising Director at a major UK fashion retail group.
You are writing a professional trade report for the weekly trading meeting, senior leadership team, and retail board.

TRADING DATA:
{data_block}

WRITING RULES — FOLLOW EXACTLY:
1. NEVER simply list or repeat numbers — interpret them, explain the WHY, identify commercial implications
2. Use professional retail language: sell-through, OTB, intake, markdown cadence, forward cover, rate of sale, open-to-buy, stock turn, ATV, GP%, working capital
3. Be specific — reference actual category names and SKUs from the data above
4. Every paragraph should contain an insight, not just a description
5. Write as if you know this business deeply — confident, authoritative, commercially minded
6. Sound like a Head of Merchandising presenting to a Board
7. Recommended actions must be specific, prioritised, and actionable — not generic advice
8. Where the data shows a positive signal, name it confidently; where it shows risk, quantify the exposure

Return ONLY valid JSON (absolutely no markdown, no code fences, no backticks) with EXACTLY these keys:

{{
  "executive_summary": "2-3 sentences. Board-level overview. Must sound like: 'Overall trading remains [positive/mixed/under pressure] with revenue [tracking X% ahead of/behind plan]. [Best category] continues to [outperform/deliver] with [specific insight]. [Risk category] remains the primary area of focus due to [specific issue].'",
  "business_overview": "4-5 sentences. Overall trading picture. Cover: how is the business performing overall, what's driving the result, what's the primary opportunity and primary risk.",
  "revenue_performance": "3-4 sentences. Interpret revenue — not just the number. WHY is it at this level? Which categories or SKUs are contributing? What does the trend say? What is the revenue quality (full-price vs markdown)?",
  "margin_performance": "3-4 sentences. Interpret margin health. What is driving current margin? Where is margin at risk (markdown exposure, intake pricing)? What actions protect margin?",
  "stock_performance": "3-4 sentences. Interpret the stock position. Is cover healthy or excessive? What is the cash implication? Which areas of the estate have stock problems?",
  "category_winners": "3-4 sentences. Name specific winning categories with confidence. Explain what is driving their success. Identify how to support and build on that momentum.",
  "category_risks": "3-4 sentences. Name specific at-risk or underperforming categories. Explain the commercial risk. Quantify where possible. Identify the decision needed.",
  "markdown_risks": "3-4 sentences. Specific markdown exposure. Which SKUs or categories are at greatest risk? What is the recommended markdown strategy? What is the cost of inaction?",
  "replenishment_opportunities": "3-4 sentences. Identify growth opportunities. Where is demand outpacing supply? What is the opportunity cost of inaction? What should be prioritised for replenishment?",
  "recommended_actions": [
    "1. [Most urgent action — specific, with clear business rationale]",
    "2. [Second priority action]",
    "3. [Third priority action]",
    "4. [Fourth priority action]",
    "5. [Fifth priority action]",
    "6. [Sixth priority action — strategic/longer-term]"
  ],
  "weekly_trade_narrative": "5-6 sentences. Weekly trade pack narrative. More operational tone. Cover rate of sale vs target, stock movements, immediate priorities for the week, and the key trading call for the week.",
  "risk_register_summary": "4-5 sentences. Risk report tone. Cover: what are the 2-3 primary risks to trading performance, what is the financial exposure, and what mitigation is recommended.",
  "category_commentary": {{
    "summary": "3-4 sentences overview of category portfolio performance",
    "winners_detail": "4-5 sentences detailed analysis of top-performing categories with specific commercial insight",
    "risks_detail": "4-5 sentences detailed analysis of underperforming categories with specific commercial implication"
  }}
}}"""

    return prompt


def _generate_report_sections(analysis):
    """Call Claude API to generate the full trade report. Returns dict of report sections."""
    prompt = _build_report_prompt(analysis)
    try:
        response = ai_client.messages.create(
            model=MODEL,
            max_tokens=3000,
            messages=[{"role": "user", "content": prompt}]
        )
        raw = response.content[0].text.strip()
        # Strip any accidental markdown fences
        if raw.startswith("```"):
            raw = "\n".join(raw.split("\n")[1:])
        if raw.endswith("```"):
            raw = "\n".join(raw.split("\n")[:-1])
        raw = raw.strip()
        return json.loads(raw)
    except json.JSONDecodeError:
        # If JSON parse fails, return a fallback with the raw text
        return _fallback_report(analysis)
    except Exception as e:
        app.logger.error(f"Report generation error: {e}", exc_info=True)
        return _fallback_report(analysis)


def _fallback_report(analysis):
    """Generate a data-driven fallback report without AI (rule-based)."""
    k    = analysis.get("kpis", {})
    cats = analysis.get("category_scorecard", [])
    ro   = analysis.get("reorder_alerts", [])
    md   = analysis.get("markdown_risk", [])
    bs   = analysis.get("best_sellers", [])

    rev    = k.get("total_revenue", 0)
    units  = k.get("total_units", 0)
    st     = k.get("avg_sell_through", 0)
    mg     = k.get("avg_margin_pct", 0)
    cover  = k.get("avg_cover_weeks", 0)
    rc     = k.get("reorder_count", 0)
    mdc    = k.get("markdown_risk_count", 0)

    top_cat   = cats[0]["category"] if cats else "the top category"
    risk_cat  = next((c["category"] for c in cats if c.get("health") in ("AT RISK","CRITICAL")), None)
    top_sku   = bs[0]["sku"] if bs else None
    crit_ro   = next((r["sku"] for r in ro if r.get("urgency") == "CRITICAL"), None)
    top_md    = md[0]["sku"] if md else None

    health = "positive" if st >= 45 and mg >= 45 else "under pressure" if st < 25 else "mixed"

    return {
        "executive_summary": (
            f"Overall trading is {health} with £{rev:,.0f} in revenue and {st:.0f}% sell-through recorded across "
            f"{analysis.get('meta',{}).get('sku_count',0)} active SKUs. "
            f"{top_cat} is the primary revenue driver, whilst "
            f"{f'{risk_cat} presents the most significant stock risk.' if risk_cat else 'stock cover requires monitoring.'} "
            f"Recommended actions focus on {'protecting full-price sell-through' if st >= 40 else 'accelerating clearance and improving rate of sale'}."
        ),
        "business_overview": (
            f"The business is generating £{rev:,.0f} in revenue from {units:,} units across "
            f"{analysis.get('meta',{}).get('sku_count',0)} SKUs, delivering {mg:.1f}% gross margin. "
            f"Sell-through at {st:.0f}% is {'ahead of expectations' if st >= 45 else 'below target'}, "
            f"with an average of {cover:.1f} weeks cover remaining across the estate. "
            f"There are {rc} SKUs requiring replenishment action and {mdc} at risk of requiring markdown support."
        ),
        "revenue_performance": (
            f"Revenue of £{rev:,.0f} is being driven by {units:,} units at an average selling price of "
            f"£{rev/max(units,1):.2f}. "
            f"{f'Top performer {top_sku} is contributing significantly to overall revenue delivery. ' if top_sku else ''}"
            f"Revenue quality {'appears healthy with strong full-price activity' if st >= 45 else 'is under pressure with elevated markdown dependency'}."
        ),
        "margin_performance": (
            f"Gross margin is running at {mg:.1f}%, delivering £{k.get('gross_profit',0):,.0f} gross profit. "
            f"{'Margin is above the 45% threshold, indicating disciplined full-price selling.' if mg >= 45 else 'Margin is below the 45% target — intake pricing or markdown activity is compressing returns.'} "
            f"With {mdc} SKUs carrying markdown risk, margin protection requires active management."
        ),
        "stock_performance": (
            f"Stock cover averaging {cover:.1f} weeks {'is within the target 4–8 week range.' if 4 <= cover <= 8 else 'is outside target, requiring rebalancing action.'} "
            f"There are {rc} SKUs with critically low cover at risk of stockout. "
            f"{'Markdown risk across ' + str(mdc) + ' SKUs represents potential working capital trapped in slow-moving lines.' if mdc > 0 else 'Markdown exposure is currently minimal.'}"
        ),
        "category_winners": (
            f"{top_cat} is the standout category, "
            f"delivering the strongest revenue contribution to the overall business. "
            f"Categories with {'strong health signals' if any(c.get('health') == 'STRONG' for c in cats) else 'healthy sell-through'} "
            f"should be prioritised for replenishment to sustain momentum."
        ),
        "category_risks": (
            f"{f'{risk_cat} is the primary area of concern' if risk_cat else 'At-risk categories require immediate attention'}, "
            f"with elevated stock cover and below-target sell-through. "
            f"If left unmanaged, this stock position will require deeper markdown to clear, "
            f"further compressing margin and tying up working capital."
        ),
        "markdown_risks": (
            f"There are {mdc} SKUs carrying markdown risk across the range. "
            f"{f'{top_md} is the highest-priority markdown candidate.' if top_md else 'Markdown candidates span multiple categories.'} "
            f"A phased markdown strategy is recommended to maximise recovery whilst protecting margin."
        ),
        "replenishment_opportunities": (
            f"With {rc} SKUs at or approaching stockout, replenishment is an immediate priority. "
            f"{f'{crit_ro} is the most critical.' if crit_ro else 'Critical SKUs should be fast-tracked for reorder.'} "
            f"Failure to replenish will result in lost revenue and weakened category performance."
        ),
        "recommended_actions": [
            f"1. Place emergency replenishment orders on {rc} SKUs showing critical/urgent stockout risk to prevent lost revenue.",
            f"2. Implement a phased markdown strategy on the {mdc} highest-risk SKUs, starting at 20% to test rate-of-sale response.",
            f"3. Review forward intake commitments for {f'{risk_cat}' if risk_cat else 'at-risk categories'} and cancel or defer non-critical orders.",
            f"4. Reforecast revenue and margin for the remainder of the season based on current sell-through trajectory.",
            f"5. Prioritise full-price promotional support for strong-performing categories to sustain momentum.",
            f"6. Conduct a supplier review for categories delivering below 40% margin to identify cost-price renegotiation opportunities.",
        ],
        "weekly_trade_narrative": (
            f"This week's trading performance reflects a {health} business position with sell-through at {st:.0f}% and "
            f"revenue of £{rev:,.0f}. The immediate priority for the trading team is to address {rc} stockout risks "
            f"and initiate markdown action on slow-moving lines. "
            f"Category mix is {'favouring stronger performers' if st >= 40 else 'showing signs of strain across multiple departments'}. "
            f"The key trading call this week is to protect margin on full-price lines whilst accelerating clearance on aged stock."
        ),
        "risk_register_summary": (
            f"The primary trading risk is concentrated in {'the ' + risk_cat + ' category' if risk_cat else 'slow-moving categories'}, "
            f"where elevated cover and declining sell-through create markdown exposure. "
            f"Stockout risk on {rc} SKUs presents a secondary revenue risk if replenishment is not actioned immediately. "
            f"Margin compression from markdown dependency on {mdc} SKUs could impact gross profit delivery for the period. "
            f"Mitigation requires immediate action on both replenishment and clearance markdown."
        ),
        "category_commentary": {
            "summary": f"The category portfolio is showing mixed performance with {top_cat} leading and {risk_cat or 'several categories'} under pressure.",
            "winners_detail": f"{top_cat} and strong-performing categories are delivering above-target sell-through and supporting overall revenue delivery.",
            "risks_detail": f"{risk_cat or 'Underperforming categories'} require urgent commercial intervention to prevent margin erosion and excess stock accumulation."
        }
    }


def _build_pptx(report, analysis):
    """Generate a PowerPoint trade report using python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import io as _io

    # Colours
    DARK   = RGBColor(0x1E, 0x1E, 0x2A)
    MINT   = RGBColor(0x7D, 0xAA, 0x75)
    MINT_L = RGBColor(0xA8, 0xC5, 0xA0)
    BLUSH  = RGBColor(0xF4, 0xC0, 0xA8)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY   = RGBColor(0x5C, 0x5C, 0x72)
    LGRAY  = RGBColor(0xF3, 0xF5, 0xF8)
    RED    = RGBColor(0xE3, 0x52, 0x3A)

    k = analysis.get("kpis", {})
    from datetime import date as _date
    today = _date.today().strftime("%d %B %Y")

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Blank

    def add_slide():
        return prs.slides.add_slide(blank_layout)

    def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None):
        from pptx.util import Inches as I
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.line.fill.background() if not line_rgb else None
        if line_rgb:
            shape.line.color.rgb = line_rgb
        if fill_rgb:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_rgb
        else:
            shape.fill.background()
        return shape

    def add_text(slide, text, left, top, width, height,
                 font_size=14, bold=False, color=None, align=PP_ALIGN.LEFT,
                 font_name="Calibri", italic=False, wrap=True):
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height))
        txBox.word_wrap = wrap
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font_name
        if color:
            run.font.color.rgb = color
        return txBox

    def add_multiline(slide, paragraphs, left, top, width, height,
                      font_size=11, color=None, bold=False):
        """Add a textbox with multiple paragraphs."""
        from pptx.util import Pt as P
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height))
        txBox.word_wrap = True
        tf = txBox.text_frame
        tf.word_wrap = True
        first = True
        for para_text in paragraphs:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.space_after = Pt(4)
            run = p.add_run()
            run.text = para_text
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.name = "Calibri"
            if color:
                run.font.color.rgb = color
        return txBox

    def slide_header(slide, title, subtitle=None):
        """Add standard header bar."""
        add_rect(slide, 0, 0, 13.33, 1.1, fill_rgb=DARK)
        add_text(slide, "SKUVVY", 0.35, 0.12, 3, 0.4,
                 font_size=10, bold=True, color=MINT_L, font_name="Calibri")
        add_text(slide, title, 0.35, 0.45, 10, 0.55,
                 font_size=22, bold=True, color=WHITE, font_name="Calibri")
        if subtitle:
            add_text(slide, subtitle, 0.35, 0.85, 8, 0.3,
                     font_size=10, color=RGBColor(0xAA,0xBB,0xCC), font_name="Calibri")
        # Date top right
        add_text(slide, today, 10.5, 0.12, 2.5, 0.4,
                 font_size=10, color=RGBColor(0x88,0x99,0xAA), align=PP_ALIGN.RIGHT)

    def kpi_box(slide, left, top, label, value, color=MINT):
        """Add a small KPI box."""
        add_rect(slide, left, top, 2.8, 1.1, fill_rgb=LGRAY)
        add_text(slide, label, left+0.12, top+0.08, 2.56, 0.35,
                 font_size=9, color=GRAY, bold=False)
        add_text(slide, value, left+0.12, top+0.38, 2.56, 0.6,
                 font_size=20, bold=True, color=color)

    # ── SLIDE 1: Title ─────────────────────────────────────────
    s = add_slide()
    add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=DARK)
    add_rect(s, 0, 5.8, 13.33, 1.7, fill_rgb=RGBColor(0x12,0x16,0x20))
    # Accent bar
    add_rect(s, 0, 2.9, 0.06, 1.8, fill_rgb=MINT)
    add_text(s, "SKUVVY", 0.35, 1.2, 12, 0.6,
             font_size=13, bold=True, color=MINT_L, font_name="Calibri")
    add_text(s, "AI Trade Report", 0.35, 2.0, 12, 1.2,
             font_size=46, bold=True, color=WHITE, font_name="Calibri")
    add_text(s, "Merchandising Intelligence | " + today, 0.35, 3.3, 12, 0.5,
             font_size=14, color=RGBColor(0x88,0x99,0xAA), font_name="Calibri")
    add_text(s, "CONFIDENTIAL — FOR INTERNAL USE ONLY", 0.35, 6.1, 12, 0.5,
             font_size=10, color=RGBColor(0x55,0x66,0x77), font_name="Calibri")

    # ── SLIDE 2: Executive Summary ─────────────────────────────
    s = add_slide()
    slide_header(s, "Executive Summary", "Board Overview")
    # Big summary text
    add_rect(s, 0.35, 1.25, 12.63, 2.2, fill_rgb=RGBColor(0xF8,0xFB,0xF8))
    add_text(s, report.get("executive_summary",""), 0.55, 1.35, 12.2, 2.0,
             font_size=14, color=DARK, font_name="Calibri", italic=True)
    # KPIs
    kpi_box(s, 0.35, 3.7,  "Revenue",       f"£{k.get('total_revenue',0)/1000:.0f}k", MINT)
    kpi_box(s, 3.25, 3.7, "Gross Margin",   f"{k.get('avg_margin_pct',0):.1f}%",      MINT)
    kpi_box(s, 6.15, 3.7, "Sell-Through",   f"{k.get('avg_sell_through',0):.0f}%",    MINT)
    kpi_box(s, 9.05, 3.7, "Stock Cover",    f"{k.get('avg_cover_weeks',0):.1f} wks",  MINT)
    # Business overview
    add_text(s, report.get("business_overview",""), 0.35, 5.1, 12.63, 2.0,
             font_size=11, color=GRAY, font_name="Calibri")

    # ── SLIDE 3: Revenue & Margin Performance ──────────────────
    s = add_slide()
    slide_header(s, "Financial Performance", "Revenue & Margin Analysis")
    # Two columns
    add_rect(s, 0.35, 1.25, 6.1, 0.38, fill_rgb=MINT)
    add_text(s, "REVENUE PERFORMANCE", 0.45, 1.27, 5.9, 0.34,
             font_size=10, bold=True, color=WHITE, font_name="Calibri")
    add_text(s, report.get("revenue_performance",""), 0.35, 1.72, 6.1, 3.0,
             font_size=11, color=DARK, font_name="Calibri")

    add_rect(s, 6.88, 1.25, 6.1, 0.38, fill_rgb=BLUSH)
    add_text(s, "MARGIN PERFORMANCE", 6.98, 1.27, 5.9, 0.34,
             font_size=10, bold=True, color=DARK, font_name="Calibri")
    add_text(s, report.get("margin_performance",""), 6.88, 1.72, 6.1, 3.0,
             font_size=11, color=DARK, font_name="Calibri")

    # KPI strip
    add_text(s, f"£{k.get('total_revenue',0)/1000:.0f}k Revenue", 0.35, 4.9, 3, 0.5,
             font_size=16, bold=True, color=MINT, font_name="Calibri")
    add_text(s, f"£{k.get('gross_profit',0)/1000:.0f}k Gross Profit", 3.5, 4.9, 3.5, 0.5,
             font_size=16, bold=True, color=MINT, font_name="Calibri")
    add_text(s, f"{k.get('avg_margin_pct',0):.1f}% GM%", 7.5, 4.9, 3, 0.5,
             font_size=16, bold=True, color=RGBColor(0xE3,0xA0,0x8E), font_name="Calibri")
    add_text(s, f"{k.get('total_units',0):,} Units", 10.5, 4.9, 2.5, 0.5,
             font_size=16, bold=True, color=GRAY, font_name="Calibri")

    # Stock performance
    add_rect(s, 0.35, 5.55, 12.63, 0.38, fill_rgb=RGBColor(0xB8,0xC4,0xE8))
    add_text(s, "STOCK PERFORMANCE", 0.45, 5.57, 12, 0.34,
             font_size=10, bold=True, color=DARK, font_name="Calibri")
    add_text(s, report.get("stock_performance",""), 0.35, 6.0, 12.63, 1.2,
             font_size=11, color=DARK, font_name="Calibri")

    # ── SLIDE 4: Category Performance ─────────────────────────
    s = add_slide()
    slide_header(s, "Category Performance", "Winners & Risks")
    cats_data = analysis.get("category_scorecard", [])

    add_rect(s, 0.35, 1.25, 5.9, 0.38, fill_rgb=MINT)
    add_text(s, "CATEGORY WINNERS", 0.45, 1.27, 5.7, 0.34,
             font_size=10, bold=True, color=WHITE)
    add_text(s, report.get("category_winners",""), 0.35, 1.72, 5.9, 1.6,
             font_size=11, color=DARK, font_name="Calibri")

    add_rect(s, 7.08, 1.25, 5.9, 0.38, fill_rgb=RED)
    add_text(s, "CATEGORY RISKS", 7.18, 1.27, 5.7, 0.34,
             font_size=10, bold=True, color=WHITE)
    add_text(s, report.get("category_risks",""), 7.08, 1.72, 5.9, 1.6,
             font_size=11, color=DARK, font_name="Calibri")

    # Category table
    if cats_data:
        headers = ["Category", "Revenue", "Sell-Through", "Margin", "Health"]
        col_w = [2.8, 1.9, 1.9, 1.6, 1.4]
        col_x = [0.35, 3.15, 5.05, 6.95, 8.55]
        row_y = 3.55
        # Header row
        for i, (h, cx, cw) in enumerate(zip(headers, col_x, col_w)):
            add_rect(s, cx, row_y, cw-0.05, 0.34, fill_rgb=DARK)
            add_text(s, h, cx+0.08, row_y+0.05, cw-0.15, 0.28,
                     font_size=9, bold=True, color=WHITE)
        for ri, cat in enumerate(cats_data[:6]):
            ry = row_y + 0.38 + ri * 0.38
            bg = RGBColor(0xF8,0xFB,0xF8) if ri % 2 == 0 else WHITE
            row_vals = [
                cat.get("category",""),
                f"£{cat.get('revenue',0)/1000:.0f}k",
                f"{cat.get('avg_sell_through','—')}%",
                f"{cat.get('avg_margin_pct','—')}%",
                cat.get("health","—")
            ]
            hcol = {"STRONG":MINT, "HEALTHY":MINT_L, "AT RISK":RGBColor(0xF0,0xC4,0xA0),
                    "CRITICAL":RED}.get(cat.get("health",""), GRAY)
            for ci, (val, cx, cw) in enumerate(zip(row_vals, col_x, col_w)):
                add_rect(s, cx, ry, cw-0.05, 0.34, fill_rgb=bg)
                add_text(s, val, cx+0.08, ry+0.04, cw-0.15, 0.3,
                         font_size=9.5, color=(hcol if ci == 4 else DARK),
                         bold=(ci == 4))

    # ── SLIDE 5: Markdown & Replenishment ─────────────────────
    s = add_slide()
    slide_header(s, "Commercial Risks & Opportunities", "Markdown & Replenishment")

    add_rect(s, 0.35, 1.25, 5.9, 0.38, fill_rgb=RED)
    add_text(s, "MARKDOWN RISKS", 0.45, 1.27, 5.7, 0.34,
             font_size=10, bold=True, color=WHITE)
    add_text(s, report.get("markdown_risks",""), 0.35, 1.72, 5.9, 1.8,
             font_size=11, color=DARK, font_name="Calibri")

    add_rect(s, 7.08, 1.25, 5.9, 0.38, fill_rgb=MINT)
    add_text(s, "REPLENISHMENT OPPORTUNITIES", 7.18, 1.27, 5.7, 0.34,
             font_size=10, bold=True, color=WHITE)
    add_text(s, report.get("replenishment_opportunities",""), 7.08, 1.72, 5.9, 1.8,
             font_size=11, color=DARK, font_name="Calibri")

    # Markdown risk SKUs
    md_skus = analysis.get("markdown_risk", [])[:5]
    if md_skus:
        add_rect(s, 0.35, 3.7, 5.9, 0.32, fill_rgb=RGBColor(0x25,0x10,0x10))
        add_text(s, "TOP MARKDOWN SKUs", 0.45, 3.72, 5.7, 0.28,
                 font_size=9, bold=True, color=RED)
        for ri, m in enumerate(md_skus):
            ry = 4.08 + ri * 0.34
            add_rect(s, 0.35, ry, 5.9, 0.3, fill_rgb=RGBColor(0xFC,0xF5,0xF5))
            add_text(s,
                     f"{m['sku']}  |  ST: {m.get('sell_through','?')}%  |  {m.get('severity','—')}  |  Depth: {m.get('recommended_depth','?')}",
                     0.5, ry+0.04, 5.6, 0.26, font_size=9.5, color=DARK)

    # Reorder SKUs
    ro_skus = analysis.get("reorder_alerts", [])[:5]
    if ro_skus:
        add_rect(s, 7.08, 3.7, 5.9, 0.32, fill_rgb=RGBColor(0x0E,0x22,0x14))
        add_text(s, "TOP REORDER SKUs", 7.18, 3.72, 5.7, 0.28,
                 font_size=9, bold=True, color=MINT)
        for ri, r in enumerate(ro_skus):
            ry = 4.08 + ri * 0.34
            add_rect(s, 7.08, ry, 5.9, 0.3, fill_rgb=RGBColor(0xF5,0xFB,0xF5))
            add_text(s,
                     f"{r['sku']}  |  Cover: {r.get('cover_weeks','?')} wks  |  {r.get('urgency','—')}",
                     7.23, ry+0.04, 5.6, 0.26, font_size=9.5, color=DARK)

    # ── SLIDE 6: Recommended Actions ──────────────────────────
    s = add_slide()
    slide_header(s, "Recommended Actions", "Priority Order")
    actions = report.get("recommended_actions", [])
    colors_act = [RED, RED, RGBColor(0xD4,0x86,0x2A), RGBColor(0xD4,0x86,0x2A), MINT, MINT]
    labels_act = ["CRITICAL","CRITICAL","HIGH","HIGH","RECOMMENDED","RECOMMENDED"]
    for ai, action in enumerate(actions[:6]):
        ay = 1.35 + ai * 0.95
        ax_l = 0.35 if ai % 2 == 0 else 6.84
        ay_r = 1.35 + (ai // 2) * 0.95
        col = colors_act[ai]
        lbl = labels_act[ai]
        lx = 0.35 if ai < 3 else 6.84
        ly = 1.35 + (ai % 3) * 0.95
        # Draw action card
        add_rect(s, lx, ly, 6.14, 0.82, fill_rgb=LGRAY)
        add_rect(s, lx, ly, 0.08, 0.82, fill_rgb=col)
        add_text(s, lbl, lx+0.18, ly+0.06, 1.8, 0.24, font_size=8, bold=True, color=col)
        add_text(s, action, lx+0.18, ly+0.28, 5.8, 0.52, font_size=10, color=DARK, font_name="Calibri")

    # ── SLIDE 7: Risk Register ─────────────────────────────────
    s = add_slide()
    slide_header(s, "Risk Register", "Period Risk Summary")
    add_rect(s, 0.35, 1.25, 12.63, 2.4, fill_rgb=RGBColor(0x22,0x10,0x10))
    add_text(s, report.get("risk_register_summary",""), 0.55, 1.35, 12.2, 2.2,
             font_size=12, color=RGBColor(0xFF,0xDD,0xD8), font_name="Calibri")
    # Risk columns
    risks = [
        ("STOCK RISK", f"{k.get('markdown_risk_count',0)} SKUs markdown risk", RED),
        ("STOCKOUT RISK", f"{k.get('reorder_count',0)} SKUs need replenishment", RGBColor(0xD4,0x86,0x2A)),
        ("MARGIN RISK", f"{k.get('avg_margin_pct',0):.1f}% current margin delivery", RGBColor(0xB8,0xC4,0xE8)),
    ]
    for ri, (title, val, col) in enumerate(risks):
        rx = 0.35 + ri * 4.33
        add_rect(s, rx, 3.85, 4.18, 1.2, fill_rgb=LGRAY)
        add_rect(s, rx, 3.85, 4.18, 0.32, fill_rgb=col)
        add_text(s, title, rx+0.1, 3.88, 3.9, 0.28, font_size=9, bold=True, color=WHITE)
        add_text(s, val, rx+0.1, 4.24, 3.9, 0.65, font_size=13, bold=True, color=DARK)
    # Mitigation
    add_rect(s, 0.35, 5.2, 12.63, 0.32, fill_rgb=MINT)
    add_text(s, "MITIGATION PRIORITIES", 0.45, 5.22, 12, 0.28, font_size=9, bold=True, color=WHITE)
    add_text(s, "  ·  ".join((actions[:3]) if actions else ["No actions defined"]),
             0.35, 5.6, 12.63, 1.6, font_size=9.5, color=DARK, font_name="Calibri")

    # ── SLIDE 8: End Slide ─────────────────────────────────────
    s = add_slide()
    add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=DARK)
    add_rect(s, 0, 3.4, 13.33, 0.06, fill_rgb=MINT)
    add_text(s, "SKUVVY", 0, 3.6, 13.33, 1.2,
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "AI Merchandising Intelligence", 0, 4.8, 13.33, 0.6,
             font_size=14, color=MINT_L, align=PP_ALIGN.CENTER)
    add_text(s, f"Report generated {today} · aeroa-ai.up.railway.app",
             0, 6.5, 13.33, 0.5, font_size=10,
             color=RGBColor(0x55,0x66,0x77), align=PP_ALIGN.CENTER)

    buf = _io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _build_xlsx(report, analysis):
    """Generate a multi-sheet Excel workbook of the analysed data using openpyxl."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter
    from datetime import date
    import io as _io

    k     = analysis.get("kpis", {}) or {}
    cats  = analysis.get("category_scorecard", []) or []
    best  = analysis.get("best_sellers", []) or []
    reord = analysis.get("reorder_alerts", []) or []
    mark  = analysis.get("markdown_risk", []) or []

    HEAD_FILL = PatternFill("solid", fgColor="1F2A44")
    HEAD_FONT = Font(bold=True, color="FFFFFF", size=11)
    TITLE_FONT = Font(bold=True, size=14, color="1F2A44")
    thin = Side(style="thin", color="D9DEE7")
    BORDER = Border(left=thin, right=thin, top=thin, bottom=thin)

    wb = Workbook()

    def style_header(ws, row, ncols):
        for c in range(1, ncols + 1):
            cell = ws.cell(row=row, column=c)
            cell.fill = HEAD_FILL
            cell.font = HEAD_FONT
            cell.alignment = Alignment(horizontal="left", vertical="center")
            cell.border = BORDER

    def autosize(ws, widths):
        for i, w in enumerate(widths, start=1):
            ws.column_dimensions[get_column_letter(i)].width = w

    def write_table(ws, start_row, headers, rows, keys, fmts=None):
        for j, h in enumerate(headers, start=1):
            ws.cell(row=start_row, column=j, value=h)
        style_header(ws, start_row, len(headers))
        r = start_row + 1
        for item in rows:
            for j, key in enumerate(keys, start=1):
                val = item.get(key)
                if fmts and key in fmts and isinstance(val, (int, float)):
                    val = fmts[key](val)
                ws.cell(row=r, column=j, value=val).border = BORDER
            r += 1
        return r

    # ── Sheet 1: Summary / KPIs ──
    ws = wb.active
    ws.title = "Summary"
    ws["A1"] = "Skuvvy — Trading Analysis Summary"
    ws["A1"].font = TITLE_FONT
    ws["A2"] = f"Generated {date.today().strftime('%d %B %Y')}"
    ws["A2"].font = Font(italic=True, color="55606E")

    kpi_rows = [
        ("Total Revenue",        f"£{k.get('total_revenue',0):,.0f}"),
        ("Total Units",          f"{k.get('total_units',0):,.0f}"),
        ("Gross Profit",         f"£{k.get('gross_profit',0):,.0f}"),
        ("Average Sell-Through", f"{k.get('avg_sell_through',0):.1f}%"),
        ("Average Margin",       f"{k.get('avg_margin_pct',0):.1f}%"
                                 if k.get('margin_available', True) else "Insufficient data"),
        ("Blended Gross Margin", f"{k.get('gross_margin_pct',0):.1f}%"
                                 if k.get('margin_available', True) else "Insufficient data"),
        ("Average Stock Cover",  f"{k.get('avg_cover_weeks',0):.1f} weeks"),
        ("Reorder Alerts",       k.get('reorder_count', 0)),
        ("Markdown Risk SKUs",   k.get('markdown_risk_count', 0)),
        ("Critical Stockouts",   k.get('critical_oos_count', 0)),
    ]
    ws.cell(row=4, column=1, value="Metric")
    ws.cell(row=4, column=2, value="Value")
    style_header(ws, 4, 2)
    for i, (name, val) in enumerate(kpi_rows, start=5):
        ws.cell(row=i, column=1, value=name).border = BORDER
        ws.cell(row=i, column=2, value=val).border = BORDER
    autosize(ws, [26, 22])

    # ── Sheet 2: Categories ──
    if cats:
        ws = wb.create_sheet("Categories")
        write_table(ws, 1,
                    ["Category", "Revenue (£)", "Units", "SKUs", "Sell-Through %", "Margin %", "Health"],
                    cats,
                    ["category", "revenue", "units", "sku_count", "avg_sell_through", "avg_margin_pct", "health"])
        autosize(ws, [22, 16, 12, 8, 16, 12, 14])

    # ── Sheet 3: Best Sellers ──
    if best:
        ws = wb.create_sheet("Best Sellers")
        write_table(ws, 1,
                    ["SKU", "Name", "Category", "Revenue (£)", "Units", "Sell-Through %", "Margin %", "Stock"],
                    best,
                    ["sku", "name", "category", "revenue", "units", "sell_through", "margin_pct", "stock"])
        autosize(ws, [14, 28, 18, 16, 10, 16, 12, 10])

    # ── Sheet 4: Reorder Alerts ──
    if reord:
        ws = wb.create_sheet("Reorder Alerts")
        write_table(ws, 1,
                    ["SKU", "Name", "Category", "Urgency", "Revenue (£)", "Stock", "Cover (wks)", "Sell-Through %"],
                    reord,
                    ["sku", "name", "category", "urgency", "revenue", "stock", "cover_weeks", "sell_through"])
        autosize(ws, [14, 28, 18, 12, 16, 10, 12, 16])

    # ── Sheet 5: Markdown Risk ──
    if mark:
        ws = wb.create_sheet("Markdown Risk")
        write_table(ws, 1,
                    ["SKU", "Name", "Category", "Severity", "Sell-Through %", "Recommended Depth", "Revenue (£)"],
                    mark,
                    ["sku", "name", "category", "severity", "sell_through", "recommended_depth", "revenue"])
        autosize(ws, [14, 28, 18, 12, 16, 20, 16])

    buf = _io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf


def _build_pdf(report, analysis):
    """Generate a clean, board-ready PDF executive report using reportlab."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import mm
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_LEFT
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, ListFlowable, ListItem
    )
    from datetime import date
    import io as _io
    import re

    k    = analysis.get("kpis", {}) or {}
    cats = analysis.get("category_scorecard", []) or []
    margin_ok = k.get("margin_available", True)

    NAVY  = colors.HexColor("#1F2A44")
    GREEN = colors.HexColor("#5B6BB5")
    GREY  = colors.HexColor("#55606E")
    LIGHT = colors.HexColor("#F3F5F9")

    styles = getSampleStyleSheet()
    h_title = ParagraphStyle("hTitle", parent=styles["Title"], textColor=NAVY,
                             fontSize=22, spaceAfter=2, alignment=TA_LEFT)
    h_sub   = ParagraphStyle("hSub", parent=styles["Normal"], textColor=GREY,
                             fontSize=10, spaceAfter=12)
    h_sec   = ParagraphStyle("hSec", parent=styles["Heading2"], textColor=NAVY,
                             fontSize=13, spaceBefore=12, spaceAfter=5)
    body    = ParagraphStyle("body", parent=styles["Normal"], fontSize=10,
                             leading=15, textColor=colors.HexColor("#243040"), spaceAfter=6)

    def esc(t):
        return (str(t or "")).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    story = []
    story.append(Paragraph("Skuvvy — Executive Trade Report", h_title))
    story.append(Paragraph(f"Generated {date.today().strftime('%d %B %Y')}", h_sub))

    # KPI band
    margin_txt = f"{k.get('avg_margin_pct',0):.1f}%" if margin_ok else "n/a"
    kpi_data = [
        ["Revenue", "Gross Margin", "Sell-Through", "Stock Cover"],
        [f"£{k.get('total_revenue',0):,.0f}", margin_txt,
         f"{k.get('avg_sell_through',0):.0f}%", f"{k.get('avg_cover_weeks',0):.1f} wks"],
    ]
    kpi_tbl = Table(kpi_data, colWidths=[42*mm]*4)
    kpi_tbl.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), NAVY),
        ("TEXTCOLOR",  (0, 0), (-1, 0), colors.white),
        ("FONTNAME",   (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE",   (0, 0), (-1, 0), 9),
        ("BACKGROUND", (0, 1), (-1, 1), LIGHT),
        ("TEXTCOLOR",  (0, 1), (-1, 1), NAVY),
        ("FONTNAME",   (0, 1), (-1, 1), "Helvetica-Bold"),
        ("FONTSIZE",   (0, 1), (-1, 1), 13),
        ("ALIGN",      (0, 0), (-1, -1), "CENTER"),
        ("VALIGN",     (0, 0), (-1, -1), "MIDDLE"),
        ("TOPPADDING", (0, 0), (-1, -1), 7),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 7),
        ("GRID",       (0, 0), (-1, -1), 0.5, colors.white),
    ]))
    story.append(kpi_tbl)
    story.append(Spacer(1, 10))

    # Narrative sections (only those present)
    sections = [
        ("Executive Summary",       report.get("executive_summary")),
        ("Business Overview",       report.get("business_overview")),
        ("Revenue Performance",     report.get("revenue_performance")),
        ("Margin Performance",      report.get("margin_performance")),
        ("Stock Performance",       report.get("stock_performance")),
        ("Markdown Risk",           report.get("markdown_risks")),
        ("Replenishment Priorities", report.get("replenishment_opportunities")),
        ("Risk Register",           report.get("risk_register_summary")),
    ]
    for title, text in sections:
        if text:
            story.append(Paragraph(title, h_sec))
            story.append(Paragraph(esc(text), body))

    # Recommended actions
    actions = report.get("recommended_actions") or []
    if actions:
        story.append(Paragraph("Recommended Actions", h_sec))
        items = [ListItem(Paragraph(esc(re.sub(r'^\d+\.\s*', '', a)), body), leftIndent=6)
                 for a in actions]
        story.append(ListFlowable(items, bulletType="1", bulletColor=GREEN,
                                  bulletFontSize=10, leftIndent=14))

    # Category table
    if cats:
        story.append(Paragraph("Category Scorecard", h_sec))
        rows = [["Category", "Revenue", "Sell-Through", "Margin", "Health"]]
        for c in cats[:12]:
            mg = c.get("avg_margin_pct")
            rows.append([
                esc(c.get("category", "—")),
                f"£{float(c.get('revenue') or 0):,.0f}",
                f"{float(c.get('avg_sell_through') or 0):.0f}%",
                (f"{float(mg):.0f}%" if mg is not None else "—"),
                esc(c.get("health", "—")),
            ])
        ctbl = Table(rows, colWidths=[45*mm, 32*mm, 30*mm, 24*mm, 28*mm])
        ctbl.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), NAVY),
            ("TEXTCOLOR",  (0, 0), (-1, 0), colors.white),
            ("FONTNAME",   (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE",   (0, 0), (-1, -1), 9),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, LIGHT]),
            ("TEXTCOLOR",  (0, 1), (-1, -1), NAVY),
            ("ALIGN",      (1, 0), (-1, -1), "CENTER"),
            ("VALIGN",     (0, 0), (-1, -1), "MIDDLE"),
            ("TOPPADDING", (0, 0), (-1, -1), 5),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ("LINEBELOW",  (0, 0), (-1, 0), 0.5, NAVY),
            ("GRID",       (0, 0), (-1, -1), 0.25, colors.HexColor("#D9DEE7")),
        ]))
        story.append(ctbl)

    story.append(Spacer(1, 14))
    story.append(Paragraph(
        "Generated by Skuvvy · AI Senior Merchandiser. Figures derived from the uploaded dataset.",
        ParagraphStyle("foot", parent=styles["Normal"], fontSize=8, textColor=GREY)))

    buf = _io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4,
                            leftMargin=16*mm, rightMargin=16*mm,
                            topMargin=16*mm, bottomMargin=16*mm,
                            title="Skuvvy Executive Trade Report")
    doc.build(story)
    buf.seek(0)
    return buf


# ── Report routes ─────────────────────────────────────────────

@app.route("/report")
@require_auth
def report_page(user):
    return send_from_directory(".", "report.html")


@app.route("/api/report")
@require_auth
def get_report(user):
    """Generate (or return cached) AI trade report."""
    db = get_user_db()
    upload_id = request.args.get("upload_id") or session.get("current_upload_id")

    analysis = None
    if upload_id:
        try:
            res = db.table("uploads").select("analysis").eq("id", upload_id).eq("user_id", user["id"]).single().execute()
            if res.data:
                analysis = res.data["analysis"]
        except Exception:
            pass

    if not analysis:
        try:
            res = db.table("uploads").select("analysis").eq("user_id", user["id"]).order("created_at", desc=True).limit(1).execute()
            if res.data:
                analysis = res.data[0]["analysis"]
        except Exception:
            pass

    if not analysis:
        return jsonify({"error": "no_data", "message": "No dataset found — please upload trading data first."}), 404

    # Use cached report from session if same upload
    cache_key = f"report_{upload_id or 'latest'}"
    cached = session.get(cache_key)
    if cached and not request.args.get("refresh"):
        return jsonify({"report": cached, "kpis": analysis.get("kpis", {}), "categories": analysis.get("category_scorecard", []), "cached": True})

    try:
        report = _generate_report_sections(analysis)
        session[cache_key] = report
        return jsonify({
            "report": report,
            "kpis": analysis.get("kpis", {}),
            "categories": analysis.get("category_scorecard", []),
            "cached": False
        })
    except Exception as e:
        app.logger.error(f"Report error: {e}", exc_info=True)
        return jsonify({"error": "generation_failed", "message": str(e)}), 500


@app.route("/api/report/stream")
@require_auth
def stream_report(user):
    """SSE endpoint — streams the report generation progress."""
    db = get_user_db()
    upload_id = session.get("current_upload_id")

    analysis = None
    try:
        res = db.table("uploads").select("analysis").eq("user_id", user["id"]).order("created_at", desc=True).limit(1).execute()
        if res.data:
            analysis = res.data[0]["analysis"]
    except Exception:
        pass

    def generate():
        if not analysis:
            yield f"data: {json.dumps({'type':'error','message':'No dataset found'})}\n\n"
            return

        yield f"data: {json.dumps({'type':'status','message':'Analysing trading data…'})}\n\n"

        try:
            prompt = _build_report_prompt(analysis)
            yield f"data: {json.dumps({'type':'status','message':'Generating commentary with AI…'})}\n\n"

            collected = ""
            with ai_client.messages.stream(
                model=MODEL,
                max_tokens=3000,
                messages=[{"role": "user", "content": prompt}]
            ) as stream:
                for chunk in stream.text_stream:
                    collected += chunk
                    yield f"data: {json.dumps({'type':'chunk','text':chunk})}\n\n"

            yield f"data: {json.dumps({'type':'status','message':'Parsing report…'})}\n\n"

            raw = collected.strip()
            if raw.startswith("```"):
                raw = "\n".join(raw.split("\n")[1:])
            if raw.endswith("```"):
                raw = "\n".join(raw.split("\n")[:-1])
            report = json.loads(raw.strip())

            cache_key = f"report_{upload_id or 'latest'}"
            session[cache_key] = report

            audit.log(REPORT_GENERATED, request=request, user_id=user["id"],
                      company_id=user.get("company_id"),
                      resource="report:stream")

            yield f"data: {json.dumps({'type':'complete','report':report,'kpis':analysis.get('kpis',{}),'categories':analysis.get('category_scorecard',[])})}\n\n"

        except Exception as e:
            app.logger.error(f"Report stream error: {e}", exc_info=True)
            fallback = _fallback_report(analysis)
            yield f"data: {json.dumps({'type':'complete','report':fallback,'kpis':analysis.get('kpis',{}),'categories':analysis.get('category_scorecard',[]),'fallback':True})}\n\n"

    return Response(generate(), mimetype="text/event-stream",
                    headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})


@app.route("/api/report/pptx")
@require_auth
def download_pptx(user):
    """Generate and stream a PowerPoint trade report."""
    db = get_user_db()
    upload_id = session.get("current_upload_id")

    analysis = None
    try:
        res = db.table("uploads").select("analysis").eq("user_id", user["id"]).order("created_at", desc=True).limit(1).execute()
        if res.data:
            analysis = res.data[0]["analysis"]
    except Exception:
        pass

    if not analysis:
        return jsonify({"error": "No data"}), 404

    # Use cached report or generate fallback quickly
    cache_key = f"report_{upload_id or 'latest'}"
    report = session.get(cache_key)
    if not report:
        report = _fallback_report(analysis)

    try:
        buf = _build_pptx(report, analysis)
        from datetime import date as _d
        fname = f"Skuvvy_Trade_Report_{_d.today().strftime('%Y%m%d')}.pptx"
        audit.log(DATA_EXPORT, request=request, user_id=user["id"],
                  company_id=user.get("company_id"),
                  resource="report:pptx")
        return Response(
            buf.getvalue(),
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={fname}"}
        )
    except Exception as e:
        app.logger.error(f"PPTX error: {e}", exc_info=True)
        return jsonify({"error": "Failed to generate PowerPoint"}), 500


def _load_report_for_export(user):
    """Shared loader: returns (analysis, report) for the current/selected upload."""
    db = get_user_db()
    upload_id = request.args.get("upload_id") or session.get("current_upload_id")

    analysis = None
    if upload_id:
        try:
            res = db.table("uploads").select("analysis").eq("id", upload_id) \
                .eq("user_id", user["id"]).single().execute()
            if res.data:
                analysis = res.data["analysis"]
        except Exception:
            pass
    if not analysis:
        try:
            res = db.table("uploads").select("analysis").eq("user_id", user["id"]) \
                .order("created_at", desc=True).limit(1).execute()
            if res.data:
                analysis = res.data[0]["analysis"]
        except Exception:
            pass
    if not analysis:
        return None, None

    cache_key = f"report_{upload_id or 'latest'}"
    report = session.get(cache_key) or _fallback_report(analysis)
    return analysis, report


@app.route("/api/report/pdf")
@require_auth
def download_pdf(user):
    """Generate and stream a board-ready PDF executive report."""
    analysis, report = _load_report_for_export(user)
    if not analysis:
        return jsonify({"error": "No data"}), 404
    try:
        buf = _build_pdf(report, analysis)
        from datetime import date as _d
        fname = f"Skuvvy_Executive_Report_{_d.today().strftime('%Y%m%d')}.pdf"
        audit.log(DATA_EXPORT, request=request, user_id=user["id"],
                  company_id=user.get("company_id"), resource="report:pdf")
        return Response(
            buf.getvalue(),
            mimetype="application/pdf",
            headers={"Content-Disposition": f"attachment; filename={fname}"}
        )
    except Exception as e:
        app.logger.error(f"PDF error: {e}", exc_info=True)
        return jsonify({"error": "Failed to generate PDF"}), 500


@app.route("/api/report/xlsx")
@require_auth
def download_xlsx(user):
    """Generate and stream an Excel workbook of the analysed data."""
    analysis, report = _load_report_for_export(user)
    if not analysis:
        return jsonify({"error": "No data"}), 404
    try:
        buf = _build_xlsx(report, analysis)
        from datetime import date as _d
        fname = f"Skuvvy_Data_Export_{_d.today().strftime('%Y%m%d')}.xlsx"
        audit.log(DATA_EXPORT, request=request, user_id=user["id"],
                  company_id=user.get("company_id"), resource="report:xlsx")
        return Response(
            buf.getvalue(),
            mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            headers={"Content-Disposition": f"attachment; filename={fname}"}
        )
    except Exception as e:
        app.logger.error(f"XLSX error: {e}", exc_info=True)
        return jsonify({"error": "Failed to generate Excel"}), 500


# ── Serve static assets ───────────────────────────────────────
@app.route("/favicon.ico")
def favicon():
    return "", 204


# ═══════════════════════════════════════════════════════════════
# SCHEDULED JOBS — File expiry & data retention
# ═══════════════════════════════════════════════════════════════

def _expire_old_uploads():
    """
    Delete raw upload records that have passed their retention deadline.
    Retains only the processed analysis JSON for reports.
    Runs every hour.
    """
    try:
        now = datetime.now(timezone.utc).isoformat()
        # Mark as deleted (soft-delete — analysis JSON preserved for reports)
        res = supabase.table("uploads") \
            .update({"deleted_at": now, "file_encrypted": False}) \
            .lt("expires_at", now) \
            .is_("deleted_at", "null") \
            .execute()
        deleted = len(res.data) if res.data else 0
        if deleted:
            app.logger.info(f"[SCHEDULER] Expired {deleted} upload(s) past retention deadline")
            audit.log(FILE_DELETED, metadata={"auto_expired": deleted, "timestamp": now})
    except Exception as e:
        app.logger.error(f"[SCHEDULER] Expiry job failed: {e}")


def _start_scheduler():
    scheduler = BackgroundScheduler(daemon=True)
    scheduler.add_job(_expire_old_uploads, "interval", hours=1, id="file_expiry")
    scheduler.start()
    app.logger.info("[SCHEDULER] Background scheduler started — file expiry every 1h")
    return scheduler


# ═══════════════════════════════════════════════════════════════
# SECURITY CENTRE ROUTES
# ═══════════════════════════════════════════════════════════════

@app.route("/security-centre")
@require_auth
@require_permission("view_security_centre")
def security_centre_page(user):
    return send_from_directory(".", "security-centre.html")


@app.route("/api/security/status")
@require_auth
@require_permission("view_security_centre")
def security_status(user):
    """Security Centre: overall security posture summary."""
    company_id = user.get("company_id")
    private_mode = _company_private_mode(company_id)
    retention = _get_retention_hours(company_id)

    stats = audit.stats(company_id) if company_id else {}

    return jsonify({
        "encryption": {
            "at_rest":   "AES-256-GCM",
            "in_transit": "TLS 1.3 (enforced)",
            "status":    "ACTIVE",
        },
        "data_retention": {
            "raw_files_hours": retention,
            "analysis_retained": True,
            "policy": f"Raw files deleted after {retention}h. Analysis metrics retained indefinitely.",
        },
        "ai_privacy": {
            "private_mode":    private_mode,
            "model_training":  False,
            "third_party_sharing": False,
            "data_masking":    "Supplier names and vendor refs pseudonymised before AI processing",
        },
        "compliance": {
            "gdpr":          "Compliant",
            "soc2_ready":    True,
            "iso27001_ready": True,
            "uk_dpa":        "Compliant",
        },
        "audit": {
            "total_events":   stats.get("total_events", 0),
            "blocked_events": stats.get("blocked_events", 0),
            "immutable":      True,
        },
        "last_checked": datetime.now(timezone.utc).isoformat(),
    })


@app.route("/api/security/audit-logs")
@require_auth
@require_permission("view_audit_logs")
def get_audit_logs(user):
    """Security Centre: paginated audit log for the company."""
    company_id = user.get("company_id")
    if not company_id:
        return jsonify({"logs": [], "message": "No company associated with this account"}), 200

    limit  = min(int(request.args.get("limit", 50)), 200)
    events = audit.recent(str(company_id), limit=limit)
    return jsonify({"logs": events, "count": len(events)})


@app.route("/api/security/events")
@require_auth
@require_permission("view_audit_logs")
def get_security_events(user):
    """Security Centre: security-relevant events (failures, blocked requests)."""
    company_id = user.get("company_id")
    if not company_id:
        return jsonify({"events": []}), 200

    events = audit.security_events(str(company_id), limit=50)
    return jsonify({"events": events, "count": len(events)})


@app.route("/api/security/users")
@require_auth
@require_permission("view_all_users")
def get_company_users(user):
    """Security Centre: all users belonging to this company."""
    company_id = user.get("company_id")
    if not company_id:
        return jsonify({"users": []}), 200
    try:
        res = supabase.table("profiles") \
            .select("id, email, full_name, role, last_login, is_active, created_at") \
            .eq("company_id", str(company_id)) \
            .execute()
        return jsonify({"users": res.data or []})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/security/retention", methods=["POST"])
@require_auth
@require_permission("change_retention_policy")
def update_retention(user):
    """Security Centre: update data retention policy (enterprise only)."""
    body  = request.get_json(silent=True) or {}
    hours = int(body.get("hours", 24))
    if hours < 1 or hours > 8760:   # 1 hour to 1 year
        return jsonify({"error": "Retention must be between 1 and 8760 hours"}), 400

    company_id = user.get("company_id")
    if not company_id:
        return jsonify({"error": "No company associated"}), 400

    try:
        supabase.table("companies") \
            .update({"retention_hours": hours}) \
            .eq("id", str(company_id)).execute()

        audit.log(SETTINGS_CHANGED, request=request, user_id=user["id"],
                  company_id=company_id,
                  metadata={"setting": "retention_hours", "value": hours})

        return jsonify({"ok": True, "retention_hours": hours})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/security/private-mode", methods=["POST"])
@require_auth
@require_permission("toggle_private_mode")
def toggle_private_mode(user):
    """Security Centre: enable/disable Private Processing Mode."""
    body    = request.get_json(silent=True) or {}
    enabled = bool(body.get("enabled", False))

    company_id = user.get("company_id")
    if not company_id:
        return jsonify({"error": "No company associated"}), 400

    try:
        supabase.table("companies") \
            .update({"private_processing_mode": enabled}) \
            .eq("id", str(company_id)).execute()

        audit.log(SETTINGS_CHANGED, request=request, user_id=user["id"],
                  company_id=company_id,
                  metadata={"setting": "private_processing_mode", "value": enabled})

        return jsonify({
            "ok": True,
            "private_mode": enabled,
            "message": (
                "Private Processing Mode enabled. AI will only receive aggregated metrics — "
                "no SKU codes, product names, or supplier data will leave Skuvvy servers."
                if enabled else
                "Private Processing Mode disabled. Standard AI analysis with supplier masking is active."
            ),
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# Start background jobs
try:
    _start_scheduler()
except Exception as e:
    app.logger.warning(f"Scheduler failed to start: {e}")


if __name__ == "__main__":
    port  = int(os.environ.get("PORT", 5000))
    debug = os.environ.get("FLASK_ENV", "development") == "development"
    print(f"\n🧵 Skuvvy running → http://localhost:{port}\n")
    app.run(host="0.0.0.0", port=port, debug=debug)
